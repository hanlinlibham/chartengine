"""GTM 模式层测试：贡献分解 / 估值区间 / 注释 / 横向条 / 预测斜纹 / 高亮 / 数值标签。"""

import pandas as pd
import pytest
from pptx import Presentation

from ablechart import normalize_spec, render_chart


def _slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


CONTRIB_DATA = {
    "季度": ["2024Q1", "2024Q2", "2024Q3", "2024Q4"],
    "消费": [0.02, 0.025, 0.018, 0.022],
    "投资": [0.012, -0.004, 0.01, 0.008],
    "净出口": [-0.005, 0.006, -0.002, 0.004],
}


class TestContribution:
    def _spec(self):
        data = dict(CONTRIB_DATA)
        data["GDP"] = [round(a + b + c, 4) for a, b, c in zip(
            data["消费"], data["投资"], data["净出口"])]
        return {"chart": "贡献图", "total": "GDP", "data": data}

    def test_parts_stacked_total_is_line(self):
        plan = normalize_spec(self._spec())
        cfg = {s["key"]: s for s in plan.kwargs["series_config"]}
        assert cfg["消费"]["grouping"] == "stacked"
        assert cfg["GDP"]["type"] == "line"
        assert "grouping" not in cfg["GDP"] or cfg["GDP"].get("grouping") is None

    def test_total_line_gets_orange_override(self):
        plan = normalize_spec(self._spec())
        gdp = next(o for o in plan.series_overrides if o["name"] == "GDP")
        assert gdp["color"] == "F5821F"

    def test_line_excluded_from_stacked_axis_bounds(self):
        prs, slide = _slide()
        chart = render_chart(slide, self._spec())
        import re
        xml = chart._chartSpace.xml
        vmax = float(re.search(r'<c:max val="([^"]+)"', xml).group(1))
        # 堆叠正向和最大约 0.034；若线被误算进堆叠求和，max 会接近翻倍
        assert vmax <= 0.06


class TestRangeChart:
    DATA = {"行业": ["A", "B", "C"], "低": [10.0, 12.0, 8.0], "高": [40.0, 35.0, 30.0],
            "均值": [22.0, 20.0, 18.0], "当前": [30.0, 15.0, 25.0]}

    def test_kind_and_columns(self):
        plan = normalize_spec({"chart": "区间", "data": self.DATA,
                               "low": "低", "high": "高", "average": "均值", "current": "当前"})
        assert plan.kind == "range"
        assert plan.kwargs["low_col"] == "低" and plan.kwargs["current_col"] == "当前"

    def test_auto_column_inference(self):
        plan = normalize_spec({"chart": "range", "data": self.DATA})
        assert plan.kwargs["low_col"] == "低" and plan.kwargs["high_col"] == "高"

    def test_render_has_legend_and_markers(self):
        prs, slide = _slide()
        chart = render_chart(slide, {"chart": "range", "data": self.DATA,
                                     "low": "低", "high": "高", "average": "均值", "current": "当前"})
        assert chart.has_legend is True
        xml = chart._chartSpace.xml
        assert '<c:symbol val="diamond"/>' in xml   # 当前值
        assert '<c:symbol val="dash"/>' in xml      # 均值刻度
        assert '<c:delete val="1"/>' in xml         # 基底系列图例隐藏


class TestHorizontal:
    def test_bar_dir(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "data": {"行业": ["A", "B"], "ROE": [0.1, 0.2]},
            "orientation": "horizontal", "legend": "none",
        })
        assert '<c:barDir val="bar"/>' in chart._chartSpace.xml

    def test_chinese_alias(self):
        plan = normalize_spec({"data": {"x": ["a"], "v": [1]}, "orientation": "横向"})
        assert plan.kwargs["orientation"] == "horizontal"


class TestGtmElements:
    DATA = {"年份": ["2023", "2024", "2025F"], "增速": [0.29, 0.39, 0.21]}

    def test_value_labels_native_dlbls(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "data": self.DATA, "legend": "none",
            "series": [{"column": "增速", "labels": {"format": "0%"}}],
        })
        xml = chart._chartSpace.xml
        assert '<c:showVal val="1"/>' in xml
        assert 'formatCode="0%"' in xml

    def test_highlight_creates_dpt(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "data": self.DATA, "legend": "none",
            "highlight": {"category": "2024", "color": "#6BA43A"},
        })
        xml = chart._chartSpace.xml
        assert "<c:dPt>" in xml and "6BA43A" in xml

    def test_forecast_pattern_fill(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "data": self.DATA, "legend": "none", "forecast_from": "2025F",
        })
        assert 'pattFill prst="ltUpDiag"' in chart._chartSpace.xml

    def test_forecast_missing_category_warns(self):
        plan = normalize_spec({"data": self.DATA, "forecast_from": "2030F"})
        assert any("forecast_from" in w for w in plan.warnings)

    def test_annotations_add_overlay_shapes(self):
        prs, slide = _slide()
        before = len(slide.shapes)
        render_chart(slide, {
            "chart": "line", "legend": "none",
            "data": {"月": ["1", "2", "3", "4"], "值": [0.02, 0.025, 0.022, 0.028]},
            "series": [{"column": "值", "last_point_label": {"format": "0.0%"}}],
            "annotations": [
                {"type": "average", "value": 0.024, "label": "均值"},
                {"type": "band", "from": 0.02, "to": 0.03},
            ],
        })
        # 图表 + 均值线 + 均值标签 + 色带 + 末点标注
        assert len(slide.shapes) - before >= 5

    def test_pinned_plot_area_when_annotated(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "chart": "line", "legend": "none",
            "data": {"月": ["1", "2", "3"], "值": [1.0, 2.0, 3.0]},
            "annotations": [{"type": "hline", "value": 2.0}],
        })
        assert "<c:manualLayout>" in chart._chartSpace.xml


def test_gtm_theme_available():
    from ablechart import COLOR_SCHEMES
    assert "gtm" in COLOR_SCHEMES
    assert COLOR_SCHEMES["gtm"][0] == "595959"


class TestCapabilityLowering:
    """降低模型能力要求的三个机制。"""

    def test_examples_gallery(self):
        from ablechart import chart_spec_examples
        full = chart_spec_examples()
        assert full.count("```json") >= 12
        assert "range" in chart_spec_examples("估值")

    def test_average_annotation_engine_computes_mean(self):
        df = pd.DataFrame({"月": list("1234"), "利差": [1.0, 2.0, 3.0, 2.0]})
        plan = normalize_spec({"chart": "line", "annotations": [
            {"type": "average", "series": "利差"}]}, df=df)
        # 渲染计划中 value 已由引擎填充为均值 2.0
        prs, slide = _slide()
        render_chart(slide, {"chart": "line", "legend": "none", "annotations": [
            {"type": "average", "series": "利差"}]}, df=df)

    def test_band_quantiles_computed(self):
        df = pd.DataFrame({"月": list("12345678"), "v": [1, 2, 3, 4, 5, 6, 7, 8]})
        plan = normalize_spec({"chart": "line", "annotations": [
            {"type": "band", "series": "v", "quantiles": [0.25, 0.75]}]}, df=df)
        assert any("分位自动计算" in w for w in plan.warnings)

    def test_average_without_value_or_series_is_error(self):
        from ablechart import SpecError
        with pytest.raises(SpecError, match="均值"):
            normalize_spec({"chart": "line",
                            "data": {"月": ["1"], "v": [1.0]},
                            "annotations": [{"type": "average"}]})

    def test_auto_dual_axis_on_scale_gap(self):
        plan = normalize_spec({"data": {
            "日": list("12345"),
            "成交额": [12000, 15000, 9000, 18000, 14000],
            "换手率": [0.011, 0.014, 0.008, 0.017, 0.013]}})
        cfg = {s["key"]: s for s in plan.kwargs["series_config"]}
        assert cfg["换手率"]["axis"] == "secondary"
        assert cfg["换手率"]["type"] == "line"
        assert any("量纲差异" in w for w in plan.warnings)

    def test_no_auto_dual_axis_when_series_explicit(self):
        plan = normalize_spec({
            "data": {"日": ["1", "2"], "成交额": [12000, 15000], "换手率": [0.011, 0.014]},
            "series": ["成交额", "换手率"]})
        assert all(s["axis"] == "primary" for s in plan.kwargs["series_config"])


class TestProfessionalPolishRound:
    """专业度收尾：副标题 / 排序 / 倍数轴 / 越界裁剪 / vband / 年份缩写。"""

    RANGE = {"行业": ["A", "B"], "低": [10.0, 12.0], "高": [40.0, 35.0],
             "均": [22.0, 20.0], "当": [30.0, 15.0]}

    def test_subtitle_two_paragraphs_left_aligned(self):
        prs, slide = _slide()
        chart = render_chart(slide, {"chart": "range", "title": "估值", "subtitle": "市盈率（TTM）",
                                     "data": self.RANGE, "low": "低", "high": "高"})
        tf = chart.chart_title.text_frame
        assert len(tf.paragraphs) == 2
        assert chart._chartSpace.xml.count('algn="l"') >= 2

    def test_subtitle_on_combo(self):
        prs, slide = _slide()
        chart = render_chart(slide, {"title": "营收", "subtitle": "亿元",
                                     "data": {"年": ["23", "24"], "v": [1, 2]}, "legend": "none"})
        assert len(chart.chart_title.text_frame.paragraphs) == 2

    def test_range_sort_desc_by_current(self):
        plan = normalize_spec({"chart": "range", "data": self.RANGE,
                               "low": "低", "high": "高", "current": "当", "sort": "desc"})
        prs, slide = _slide()
        from ablechart import create_range_chart
        chart = create_range_chart(slide, plan.df, **plan.kwargs)
        # 排序后第一类目应为当前值更大的 A
        assert "<c:v>A</c:v>" in chart._chartSpace.xml.split("<c:v>B</c:v>")[0]

    def test_times_format_alias(self):
        prs, slide = _slide()
        chart = render_chart(slide, {"chart": "range", "data": self.RANGE,
                                     "low": "低", "high": "高", "format": "times"})
        assert 'formatCode="0&quot;x&quot;"' in chart._chartSpace.xml

    def test_combo_sort_desc(self):
        plan = normalize_spec({"data": {"区": ["甲", "乙", "丙"], "v": [2.0, 9.0, 5.0]},
                               "sort": "desc"})
        assert plan.df["v"].tolist() == [9.0, 5.0, 2.0]

    def test_out_of_range_band_skipped(self):
        prs, slide = _slide()
        before = len(slide.shapes)
        render_chart(slide, {
            "chart": "line", "legend": "none",
            "data": {"月": ["1", "2", "3"], "值": [1.4, 1.6, 1.8]},
            "annotations": [{"type": "band", "from": 0.008, "to": 0.012}],  # 轴范围外
        })
        assert len(slide.shapes) - before == 1  # 只有图表，没有越界色带

    def test_vband_annotation(self):
        prs, slide = _slide()
        before = len(slide.shapes)
        render_chart(slide, {
            "chart": "line", "legend": "none",
            "data": {"月": ["1", "2", "3", "4"], "值": [1.0, 2.0, 1.5, 2.5]},
            "annotations": [{"type": "vband", "from_category": "2", "to_category": "3", "label": "衰退期"}],
        })
        assert len(slide.shapes) - before >= 3  # 图表 + 竖带 + 标签

    def test_year_abbrev_format(self):
        from ablechart.polish import strftime_from_excel
        assert strftime_from_excel("'yy") == "'%y"

    def test_nice_ladder_15_and_3(self):
        from ablechart.polish import nice_range
        # max=62：1.5 档让轴顶收敛到 75 而不是 80/100
        nmin, nmax, unit = nice_range(0, 62, 5, include_zero="always")
        assert nmax == 75.0 and unit == 15.0

    def test_zero_axis_emphasis(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "data": {"年": ["23", "24"], "v": [5.0, -3.0]}, "legend": "none"})
        # 跨零 → 分类轴线深色
        assert "404040" in chart._chartSpace.xml

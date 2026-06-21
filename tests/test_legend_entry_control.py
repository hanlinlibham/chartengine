"""Per-series legend-entry control (issue #9): hide one series from the legend.

`show_in_legend=False` on a series (config) or `legend: false` (spec) emits a
`c:legendEntry` with `delete=1` for that series, so it disappears from the
legend while staying on the plot.
"""
import pandas as pd
from pptx import Presentation

from ablechart import create_combo_chart, render_chart, ChartLayoutConfig, LegendConfig

_C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"


def _deleted_legend_indices(chart_el):
    out = set()
    legend = chart_el.find(f".//{_C}legend")
    if legend is None:
        return out
    for entry in legend.findall(f"{_C}legendEntry"):
        idx = entry.find(f"{_C}idx")
        delete = entry.find(f"{_C}delete")
        if idx is not None and delete is not None and delete.get("val") in ("1", "true"):
            out.add(int(idx.get("val")))
    return out


def _df():
    return pd.DataFrame({"FY": ["2022", "2023", "2024"],
                         "Revenue": [328.6, 400.9, 362.0],
                         "Target": [350.0, 380.0, 400.0]})


def test_config_hides_series_from_legend(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide, df=_df(), categories_col="FY",
        series_config=[
            {"key": "Revenue", "name": "Revenue", "type": "column", "axis": "primary"},
            # a reference/target line we don't want cluttering the legend
            {"key": "Target", "name": "Target", "type": "line", "axis": "primary",
             "show_in_legend": False},
        ],
        layout_config=ChartLayoutConfig(legend_config=LegendConfig()),
    )
    out = tmp_path / "c.pptx"
    prs.save(out)
    chart = Presentation(str(out)).slides[0].shapes[0].chart
    # Target is series index 1
    assert 1 in _deleted_legend_indices(chart._element)


def test_spec_legend_false_hides_series(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_chart(slide, {
        "data": None,
        "categories": "FY",
        "series": [
            {"column": "Revenue", "type": "column"},
            {"column": "Target", "type": "line", "legend": False},
        ],
    }, _df())
    out = tmp_path / "s.pptx"
    prs.save(out)
    chart = Presentation(str(out)).slides[0].shapes[0].chart
    assert 1 in _deleted_legend_indices(chart._element)


def test_no_hiding_by_default(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide, df=_df(), categories_col="FY",
        series_config=[
            {"key": "Revenue", "name": "Revenue", "type": "column", "axis": "primary"},
            {"key": "Target", "name": "Target", "type": "line", "axis": "primary"},
        ],
        layout_config=ChartLayoutConfig(legend_config=LegendConfig()),
    )
    out = tmp_path / "n.pptx"
    prs.save(out)
    chart = Presentation(str(out)).slides[0].shapes[0].chart
    assert _deleted_legend_indices(chart._element) == set()

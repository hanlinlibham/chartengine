"""Contract tests verifying inspect / replace work on charts NOT authored by
this engine.

Why a separate suite: tests in ``test_inspect.py`` / ``test_replace.py``
inline-generate fixtures via ``pptchartengine.create_*_chart``, which means
we're testing the engine against its own output. Real **template update**
scenarios (ADR-0001 northern star) involve ``.pptx`` authored in PowerPoint
desktop, Keynote, Google Slides export, or other engines.

This file uses python-pptx's native ``slide.shapes.add_chart(...)`` API
directly, bypassing **all** pptchartengine code paths during fixture
generation. If inspect/replace need engine-written metadata to work, these
tests will fail.
"""

from __future__ import annotations

import os

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from pptchartengine import (
    SeriesData,
    inspect_pptx_charts,
    replace_pptx_chart_data,
)


# ---------------------------------------------------------------------------
# Fixture helper — pure python-pptx, no pptchartengine
# ---------------------------------------------------------------------------


def _make_native_pptx_chart(
    tmp_path,
    *,
    filename: str = "native.pptx",
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.BAR_CLUSTERED,
):
    """Build a .pptx with one chart, using **only** python-pptx native API."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Sales", (100.0, 200.0, 150.0, 250.0))
    chart_data.add_series("Costs", (80.0, 130.0, 110.0, 180.0))

    slide.shapes.add_chart(
        chart_type,
        Inches(1), Inches(1), Inches(8), Inches(5),
        chart_data,
    )

    out = tmp_path / filename
    prs.save(str(out))
    return str(out)


# ---------------------------------------------------------------------------
# inspect on non-engine charts
# ---------------------------------------------------------------------------


def test_inspect_non_engine_authored_bar_chart(tmp_path):
    """python-pptx-authored bar chart inspects correctly without engine metadata."""
    input_pptx = _make_native_pptx_chart(tmp_path)
    inv = inspect_pptx_charts(input_pptx)

    assert len(inv) == 1
    item = inv[0]
    assert item.chart_type == "bar"
    assert item.category_count == 4
    assert item.series_count == 2
    assert "Sales" in item.series_names
    assert "Costs" in item.series_names
    assert item.has_embedded_workbook is True
    assert item.replaceable is True
    assert item.warnings == []


def test_inspect_non_engine_authored_line_chart(tmp_path):
    """python-pptx-authored line chart classified as 'line'."""
    input_pptx = _make_native_pptx_chart(
        tmp_path, chart_type=XL_CHART_TYPE.LINE, filename="line.pptx",
    )
    inv = inspect_pptx_charts(input_pptx)
    assert len(inv) == 1
    assert inv[0].chart_type == "line"
    assert inv[0].replaceable is True


# ---------------------------------------------------------------------------
# replace on non-engine charts
# ---------------------------------------------------------------------------


def test_replace_non_engine_authored_bar_chart_end_to_end(tmp_path):
    """python-pptx-authored chart can be replaced; new data round-trips through inspect."""
    input_pptx = _make_native_pptx_chart(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    output_pptx = str(tmp_path / "out.pptx")

    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=["2026Q1", "2026Q2"],
        series=[
            SeriesData(name="Revenue", values=[500.0, 600.0]),
        ],
    )

    assert result.status == "ok", f"{result.error_code}: {result.error_detail}"
    assert result.categories_replaced == 2
    assert result.series_replaced == 1

    inv_after = inspect_pptx_charts(output_pptx)
    assert len(inv_after) == 1
    item = inv_after[0]
    assert item.category_count == 2
    assert item.series_count == 1
    assert "Revenue" in item.series_names
    assert item.chart_type == "bar"  # type preserved
    # shape identity preserved across non-engine origin
    assert item.selector.shape_id == inv[0].selector.shape_id


# ---------------------------------------------------------------------------
# Sanity assertion on project sample template (observed fact, locks behaviour)
# ---------------------------------------------------------------------------


def test_project_sample_template_has_zero_native_charts():
    """The project's ``financial_charts_template_sample.pptx`` actually contains
    zero native chart shapes (text/placeholder/image only).

    This is documented as an observed fact during 2026-05-24 inspect inventory:
    despite the filename saying "charts_template", the file has no native chart
    shapes that ``shape.has_chart`` returns True for. The know-how-brief §2
    rule "agent cannot guess from filename" applies — only inspection settles
    the question.

    This test pins the observation so we notice if:
    - the sample template gets edited to add real charts
    - someone replaces the sample file with a different one
    - inspect_pptx_charts changes its definition of "what counts as a chart"
    """
    sample = "/Users/jameslee/pension_plan/ppt-project/financial_charts_template_sample.pptx"
    if not os.path.exists(sample):
        pytest.skip(f"sample template not present: {sample}")

    inv = inspect_pptx_charts(sample)
    assert inv == [], (
        f"Sample template now reports charts (was empty when contract was written): {inv}. "
        f"Either the file was edited or inspect's definition changed — investigate before updating."
    )

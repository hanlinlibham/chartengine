"""Contract tests for ``inspect_pptx_charts`` (ADR-0006 + ADR-0007 §1 inspect lifecycle).

TDD spec per ADR-0005 §1: these tests describe the inventory contract every
chart must satisfy when scanned by ``inspect_pptx_charts``, independent of how
the chart was authored. Tests are written **before** the implementation;
they fail with ``NotImplementedError`` until ``inspect_pptx_charts`` lands.

Fixture strategy mirrors ``test_package_contract.py``: inline-generate .pptx
in ``tmp_path``, no external fixture files.
"""

from __future__ import annotations

import pandas as pd
import pytest
from pptx import Presentation

from pptchartengine import create_combo_chart
from pptchartengine.inspect import (
    ChartInventoryItem,
    ChartSelector,
    inspect_pptx_charts,
)


# ---------------------------------------------------------------------------
# Fixture helpers (in-test generation, no external .pptx)
# ---------------------------------------------------------------------------


def _make_combo_pptx(tmp_path, *, n_slides: int = 1, charts_per_slide: int = 1):
    """Build a .pptx with ``n_slides`` × ``charts_per_slide`` combo charts.

    Each combo = bar(revenue) + line(growth_pct) on dual axes,
    3 monthly categories. Series names embed (slide_i, chart_i) so assertions
    can verify per-chart identity.
    """
    prs = Presentation()
    for slide_i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for chart_i in range(charts_per_slide):
            df = pd.DataFrame({
                "month": ["1月", "2月", "3月"],
                "revenue": [100 + chart_i, 120, 110],
                "growth_pct": [5.0, 8.0, 4.5],
            })
            create_combo_chart(
                slide=slide,
                df=df,
                categories_col="month",
                series_config=[
                    {"key": "revenue", "name": f"Revenue {slide_i}.{chart_i}",
                     "type": "bar", "axis": "primary"},
                    {"key": "growth_pct", "name": f"Growth {slide_i}.{chart_i}",
                     "type": "line", "axis": "secondary"},
                ],
            )
    out = tmp_path / "fixture.pptx"
    prs.save(str(out))
    return str(out)


# ---------------------------------------------------------------------------
# Empty / no-chart edge cases (must not raise — ADR-0006 §1 implicit)
# ---------------------------------------------------------------------------


def test_inspect_empty_pptx_returns_empty_list(tmp_path):
    """Empty .pptx (no slides) → ``[]``. Must not raise."""
    prs = Presentation()
    out = tmp_path / "empty.pptx"
    prs.save(str(out))

    inventory = inspect_pptx_charts(str(out))
    assert inventory == []


def test_inspect_slides_without_charts_returns_empty_list(tmp_path):
    """Slides exist but no charts → ``[]``."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    out = tmp_path / "no_charts.pptx"
    prs.save(str(out))

    inventory = inspect_pptx_charts(str(out))
    assert inventory == []


# ---------------------------------------------------------------------------
# Single chart — full inventory contract
# ---------------------------------------------------------------------------


def test_inspect_single_combo_chart_returns_one_item(tmp_path):
    """One combo chart → 1 inventory item with all ADR-0006 §1 facts."""
    pptx_path = _make_combo_pptx(tmp_path, n_slides=1, charts_per_slide=1)

    inventory = inspect_pptx_charts(pptx_path)

    assert len(inventory) == 1
    item = inventory[0]
    assert isinstance(item, ChartInventoryItem)
    assert isinstance(item.selector, ChartSelector)

    # selector facts (ADR-0006 §1)
    assert item.selector.slide_index == 0
    assert item.selector.shape_id > 0
    assert item.selector.chart_part.startswith("ppt/charts/chart")
    assert item.selector.chart_part.endswith(".xml")

    # per-chart facts
    assert item.chart_index_on_slide == 0
    assert item.category_count == 3
    assert item.series_count == 2
    assert "Revenue 0.0" in item.series_names
    assert "Growth 0.0" in item.series_names

    # ADR-0006 §3: combo (bar+line) is in the first-batch supported set
    assert item.has_embedded_workbook is True
    assert item.replaceable is True
    assert item.warnings == []


# ---------------------------------------------------------------------------
# Multi-chart ordering
# ---------------------------------------------------------------------------


def test_inspect_multiple_slides_charts_ordered_by_slide_then_chart_index(tmp_path):
    """Multi-slide multi-chart → ordered by ``(slide_index, chart_index_on_slide)``."""
    pptx_path = _make_combo_pptx(tmp_path, n_slides=2, charts_per_slide=2)

    inventory = inspect_pptx_charts(pptx_path)

    assert len(inventory) == 4
    indices = [(i.selector.slide_index, i.chart_index_on_slide) for i in inventory]
    assert indices == [(0, 0), (0, 1), (1, 0), (1, 1)]


# ---------------------------------------------------------------------------
# Selector stability (ADR-0006 §1 — selector must survive re-inspection)
# ---------------------------------------------------------------------------


def test_inspect_selector_chart_part_is_stable_across_inspect_calls(tmp_path):
    """Two consecutive inspects on the same .pptx return identical chart_part for each chart."""
    pptx_path = _make_combo_pptx(tmp_path, n_slides=1, charts_per_slide=2)

    inv_a = inspect_pptx_charts(pptx_path)
    inv_b = inspect_pptx_charts(pptx_path)

    parts_a = [i.selector.chart_part for i in inv_a]
    parts_b = [i.selector.chart_part for i in inv_b]
    assert parts_a == parts_b
    assert len(set(parts_a)) == 2  # two distinct chart parts


# ---------------------------------------------------------------------------
# Read-only contract (ADR-0007 §1 — inspect must not mutate the .pptx)
# ---------------------------------------------------------------------------


def test_inspect_does_not_mutate_pptx_file(tmp_path):
    """File bytes must be byte-identical before vs after inspect (read-only contract)."""
    pptx_path = _make_combo_pptx(tmp_path, n_slides=1, charts_per_slide=1)
    before = open(pptx_path, "rb").read()

    inspect_pptx_charts(pptx_path)

    after = open(pptx_path, "rb").read()
    assert before == after, "inspect_pptx_charts must not mutate the input .pptx"

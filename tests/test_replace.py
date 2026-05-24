"""Contract tests for ``replace_pptx_chart_data`` (ADR-0006 §2/§5 + ADR-0007 §4).

TDD spec per ADR-0005 §1: each test maps to one of the ADR-0006 §5 invariants.
Tests are written **before** the implementation; they fail with
``NotImplementedError`` until ``replace_pptx_chart_data`` lands.

Fixture strategy mirrors :mod:`tests.test_inspect`: inline-generate the
input .pptx in ``tmp_path``, no external fixture files.
"""

from __future__ import annotations

import pandas as pd
import pytest
from pptx import Presentation

from pptchartengine import (
    create_bubble_chart,
    create_combo_chart,
    create_scatter_chart,
    inspect_pptx_charts,
)
from pptchartengine.replace import (
    ReplaceResult,
    SeriesData,
    replace_pptx_chart_data,
)


# ---------------------------------------------------------------------------
# Fixture helper
# ---------------------------------------------------------------------------


def _make_combo_pptx(tmp_path, filename="input.pptx"):
    """Build a single-slide single combo (bar+line) chart for replace tests."""
    df = pd.DataFrame({
        "month": ["1月", "2月", "3月"],
        "revenue": [100, 120, 110],
        "growth_pct": [5.0, 8.0, 4.5],
    })
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide,
        df=df,
        categories_col="month",
        series_config=[
            {"key": "revenue", "name": "Revenue", "type": "bar", "axis": "primary"},
            {"key": "growth_pct", "name": "Growth %", "type": "line", "axis": "secondary"},
        ],
    )
    out = tmp_path / filename
    prs.save(str(out))
    return str(out)


# ---------------------------------------------------------------------------
# 1. Happy path
# ---------------------------------------------------------------------------


def test_replace_combo_returns_ok_status(tmp_path):
    """Basic happy path: replace combo data → ``status="ok"`` + counts."""
    input_pptx = _make_combo_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    selector = inv[0].selector

    output_pptx = str(tmp_path / "output.pptx")
    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=selector,
        categories=["4月", "5月", "6月"],
        series=[
            SeriesData(name="Revenue", values=[200.0, 220.0, 210.0]),
            SeriesData(name="Growth %", values=[10.0, 12.0, 8.5]),
        ],
    )

    assert isinstance(result, ReplaceResult)
    assert result.status == "ok"
    assert result.series_replaced == 2
    assert result.categories_replaced == 3
    assert result.error_code is None


# ---------------------------------------------------------------------------
# 2. Input not modified (ADR-0001 «不允许覆盖用户原始模板»)
# ---------------------------------------------------------------------------


def test_replace_does_not_modify_input_pptx(tmp_path):
    """Input file bytes must be byte-identical before vs after replace."""
    input_pptx = _make_combo_pptx(tmp_path)
    before = open(input_pptx, "rb").read()

    inv = inspect_pptx_charts(input_pptx)
    output_pptx = str(tmp_path / "output.pptx")
    replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=["A", "B", "C"],
        series=[
            SeriesData(name="Revenue", values=[1.0, 2.0, 3.0]),
            SeriesData(name="Growth %", values=[4.0, 5.0, 6.0]),
        ],
    )

    after = open(input_pptx, "rb").read()
    assert before == after, "replace must not modify the input .pptx"


# ---------------------------------------------------------------------------
# 3. Shape identity preserved (ADR-0006 §5.3 + ADR-0007 §4)
# ---------------------------------------------------------------------------


def test_replace_preserves_shape_identity(tmp_path):
    """shape_id, left, top, width, height unchanged after replace."""
    input_pptx = _make_combo_pptx(tmp_path)
    inv_before = inspect_pptx_charts(input_pptx)
    shape_id_before = inv_before[0].selector.shape_id

    prs_in = Presentation(input_pptx)
    shape_in = next(s for s in prs_in.slides[0].shapes if s.has_chart)
    pos_before = (shape_in.left, shape_in.top, shape_in.width, shape_in.height)

    output_pptx = str(tmp_path / "output.pptx")
    replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv_before[0].selector,
        categories=["X", "Y", "Z"],
        series=[
            SeriesData(name="Revenue", values=[50.0, 60.0, 70.0]),
            SeriesData(name="Growth %", values=[1.0, 2.0, 3.0]),
        ],
    )

    inv_after = inspect_pptx_charts(output_pptx)
    prs_out = Presentation(output_pptx)
    shape_out = next(s for s in prs_out.slides[0].shapes if s.has_chart)
    pos_after = (shape_out.left, shape_out.top, shape_out.width, shape_out.height)

    assert inv_after[0].selector.shape_id == shape_id_before, \
        "ADR-0007 §4: shape_id must survive replace (forbids delete+reinsert)"
    assert pos_after == pos_before, "ADR-0006 §5.3: position/size must not change"


# ---------------------------------------------------------------------------
# 4. Chart type preserved (ADR-0006 §5.4)
# ---------------------------------------------------------------------------


def test_replace_preserves_chart_type(tmp_path):
    """chart_type unchanged after replace (combo stays combo)."""
    input_pptx = _make_combo_pptx(tmp_path)
    inv_before = inspect_pptx_charts(input_pptx)
    chart_type_before = inv_before[0].chart_type

    output_pptx = str(tmp_path / "output.pptx")
    replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv_before[0].selector,
        categories=["X", "Y", "Z"],
        series=[
            SeriesData(name="Revenue", values=[1.0, 2.0, 3.0]),
            SeriesData(name="Growth %", values=[4.0, 5.0, 6.0]),
        ],
    )

    inv_after = inspect_pptx_charts(output_pptx)
    assert inv_after[0].chart_type == chart_type_before


# ---------------------------------------------------------------------------
# 5. Structured failure on shape mismatch (ADR-0006 §3 error code)
# ---------------------------------------------------------------------------


def test_replace_categories_length_mismatch_fails_loud(tmp_path):
    """categories len != series.values len → ``error_code="category_length_mismatch"``."""
    input_pptx = _make_combo_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    output_pptx = str(tmp_path / "output.pptx")

    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=["A", "B", "C"],  # 3 categories
        series=[SeriesData(name="X", values=[1.0, 2.0])],  # 2 values
    )

    assert result.status == "failed"
    assert result.error_code == "category_length_mismatch"
    # ADR-0006 §5.8: failed run must not produce output file
    import os
    assert not os.path.exists(output_pptx), \
        "ADR-0006 §5.8: failed replace must not produce a half-written output"


# ---------------------------------------------------------------------------
# 6. Data round-trip: replace → inspect → new data (ADR-0006 §5.5/§5.7)
# ---------------------------------------------------------------------------


def test_replace_then_inspect_reflects_new_data(tmp_path):
    """After replace, inspect reads back new categories / series_names / counts."""
    input_pptx = _make_combo_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    output_pptx = str(tmp_path / "output.pptx")

    replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=["Q1", "Q2", "Q3", "Q4"],  # changed from 3 → 4
        series=[
            SeriesData(name="Revenue", values=[100.0, 200.0, 300.0, 400.0]),
            SeriesData(name="Growth %", values=[5.0, 10.0, 15.0, 20.0]),
        ],
    )

    inv_after = inspect_pptx_charts(output_pptx)
    assert len(inv_after) == 1
    item = inv_after[0]
    assert item.category_count == 4
    assert item.series_count == 2
    assert "Revenue" in item.series_names
    assert "Growth %" in item.series_names


# ---------------------------------------------------------------------------
# 7. selector_resolved returned correctly
# ---------------------------------------------------------------------------


def test_replace_returns_resolved_selector(tmp_path):
    """``ReplaceResult.selector_resolved`` echoes the resolved selector."""
    input_pptx = _make_combo_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    input_selector = inv[0].selector
    output_pptx = str(tmp_path / "output.pptx")

    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=input_selector,
        categories=["A", "B", "C"],
        series=[
            SeriesData(name="Revenue", values=[1.0, 2.0, 3.0]),
            SeriesData(name="Growth %", values=[4.0, 5.0, 6.0]),
        ],
    )

    assert result.selector_resolved == input_selector
    assert result.chart_part == input_selector.chart_part


# ===========================================================================
# Scatter / Bubble replace (ADR-0006 §3 first-batch remaining types)
# ===========================================================================


def _make_scatter_pptx(tmp_path, filename="scatter.pptx"):
    """Build a single-slide scatter chart fixture."""
    df = pd.DataFrame({
        "volatility": [8.1, 9.2, 7.4],
        "return_a": [10.5, 12.0, 8.8],
    })
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_scatter_chart(
        slide=slide,
        df=df,
        x_col="volatility",
        y_col="return_a",
        series_name="Risk/Return",
    )
    out = tmp_path / filename
    prs.save(str(out))
    return str(out)


def _make_bubble_pptx(tmp_path, filename="bubble.pptx"):
    """Build a single-slide bubble chart fixture."""
    df = pd.DataFrame({
        "volatility": [8.1, 9.2, 7.4],
        "return": [10.5, 12.0, 8.8],
        "aum": [120.0, 210.0, 95.0],
    })
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_bubble_chart(
        slide=slide,
        df=df,
        x_col="volatility",
        y_col="return",
        size_col="aum",
        series_name="Risk/Return/AUM",
    )
    out = tmp_path / filename
    prs.save(str(out))
    return str(out)


def test_replace_scatter_returns_ok(tmp_path):
    """Scatter chart replace with new (x, y) → status='ok'."""
    input_pptx = _make_scatter_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    assert inv[0].chart_type == "scatter", "fixture sanity: scatter detected"

    output_pptx = str(tmp_path / "output.pptx")
    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=None,
        series=[
            SeriesData(
                name="Risk/Return",
                values=[15.0, 18.0, 12.0, 20.0],
                x_values=[10.0, 12.0, 9.0, 14.0],
            ),
        ],
    )

    assert result.status == "ok", f"failed: {result.error_code} / {result.error_detail}"
    assert result.series_replaced == 1
    assert result.categories_replaced == 0  # scatter has no categories
    assert result.data_points_replaced == 4


def test_replace_bubble_returns_ok(tmp_path):
    """Bubble chart replace with new (x, y, size) → status='ok'."""
    input_pptx = _make_bubble_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    assert inv[0].chart_type == "bubble", "fixture sanity: bubble detected"

    output_pptx = str(tmp_path / "output.pptx")
    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=None,
        series=[
            SeriesData(
                name="Risk/Return/AUM",
                values=[20.0, 25.0, 15.0],
                x_values=[10.0, 12.0, 8.0],
                size_values=[150.0, 250.0, 100.0],
            ),
        ],
    )

    assert result.status == "ok", f"failed: {result.error_code} / {result.error_detail}"
    assert result.series_replaced == 1
    assert result.categories_replaced == 0
    assert result.data_points_replaced == 3


def test_replace_scatter_preserves_chart_type(tmp_path):
    """Scatter stays scatter after replace (ADR-0006 §5.4)."""
    input_pptx = _make_scatter_pptx(tmp_path)
    inv_before = inspect_pptx_charts(input_pptx)

    output_pptx = str(tmp_path / "output.pptx")
    replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv_before[0].selector,
        categories=None,
        series=[
            SeriesData(name="Risk/Return", values=[1.0, 2.0], x_values=[3.0, 4.0]),
        ],
    )

    inv_after = inspect_pptx_charts(output_pptx)
    assert inv_after[0].chart_type == "scatter"
    assert inv_after[0].selector.shape_id == inv_before[0].selector.shape_id


def test_replace_bubble_preserves_chart_type(tmp_path):
    """Bubble stays bubble after replace (ADR-0006 §5.4)."""
    input_pptx = _make_bubble_pptx(tmp_path)
    inv_before = inspect_pptx_charts(input_pptx)

    output_pptx = str(tmp_path / "output.pptx")
    replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv_before[0].selector,
        categories=None,
        series=[
            SeriesData(
                name="Risk/Return/AUM",
                values=[1.0, 2.0],
                x_values=[3.0, 4.0],
                size_values=[5.0, 6.0],
            ),
        ],
    )

    inv_after = inspect_pptx_charts(output_pptx)
    assert inv_after[0].chart_type == "bubble"
    assert inv_after[0].selector.shape_id == inv_before[0].selector.shape_id


def test_replace_scatter_without_x_values_fails_loud(tmp_path):
    """Scatter requires x_values; missing → error_code=x_values_required_for_scatter."""
    input_pptx = _make_scatter_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    output_pptx = str(tmp_path / "output.pptx")

    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=None,
        series=[SeriesData(name="X", values=[1.0, 2.0, 3.0])],  # missing x_values
    )

    assert result.status == "failed"
    assert result.error_code == "x_values_required_for_scatter"
    import os
    assert not os.path.exists(output_pptx), \
        "ADR-0006 §5.8: failed replace must not produce partial output"


def test_replace_bubble_without_size_values_fails_loud(tmp_path):
    """Bubble requires size_values; missing → error_code=size_values_required_for_bubble."""
    input_pptx = _make_bubble_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    output_pptx = str(tmp_path / "output.pptx")

    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=None,
        series=[
            SeriesData(name="X", values=[1.0, 2.0, 3.0], x_values=[4.0, 5.0, 6.0]),
        ],  # size_values missing
    )

    assert result.status == "failed"
    assert result.error_code == "size_values_required_for_bubble"


def test_replace_category_chart_without_categories_fails_loud(tmp_path):
    """Category-based chart (combo) requires categories; None → error."""
    # Reuse combo fixture from earlier (combo is category-based)
    input_pptx = _make_combo_pptx(tmp_path)
    inv = inspect_pptx_charts(input_pptx)
    output_pptx = str(tmp_path / "output.pptx")

    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv[0].selector,
        categories=None,  # missing for combo
        series=[
            SeriesData(name="Revenue", values=[1.0, 2.0, 3.0]),
            SeriesData(name="Growth %", values=[4.0, 5.0, 6.0]),
        ],
    )

    assert result.status == "failed"
    assert result.error_code == "categories_required_for_category_chart"

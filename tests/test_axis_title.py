"""Axis-title support (issue #9, Gap A): primary/secondary value + category axes.

Titles are written as native ``c:catAx/c:valAx > c:title`` and must survive a
save -> reopen round-trip and land on the correct axis.
"""
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from ablechart import (
    create_combo_chart,
    ChartLayoutConfig,
    CategoryAxisConfig,
    ValueAxisConfig,
)

_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _axis_titles(chart_el):
    """Return {axis_tag: title_text} for every axis carrying a c:title."""
    out = {}
    for tag in ("catAx", "valAx"):
        for ax in chart_el.findall(f".//{{{_C}}}{tag}"):
            title = ax.find(f"{{{_C}}}title")
            if title is None:
                continue
            texts = [t.text for t in title.findall(f".//{{{_A}}}t") if t.text]
            pos = ax.find(f"{{{_C}}}axPos")
            key = f"{tag}:{pos.get('val') if pos is not None else '?'}"
            out[key] = "".join(texts)
    return out


def _build(tmp_path):
    df = pd.DataFrame({"FY": ["2022", "2023", "2024"],
                       "Revenue": [328.6, 400.9, 362.0],
                       "Margin": [9.4, 11.0, 14.0]})
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide, df=df, categories_col="FY",
        series_config=[
            {"key": "Revenue", "name": "Revenue", "type": "column", "axis": "primary"},
            {"key": "Margin", "name": "Margin", "type": "line", "axis": "secondary"},
        ],
        layout_config=ChartLayoutConfig(
            category_axis_config=CategoryAxisConfig(axis_title="Fiscal year"),
            value_axis_config=ValueAxisConfig(number_format='"¥"0" bn"', axis_title="RMB bn"),
            secondary_value_axis_config=ValueAxisConfig(number_format='0"%"', axis_title="Net margin (%)"),
        ),
    )
    out = tmp_path / "axis-titles.pptx"
    prs.save(out)
    return out


def test_axis_titles_survive_save_and_reopen(tmp_path):
    out = _build(tmp_path)

    reopened = Presentation(str(out))
    chart = reopened.slides[0].shapes[0].chart
    titles = _axis_titles(chart._element)

    # category axis (bottom) + primary value axis (left) + secondary (right)
    assert titles.get("catAx:b") == "Fiscal year"
    assert titles.get("valAx:l") == "RMB bn"
    assert titles.get("valAx:r") == "Net margin (%)"


def test_axis_title_is_optional_and_off_by_default(tmp_path):
    df = pd.DataFrame({"FY": ["2023", "2024"], "Revenue": [400.9, 362.0]})
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide, df=df, categories_col="FY",
        series_config=[{"key": "Revenue", "name": "Revenue", "type": "column", "axis": "primary"}],
        layout_config=ChartLayoutConfig(value_axis_config=ValueAxisConfig(number_format="0")),
    )
    out = tmp_path / "no-titles.pptx"
    prs.save(out)
    chart = Presentation(str(out)).slides[0].shapes[0].chart
    assert _axis_titles(chart._element) == {}

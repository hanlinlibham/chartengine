"""Generate the persistent external-sample .pptx used by external-fixture tests.

Run on demand to refresh ``external_sample.pptx``:

    python tests/fixtures/_generate_external_sample.py

The generated file is checked in to git so:

- tests don't regenerate it every run (faster, deterministic, no test ordering coupling)
- byte-level diffs catch python-pptx / OOXML behaviour drift across versions
- simulates a real "user-uploaded template" scenario — persistent artifact
  authored without engine knowledge

**Naming**: ADR-0003 hard rule — no theme/company/brand names. Categories
and series names are generic / role-based only.

**Filename underscore prefix** prevents pytest from collecting this file
as a test module.

Layout:

- slide 0: 1 bar chart, shape explicitly named ``quarterly_performance_chart``
  (verifies ``explicit_name`` extraction)
- slide 1: 1 line chart + 1 scatter chart on same slide
  (verifies ``chart_index_on_slide`` ordering across charts)
- slide 2: 1 pie chart + 1 plain text box
  (verifies non-chart shapes are skipped by ``inspect_pptx_charts``)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def build() -> Path:
    prs = Presentation()

    # --- Slide 0: bar chart with explicit business name ----------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar_data = CategoryChartData()
    bar_data.categories = ["2024Q1", "2024Q2", "2024Q3", "2024Q4"]
    bar_data.add_series("Revenue", (100.0, 120.0, 110.0, 140.0))
    bar_data.add_series("Costs", (60.0, 70.0, 65.0, 80.0))
    bar = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(1), Inches(1), Inches(8), Inches(4),
        bar_data,
    )
    bar.name = "quarterly_performance_chart"  # explicit business tag

    # --- Slide 1: line + scatter on same slide -------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    line_data = CategoryChartData()
    line_data.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    line_data.add_series("Index A", (100.0, 102.5, 105.0, 103.0, 107.5, 110.0))
    line_data.add_series("Index B", (100.0, 99.5, 101.0, 102.5, 104.0, 105.5))
    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.5), Inches(0.5), Inches(4.5), Inches(3),
        line_data,
    )

    scatter_data = XyChartData()
    s = scatter_data.add_series("Risk-Return")
    for x, y in [(5.0, 8.0), (7.0, 10.0), (9.0, 11.0), (12.0, 13.0), (15.0, 12.0)]:
        s.add_data_point(x, y)
    slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(5.0), Inches(0.5), Inches(4.5), Inches(3),
        scatter_data,
    )

    # --- Slide 2: pie chart + non-chart text box -----------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pie_data = CategoryChartData()
    pie_data.categories = ["Equity", "Bond", "Cash", "Alt"]
    pie_data.add_series("Allocation", (40.0, 35.0, 15.0, 10.0))
    slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(1), Inches(1), Inches(5), Inches(5),
        pie_data,
    )
    tb = slide.shapes.add_textbox(Inches(7), Inches(1), Inches(3), Inches(1))
    tb.text_frame.text = "Asset allocation overview — quarterly snapshot"

    out = Path(__file__).parent / "external_sample.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build()
    print(f"Generated: {path}")

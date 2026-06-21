"""Tick-label rotation (issue #9, Gap A): a:bodyPr@rot on category/value axes."""
import pandas as pd
from pptx import Presentation

from ablechart import (
    create_combo_chart,
    ChartLayoutConfig,
    CategoryAxisConfig,
    ValueAxisConfig,
)

_C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _rotations(chart_el):
    """{axis_key: rot_value} for axes whose tick labels carry a bodyPr@rot."""
    out = {}
    for tag in ("catAx", "valAx"):
        for ax in chart_el.findall(f".//{_C}{tag}"):
            txPr = ax.find(f"{_C}txPr")
            if txPr is None:
                continue
            bodyPr = txPr.find(f"{_A}bodyPr")
            if bodyPr is None or bodyPr.get("rot") is None:
                continue
            pos = ax.find(f"{_C}axPos")
            out[f"{tag}:{pos.get('val') if pos is not None else '?'}"] = bodyPr.get("rot")
    return out


def test_tick_label_rotation_round_trips(tmp_path):
    df = pd.DataFrame({"Quarter": ["2024Q1", "2024Q2", "2024Q3"],
                       "Revenue": [328.6, 400.9, 362.0],
                       "Margin": [9.4, 11.0, 14.0]})
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide, df=df, categories_col="Quarter",
        series_config=[
            {"key": "Revenue", "name": "Revenue", "type": "column", "axis": "primary"},
            {"key": "Margin", "name": "Margin", "type": "line", "axis": "secondary"},
        ],
        layout_config=ChartLayoutConfig(
            category_axis_config=CategoryAxisConfig(tick_label_rotation=-45),
            value_axis_config=ValueAxisConfig(number_format="0", tick_label_rotation=0),
            secondary_value_axis_config=ValueAxisConfig(number_format='0"%"', tick_label_rotation=90),
        ),
    )
    out = tmp_path / "rot.pptx"
    prs.save(out)

    chart = Presentation(str(out)).slides[0].shapes[0].chart
    rots = _rotations(chart._element)
    # OOXML rot is 60000ths of a degree
    assert rots.get("catAx:b") == str(-45 * 60000)
    assert rots.get("valAx:l") == "0"
    assert rots.get("valAx:r") == str(90 * 60000)


def test_rotation_off_by_default(tmp_path):
    df = pd.DataFrame({"Q": ["a", "b"], "V": [1.0, 2.0]})
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide, df=df, categories_col="Q",
        series_config=[{"key": "V", "name": "V", "type": "column", "axis": "primary"}],
        layout_config=ChartLayoutConfig(category_axis_config=CategoryAxisConfig()),
    )
    out = tmp_path / "norot.pptx"
    prs.save(out)
    chart = Presentation(str(out)).slides[0].shapes[0].chart
    assert _rotations(chart._element) == {}

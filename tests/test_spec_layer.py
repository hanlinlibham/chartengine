"""声明式 Spec 层测试：容错归一化、智能推断、友好报错、端到端渲染。"""



import pandas as pd
import pytest
from pptx import Presentation

from ablechart import (
    SpecError,
    chart_spec_reference,
    normalize_spec,
    parse_chart_from_pptx,
    parse_waterfall_from_pptx,
    render_chart,
    validate_spec,
)
from ablechart.spec import normalize_color


def _slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _save_and_reload(prs, tmp_path):
    path = str(tmp_path / "out.pptx")
    prs.save(path)
    return path


COMBO_DATA = {
    "年份": ["2021", "2022", "2023", "2024"],
    "营收": [100, 110, 120, 140],
    "利润": [10, 12, 15, 18],
    "增速": [0.10, 0.10, 0.09, 0.17],
}


# ============================================================================
# 归一化 / 推断
# ============================================================================

class TestNormalization:
    def test_minimal_spec_infers_everything(self):
        plan = normalize_spec({"data": COMBO_DATA})
        assert plan.kind == "combo"
        assert plan.kwargs["categories_col"] == "年份"
        keys = [s["key"] for s in plan.kwargs["series_config"]]
        assert keys == ["营收", "利润", "增速"]
        assert any("series 未指定" in w for w in plan.warnings)

    def test_string_series_expand_to_dicts(self):
        plan = normalize_spec({"data": COMBO_DATA, "series": ["营收", "利润"]})
        cfg = plan.kwargs["series_config"]
        assert cfg[0] == {"key": "营收", "name": "营收", "type": "bar", "axis": "primary"}

    def test_aliases_type_axis_chinese(self):
        plan = normalize_spec({
            "data": COMBO_DATA,
            "series": [
                {"column": "营收", "type": "柱状图", "axis": "左轴"},
                {"column": "增速", "kind": "折线", "side": "right"},
            ],
        })
        cfg = plan.kwargs["series_config"]
        assert (cfg[0]["type"], cfg[0]["axis"]) == ("bar", "primary")
        assert (cfg[1]["type"], cfg[1]["axis"]) == ("line", "secondary")

    def test_chart_alias_sets_default_series_type(self):
        plan = normalize_spec({"chart": "line", "data": COMBO_DATA, "series": ["营收"]})
        assert plan.kwargs["series_config"][0]["type"] == "line"

    def test_stacked_shortcut(self):
        plan = normalize_spec({"data": COMBO_DATA, "series": ["营收", "利润"], "stacked": True})
        assert all(s["grouping"] == "stacked" for s in plan.kwargs["series_config"])

    def test_percent_grouping_alias(self):
        plan = normalize_spec({"data": COMBO_DATA, "series": ["营收"], "grouping": "100%"})
        assert plan.kwargs["series_config"][0]["grouping"] == "percent_stacked"

    def test_numeric_string_column_coerced(self):
        data = dict(COMBO_DATA, 利润=["10", "12", "15", "18"])
        plan = normalize_spec({"data": data, "series": ["利润"]})
        assert plan.df["利润"].dtype.kind in "if"

    def test_per_series_override_collected(self):
        plan = normalize_spec({
            "data": COMBO_DATA,
            "series": [{"column": "增速", "type": "line", "color": "#c00000", "line_width": 2}],
        })
        assert plan.series_overrides == [
            {"name": "增速", "color": "C00000", "line_width_pt": 2.0}
        ]

    def test_datetime_categories_get_date_axis(self):
        data = {
            "日期": pd.date_range("2024-01-01", periods=30, freq="D"),
            "净值": list(range(30)),
        }
        plan = normalize_spec({"data": pd.DataFrame(data)})
        assert plan.kwargs["layout_config"].date_axis_config is not None

    def test_date_format_converts_string_categories(self):
        data = {"月份": ["2024-01", "2024-02", "2024-03"], "净值": [1, 2, 3]}
        plan = normalize_spec({"data": data, "layout": {"x_axis": {"date_format": "yyyy/mm"}}})
        assert pd.api.types.is_datetime64_any_dtype(plan.df["月份"])


# ============================================================================
# 数据形态
# ============================================================================

class TestDataCoercion:
    def test_records_list(self):
        records = [{"年份": "2021", "营收": 100}, {"年份": "2022", "营收": 110}]
        plan = normalize_spec({"data": records})
        assert plan.kwargs["categories_col"] == "年份"

    def test_columns_rows(self):
        data = {"columns": ["年份", "营收"], "rows": [["2021", 100], ["2022", 110]]}
        plan = normalize_spec({"data": data})
        assert list(plan.df.columns) == ["年份", "营收"]

    def test_external_df_param(self):
        plan = normalize_spec({"series": ["营收"]}, df=pd.DataFrame(COMBO_DATA))
        assert plan.kwargs["categories_col"] == "年份"

    def test_missing_data_is_clear_error(self):
        with pytest.raises(SpecError, match="缺少数据"):
            normalize_spec({"series": ["营收"]})


# ============================================================================
# 报错质量
# ============================================================================

class TestErrors:
    def test_column_typo_did_you_mean(self):
        with pytest.raises(SpecError) as exc:
            normalize_spec({"data": COMBO_DATA, "series": ["利浬"]})
        assert "利润" in str(exc.value)

    def test_errors_collected_not_first_only(self):
        with pytest.raises(SpecError) as exc:
            normalize_spec({
                "data": COMBO_DATA,
                "series": [{"column": "不存在"}, {"column": "营收", "axis": "middle"}],
            })
        assert len(exc.value.errors) >= 2

    def test_unknown_top_field_is_warning_not_error(self):
        plan = normalize_spec({"data": COMBO_DATA, "serie": ["营收"]})
        assert any("serie" in w for w in plan.warnings)

    def test_validate_spec_reports_without_raising(self):
        report = validate_spec({"data": COMBO_DATA, "series": ["利浬"]})
        assert report["ok"] is False
        assert any("利润" in e for e in report["errors"])

    def test_validate_spec_ok_with_summary(self):
        report = validate_spec({"data": COMBO_DATA, "series": ["营收"]})
        assert report["ok"] is True
        assert report["normalized"]["categories"] == "年份"

    def test_unknown_theme_falls_back_with_warning(self):
        # 品牌盘（able_finance 等）由上层注册；引擎独立时 did-you-mean 针对自带方案
        plan = normalize_spec({"data": COMBO_DATA, "style": {"theme": "categorica"}})
        assert any("categorical" in w for w in plan.warnings)


# ============================================================================
# 颜色
# ============================================================================

class TestColors:
    @pytest.mark.parametrize("raw,expected", [
        ("#C00000", "C00000"),
        ("c00000", "C00000"),
        ("#f00", "FF0000"),
        ("red", "C00000"),
        ("红", "C00000"),
    ])
    def test_normalize_color(self, raw, expected):
        assert normalize_color(raw) == expected

    def test_custom_palette(self):
        plan = normalize_spec({
            "data": COMBO_DATA,
            "series": ["营收"],
            "style": {"colors": ["#1B3D6E", "gold"]},
        })
        assert plan.kwargs["style_config"].colors == ["1B3D6E", "C9A84C"]


# ============================================================================
# 端到端渲染
# ============================================================================

class TestRender:
    def test_combo_render_and_roundtrip(self, tmp_path):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "title": "营收与增速",
            "data": COMBO_DATA,
            "series": ["营收", {"column": "增速", "type": "line", "axis": "right", "color": "#C00000"}],
            "layout": {"legend": "top", "y2_axis": {"format": "percent"}},
        })
        assert chart is not None
        buf = _save_and_reload(prs, tmp_path)
        series_config, _df, categories_col, _layout = parse_chart_from_pptx(buf)
        keys = [s["key"] for s in series_config]
        assert keys == ["营收", "增速"]
        assert categories_col == "年份"

    def test_legend_off(self):
        prs, slide = _slide()
        chart = render_chart(slide, {"data": COMBO_DATA, "series": ["营收"], "legend": "none"})
        assert chart.has_legend is False

    def test_waterfall_via_spec(self, tmp_path):
        prs, slide = _slide()
        render_chart(slide, {
            "chart": "瀑布图",
            "data": {
                "阶段": ["期初", "权益", "债券", "汇率", "期末"],
                "贡献": [8.5, 2.1, 1.3, -1.8, 10.1],
            },
            "totals": ["期初", "期末"],
        })
        buf = _save_and_reload(prs, tmp_path)
        result = parse_waterfall_from_pptx(buf)
        assert result.categories_col == "阶段"
        assert result.value_col == "贡献"
        restored = dict(zip(result.df[result.categories_col], result.df[result.value_col]))
        assert restored["期末"] == pytest.approx(10.1)

    def test_waterfall_infers_measure_col(self):
        plan = normalize_spec({
            "chart": "waterfall",
            "data": {
                "阶段": ["期初", "权益", "期末"],
                "贡献": [8.5, 2.1, 10.6],
                "度量": ["total", "relative", "total"],
            },
        })
        assert plan.kwargs["measure_col"] == "度量"

    def test_bubble_via_spec_with_size_inference(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "chart": "bubble",
            "data": {"波动率": [8.1, 9.2], "收益率": [10.5, 12.0], "规模": [50, 80]},
        })
        assert chart is not None

    def test_scatter_explicit_xy(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "chart": "散点图",
            "data": {"波动率": [8.1, 9.2, 7.4], "收益率": [10.5, 12.0, 8.8]},
            "x": "波动率", "y": "收益率", "name": "风险收益",
        })
        assert chart is not None

    def test_quiet_suppresses_engine_prints(self, capsys):
        prs, slide = _slide()
        render_chart(slide, {"data": COMBO_DATA, "series": ["营收"]})
        assert capsys.readouterr().out == ""

    def test_per_series_color_override_lands_in_xml(self):
        prs, slide = _slide()
        chart = render_chart(slide, {
            "data": COMBO_DATA,
            "series": [
                "营收",
                {"column": "增速", "type": "line", "axis": "right", "color": "#A1B2C3"},
            ],
        })
        xml = chart._chartSpace.xml
        assert "A1B2C3" in xml


# ============================================================================
# 文档
# ============================================================================

def test_spec_reference_mentions_all_kinds():
    doc = chart_spec_reference()
    for kind in ("combo", "waterfall", "scatter", "bubble"):
        assert kind in doc

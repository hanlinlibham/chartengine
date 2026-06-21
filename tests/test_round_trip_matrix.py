"""Round-trip safety net: matrix harness across multiple chart_types.

Systematic enforcement of ADR-0004 §5 + ADR-0006 §5 invariants. Earlier
``test_inspect.py`` / ``test_replace.py`` validated contracts with combo
fixtures only; this file uses ``pytest.parametrize`` to apply the same
invariants across N chart_types.

Coverage target (v1, 2026-05-24): ADR-0006 §3 first-batch 3 chart_types
most distinct in OOXML semantics:

- ``combo`` (bar+line, dual axes) — category-based, multi-series, mixed plots
- ``scatter`` — XY-based, no categories
- ``bubble`` — XYZ-based, no categories, with size

Future expansion (PCE-007): area, pie, plain bar, plain line, then
semantic_family.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Callable, List, Optional

import pandas as pd
import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from ablechart import (
    SeriesData,
    create_bubble_chart,
    create_combo_chart,
    create_scatter_chart,
    inspect_pptx_charts,
    replace_pptx_chart_data,
)


# ---------------------------------------------------------------------------
# Family spec — captures (fixture_factory, expected inspect facts, replace input)
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class FamilySpec:
    """All metadata needed to drive the round-trip matrix on one chart type."""

    id: str
    fixture: Callable                       # tmp_path -> pptx_path
    expected_chart_type: str
    expected_categories: int                # 0 for xy / bubble
    expected_series_count: int
    expected_series_names: List[str]
    replace_categories: Optional[List[Any]]
    replace_series: List[SeriesData]
    expected_categories_after: int


# ---------------------------------------------------------------------------
# Fixture factories (engine-authored)
# ---------------------------------------------------------------------------


def _make_combo(tmp_path):
    df = pd.DataFrame({
        "month": ["1月", "2月", "3月"],
        "revenue": [100, 120, 110],
        "growth_pct": [5.0, 8.0, 4.5],
    })
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide, df=df, categories_col="month",
        series_config=[
            {"key": "revenue", "name": "Revenue", "type": "bar", "axis": "primary"},
            {"key": "growth_pct", "name": "Growth %", "type": "line", "axis": "secondary"},
        ],
    )
    out = tmp_path / "combo.pptx"
    prs.save(str(out))
    return str(out)


def _make_scatter(tmp_path):
    df = pd.DataFrame({"vol": [8.0, 9.0, 7.0], "ret": [10.0, 12.0, 8.0]})
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_scatter_chart(slide=slide, df=df, x_col="vol", y_col="ret",
                         series_name="Risk/Return")
    out = tmp_path / "scatter.pptx"
    prs.save(str(out))
    return str(out)


def _make_bubble(tmp_path):
    df = pd.DataFrame({
        "vol": [8.0, 9.0, 7.0],
        "ret": [10.0, 12.0, 8.0],
        "aum": [120.0, 210.0, 95.0],
    })
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_bubble_chart(slide=slide, df=df, x_col="vol", y_col="ret", size_col="aum",
                        series_name="Risk/Return/AUM")
    out = tmp_path / "bubble.pptx"
    prs.save(str(out))
    return str(out)


def _make_native_category_chart(
    tmp_path,
    *,
    chart_type: XL_CHART_TYPE,
    filename: str,
    single_series: bool = False,
):
    """Native python-pptx chart fixture for category-based charts
    (bar / line / area / pie). Bypasses ``ablechart.create_*`` —
    simulates "external author" so the test verifies that inspect+replace
    work without engine-written metadata."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Sales", (100.0, 200.0, 150.0, 250.0))
    if not single_series:
        chart_data.add_series("Costs", (80.0, 130.0, 110.0, 180.0))
    slide.shapes.add_chart(
        chart_type,
        Inches(1), Inches(1), Inches(8), Inches(5),
        chart_data,
    )
    out = tmp_path / filename
    prs.save(str(out))
    return str(out)


def _make_plain_bar(tmp_path):
    return _make_native_category_chart(
        tmp_path, chart_type=XL_CHART_TYPE.BAR_CLUSTERED,
        filename="plain_bar.pptx", single_series=False,
    )


def _make_plain_line(tmp_path):
    return _make_native_category_chart(
        tmp_path, chart_type=XL_CHART_TYPE.LINE,
        filename="plain_line.pptx", single_series=False,
    )


def _make_area(tmp_path):
    return _make_native_category_chart(
        tmp_path, chart_type=XL_CHART_TYPE.AREA,
        filename="area.pptx", single_series=False,
    )


def _make_pie(tmp_path):
    # pie supports only one series — single_series=True
    return _make_native_category_chart(
        tmp_path, chart_type=XL_CHART_TYPE.PIE,
        filename="pie.pptx", single_series=True,
    )


# ---------------------------------------------------------------------------
# Matrix
# ---------------------------------------------------------------------------


FAMILY_MATRIX: List[FamilySpec] = [
    FamilySpec(
        id="combo",
        fixture=_make_combo,
        expected_chart_type="combo",
        expected_categories=3,
        expected_series_count=2,
        expected_series_names=["Revenue", "Growth %"],
        replace_categories=["Q1", "Q2", "Q3", "Q4"],
        replace_series=[
            SeriesData(name="Revenue", values=[200.0, 250.0, 220.0, 280.0]),
            SeriesData(name="Growth %", values=[10.0, 12.0, 8.0, 15.0]),
        ],
        expected_categories_after=4,
    ),
    FamilySpec(
        id="scatter",
        fixture=_make_scatter,
        expected_chart_type="scatter",
        expected_categories=0,
        expected_series_count=1,
        expected_series_names=["Risk/Return"],
        replace_categories=None,
        replace_series=[
            SeriesData(
                name="Risk/Return",
                values=[15.0, 20.0, 25.0, 30.0],
                x_values=[5.0, 7.0, 9.0, 11.0],
            ),
        ],
        expected_categories_after=0,
    ),
    FamilySpec(
        id="bubble",
        fixture=_make_bubble,
        expected_chart_type="bubble",
        expected_categories=0,
        expected_series_count=1,
        expected_series_names=["Risk/Return/AUM"],
        replace_categories=None,
        replace_series=[
            SeriesData(
                name="Risk/Return/AUM",
                values=[15.0, 25.0],
                x_values=[5.0, 10.0],
                size_values=[100.0, 200.0],
            ),
        ],
        expected_categories_after=0,
    ),
    # ---- v2 expansion 2026-05-24: native (non-engine) fixtures, full first batch ----
    FamilySpec(
        id="plain_bar",
        fixture=_make_plain_bar,
        expected_chart_type="bar",
        expected_categories=4,
        expected_series_count=2,
        expected_series_names=["Sales", "Costs"],
        replace_categories=["A", "B", "C"],
        replace_series=[
            SeriesData(name="Revenue", values=[10.0, 20.0, 30.0]),
            SeriesData(name="Costs", values=[5.0, 15.0, 25.0]),
        ],
        expected_categories_after=3,
    ),
    FamilySpec(
        id="plain_line",
        fixture=_make_plain_line,
        expected_chart_type="line",
        expected_categories=4,
        expected_series_count=2,
        expected_series_names=["Sales", "Costs"],
        replace_categories=["Jan", "Feb", "Mar", "Apr", "May"],
        replace_series=[
            SeriesData(name="Active", values=[1.0, 2.0, 3.0, 4.0, 5.0]),
            SeriesData(name="Inactive", values=[5.0, 4.0, 3.0, 2.0, 1.0]),
        ],
        expected_categories_after=5,
    ),
    FamilySpec(
        id="area",
        fixture=_make_area,
        expected_chart_type="area",
        expected_categories=4,
        expected_series_count=2,
        expected_series_names=["Sales", "Costs"],
        replace_categories=["W1", "W2"],
        replace_series=[
            SeriesData(name="Inflow", values=[10.0, 20.0]),
            SeriesData(name="Outflow", values=[5.0, 8.0]),
        ],
        expected_categories_after=2,
    ),
    FamilySpec(
        id="pie",
        fixture=_make_pie,
        expected_chart_type="pie",
        expected_categories=4,
        expected_series_count=1,
        expected_series_names=["Sales"],
        replace_categories=["EquityA", "EquityB", "Bond", "Cash"],
        replace_series=[
            SeriesData(name="Allocation", values=[40.0, 30.0, 20.0, 10.0]),
        ],
        expected_categories_after=4,
    ),
]


_PARAMS = pytest.mark.parametrize(
    "spec", FAMILY_MATRIX, ids=[s.id for s in FAMILY_MATRIX],
)


# ---------------------------------------------------------------------------
# Round-trip contracts (one parametrized test = one column in the matrix)
# ---------------------------------------------------------------------------


@_PARAMS
def test_create_then_inspect_facts_consistent(tmp_path, spec: FamilySpec):
    """create → save → inspect: inventory matches the create-time input.

    Enforces ADR-0006 §1 inventory accuracy on engine-authored charts.
    """
    pptx_path = spec.fixture(tmp_path)
    inv = inspect_pptx_charts(pptx_path)

    assert len(inv) == 1, f"{spec.id}: expected 1 chart, got {len(inv)}"
    item = inv[0]
    assert item.chart_type == spec.expected_chart_type
    assert item.category_count == spec.expected_categories
    assert item.series_count == spec.expected_series_count
    for name in spec.expected_series_names:
        assert name in item.series_names, f"{spec.id}: '{name}' not in {item.series_names}"
    assert item.has_embedded_workbook is True
    assert item.replaceable is True


@_PARAMS
def test_create_replace_preserves_identity(tmp_path, spec: FamilySpec):
    """create → replace: shape_id + chart_type preserved (ADR-0006 §5.3/§5.4)."""
    input_pptx = spec.fixture(tmp_path)
    inv_before = inspect_pptx_charts(input_pptx)
    shape_id_before = inv_before[0].selector.shape_id
    chart_part_before = inv_before[0].selector.chart_part

    output_pptx = str(tmp_path / "out.pptx")
    result = replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv_before[0].selector,
        categories=spec.replace_categories,
        series=spec.replace_series,
    )
    assert result.status == "ok", f"{spec.id}: {result.error_code} / {result.error_detail}"

    inv_after = inspect_pptx_charts(output_pptx)
    assert inv_after[0].selector.shape_id == shape_id_before, \
        f"{spec.id}: shape_id changed across replace"
    assert inv_after[0].selector.chart_part == chart_part_before, \
        f"{spec.id}: chart_part changed across replace"
    assert inv_after[0].chart_type == spec.expected_chart_type, \
        f"{spec.id}: chart_type changed across replace"


@_PARAMS
def test_create_replace_then_inspect_reflects_new_data(tmp_path, spec: FamilySpec):
    """create → replace → inspect: new categories / series visible (ADR-0006 §5.5/§5.7)."""
    input_pptx = spec.fixture(tmp_path)
    inv_before = inspect_pptx_charts(input_pptx)

    output_pptx = str(tmp_path / "out.pptx")
    replace_pptx_chart_data(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        selector=inv_before[0].selector,
        categories=spec.replace_categories,
        series=spec.replace_series,
    )

    inv_after = inspect_pptx_charts(output_pptx)
    item = inv_after[0]
    assert item.category_count == spec.expected_categories_after, \
        f"{spec.id}: category_count expected {spec.expected_categories_after}, got {item.category_count}"
    assert item.series_count == len(spec.replace_series)
    for s in spec.replace_series:
        assert s.name in item.series_names, \
            f"{spec.id}: replaced series '{s.name}' not in {item.series_names}"

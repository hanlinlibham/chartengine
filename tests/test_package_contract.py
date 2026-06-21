import pandas as pd
import pytest
from pptx import Presentation

from pptchartengine import (
    ATTRIBUTION_DECOMPOSITION_FAMILY,
    AWARD_TIMELINE_PANEL_FAMILY,
    BIWEEKLY_TICKS,
    DAILY_TICKS,
    DUAL_CHART_PANEL_FAMILY,
    FACTOR_ATTRIBUTION_PANEL_FAMILY,
    HEATMAP_MATRIX_FAMILY,
    HOLDING_DETAIL_FAMILY,
    MANAGER_TIMELINE_PROFILE_FAMILY,
    MONTHLY_TICKS,
    PERFORMANCE_COMPARE_FAMILY,
    QUARTERLY_TICKS,
    REGIME_TABLE_PANEL_FAMILY,
    RANKED_TILE_MATRIX_FAMILY,
    SELECTION_TIMING_GRID_FAMILY,
    TABLE_PLUS_CHART_COMPOSITE_FAMILY,
    WEEKLY_TICKS,
    YEARLY_TICKS,
    CHART_PRESET_FUNCTIONS,
    ChartLayoutConfig,
    ChartParser,
    DateAxisConfig,
    FINANCE_PRESET_FUNCTIONS,
    RANGE_SNAPSHOT_PRESET_FUNCTIONS,
    VERTICAL_VALUATION_PRESET_FUNCTIONS,
    StyleConfig,
    WaterfallParseResult,
    ScatterParseResult,
    RangeSnapshotParseResult,
    build_range_snapshot_preset,
    create_distribution_history_chart,
    create_dual_chart_panel,
    create_event_timeline_chart,
    create_factor_attribution_panel,
    create_award_timeline_panel,
    create_holding_detail_panel,
    create_performance_compare_chart,
    create_manager_timeline_profile,
    create_regime_table_panel,
    create_heatmap_matrix_chart,
    create_ranked_tile_matrix_chart,
    create_selection_timing_grid,
    create_score_overlay_chart,
    create_semantic_chart,
    create_style_box_chart,
    create_table_plus_chart_composite,
    get_asx200_sector_valuation_snapshot_preset,
    create_bubble_chart,
    create_combo_chart,
    create_range_snapshot_chart,
    create_scatter_chart,
    create_waterfall_chart,
    get_waterfall_spec,
    get_range_snapshot_spec,
    get_chart1_config,
    get_chart4_config,
    get_msci_emu_sector_valuation_snapshot_preset,
    get_msci_japan_sector_valuation_snapshot_preset,
    get_sp500_sector_valuation_snapshot_preset,
    get_vertical_global_valuation_snapshot_preset,
    get_vertical_sector_valuation_snapshot_preset,
    parse_all_charts_from_pptx,
    parse_all_semantic_components_from_pptx,
    parse_bubble_from_pptx,
    parse_chart_from_pptx,
    parse_range_snapshot_chart,
    parse_range_snapshot_from_pptx,
    parse_semantic_component_from_pptx,
    parse_scatter_from_pptx,
    parse_waterfall_from_pptx,
    list_semantic_families,
    get_semantic_chart_spec,
    parse_semantic_chart_from_layout_info,
    prepare_range_snapshot_dataframe,
    prepare_waterfall_dataframe,
    restore_range_snapshot_dataframe,
    restore_waterfall_dataframe,
)
from pptchartengine.date_axis import format_category_label


def _sample_finance_df() -> pd.DataFrame:
    return pd.DataFrame(
        {
            "分类": pd.date_range("2024-01-01", periods=14, freq="D"),
            "沪深300指数(收盘价)": [3500 + i * 10 for i in range(14)],
            "组合收益率（左轴）": [0.01 + i * 0.001 for i in range(14)],
            "组合规模(万元)": [100000 + i * 2000 for i in range(14)],
            "累计收益率": [0.10 + i * 0.002 for i in range(14)],
            "权益仓位（人社部口径）": [0.18 + i * 0.003 for i in range(14)],
            "中债新综合总财富指数(收盘价)": [210.5 + i * 0.8 for i in range(14)],
            "久期": [0.45 + i * 0.02 for i in range(14)],
        }
    )


def test_public_create_and_parse_are_quiet_by_default(tmp_path, capsys):
    df = pd.DataFrame(
        {
            "年份": [2021, 2022, 2023],
            "营收": [100, 120, 130],
            "利润": [10, 14, 16],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_combo_chart(
        slide=slide,
        df=df,
        categories_col="年份",
        series_config=[
            {"key": "营收", "name": "营收", "type": "bar", "axis": "primary"},
            {"key": "利润", "name": "利润", "type": "line", "axis": "secondary"},
        ],
    )
    output = tmp_path / "combo-smoke.pptx"
    prs.save(output)

    parse_chart_from_pptx(output)

    captured = capsys.readouterr()
    assert captured.out == ""


def test_public_contract_exports():
    assert callable(create_combo_chart)
    assert callable(create_performance_compare_chart)
    assert callable(create_distribution_history_chart)
    assert callable(create_dual_chart_panel)
    assert callable(create_score_overlay_chart)
    assert callable(create_style_box_chart)
    assert callable(create_event_timeline_chart)
    assert callable(create_factor_attribution_panel)
    assert callable(create_manager_timeline_profile)
    assert callable(create_award_timeline_panel)
    assert callable(create_holding_detail_panel)
    assert callable(create_regime_table_panel)
    assert callable(create_ranked_tile_matrix_chart)
    assert callable(create_heatmap_matrix_chart)
    assert callable(create_selection_timing_grid)
    assert callable(create_table_plus_chart_composite)
    assert callable(create_semantic_chart)
    assert callable(parse_chart_from_pptx)
    assert callable(parse_all_charts_from_pptx)
    assert callable(parse_semantic_component_from_pptx)
    assert callable(parse_all_semantic_components_from_pptx)
    assert ChartParser.__name__ == "ChartParser"
    assert StyleConfig.__name__ == "StyleConfig"
    assert ChartLayoutConfig.__name__ == "ChartLayoutConfig"
    assert DateAxisConfig.__name__ == "DateAxisConfig"
    assert WaterfallParseResult.__name__ == "WaterfallParseResult"
    assert ScatterParseResult.__name__ == "ScatterParseResult"
    assert RangeSnapshotParseResult.__name__ == "RangeSnapshotParseResult"
    assert callable(build_range_snapshot_preset)
    assert callable(list_semantic_families)
    assert callable(get_semantic_chart_spec)
    assert callable(parse_semantic_chart_from_layout_info)


def test_semantic_family_registry_contains_demo01_families():
    families = list_semantic_families()
    assert PERFORMANCE_COMPARE_FAMILY in families
    assert ATTRIBUTION_DECOMPOSITION_FAMILY in families
    assert DUAL_CHART_PANEL_FAMILY in families
    assert RANKED_TILE_MATRIX_FAMILY in families
    assert HEATMAP_MATRIX_FAMILY in families
    assert TABLE_PLUS_CHART_COMPOSITE_FAMILY in families
    assert FACTOR_ATTRIBUTION_PANEL_FAMILY in families
    assert REGIME_TABLE_PANEL_FAMILY in families
    assert MANAGER_TIMELINE_PROFILE_FAMILY in families
    assert AWARD_TIMELINE_PANEL_FAMILY in families
    assert SELECTION_TIMING_GRID_FAMILY in families
    assert HOLDING_DETAIL_FAMILY in families
    assert families[PERFORMANCE_COMPARE_FAMILY]["renderable"] is True
    assert families["holding_detail"]["renderable"] is True


def test_semantic_combo_chart_metadata_round_trip(tmp_path):
    df = pd.DataFrame(
        {
            "日期": pd.date_range("2025-01-01", periods=5, freq="ME"),
            "基金": [0.01, 0.03, 0.02, 0.04, 0.06],
            "沪深300": [0.00, 0.01, 0.015, 0.02, 0.03],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_performance_compare_chart(
        slide=slide,
        df=df,
        categories_col="日期",
        series_entries=[
            {"key": "基金", "name": "基金", "role": "fund", "type": "line"},
            {"key": "沪深300", "name": "沪深300", "role": "benchmark", "type": "line"},
        ],
        title="绩效对比",
    )
    output = tmp_path / "semantic-performance.pptx"
    prs.save(output)

    _, _, _, layout_info = parse_chart_from_pptx(output)
    spec = get_semantic_chart_spec(layout_info)
    parsed = parse_semantic_chart_from_layout_info(layout_info)

    assert spec["chart_family"] == PERFORMANCE_COMPARE_FAMILY
    assert spec["series_roles"][0]["role"] == "fund"
    assert parsed.family == PERFORMANCE_COMPARE_FAMILY


def test_shape_only_semantic_family_is_discoverable_after_save(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_holding_detail_panel(
        slide=slide,
        headers=["证券", "权重"],
        rows=[["贵州茅台", "9.2%"], ["宁德时代", "5.6%"]],
        subtitle="前十大重仓股",
        summary_text="示例持仓明细",
    )

    output = tmp_path / "holding-detail-panel.pptx"
    prs.save(output)

    components = parse_all_semantic_components_from_pptx(str(output))
    assert len(components) == 1
    component = components[0]
    assert component["source"] == "anchor"
    assert component["family"] == HOLDING_DETAIL_FAMILY
    assert component["metadata"]["headers"] == ["证券", "权重"]
    assert component["metadata"]["rows"][0][0] == "贵州茅台"


def test_composite_semantic_family_identity_survives_reopen(tmp_path):
    df = pd.DataFrame(
        {
            "日期": pd.date_range("2025-01-01", periods=4, freq="ME"),
            "基金": [0.01, 0.02, 0.015, 0.03],
            "沪深300": [0.0, 0.01, 0.012, 0.018],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    create_table_plus_chart_composite(
        slide=slide,
        chart_family=PERFORMANCE_COMPARE_FAMILY,
        chart_kwargs={
            "df": df,
            "categories_col": "日期",
            "series_entries": [
                {"key": "基金", "name": "基金", "role": "fund", "type": "line"},
                {"key": "沪深300", "name": "沪深300", "role": "benchmark", "type": "line"},
            ],
            "title": "绩效对比",
        },
        headers=["指标", "数值"],
        rows=[["近1月", "3.0%"], ["近3月", "5.6%"]],
        left_title="左图",
        right_title="右表",
    )

    output = tmp_path / "table-plus-chart-composite.pptx"
    prs.save(output)

    component = parse_semantic_component_from_pptx(str(output))
    components = parse_all_semantic_components_from_pptx(str(output))
    assert len(components) == 1
    assert component["source"] == "anchor"
    assert component["family"] == TABLE_PLUS_CHART_COMPOSITE_FAMILY
    assert component["metadata"]["nested_chart_family"] == PERFORMANCE_COMPARE_FAMILY
    assert component["metadata"]["rows"][1][0] == "近3月"


def test_date_axis_presets_exposed():
    presets = [
        DAILY_TICKS,
        WEEKLY_TICKS,
        BIWEEKLY_TICKS,
        MONTHLY_TICKS,
        QUARTERLY_TICKS,
        YEARLY_TICKS,
    ]
    assert all(isinstance(preset, DateAxisConfig) for preset in presets)


def test_format_category_label_drops_time_of_day_for_date_like_values():
    assert format_category_label("2026-04-21 00:00:00") == "2026-04-21"
    assert format_category_label("2026-04-21T15:30:00") == "2026-04-21"
    assert format_category_label("2026-04", "yyyy-mm") == "2026-04"


def test_financial_chart1_preset_shape():
    config = get_chart1_config(_sample_finance_df())

    assert config["categories_col"] == "分类"
    assert [series["type"] for series in config["series_config"]] == ["bar", "line"]
    assert [series["axis"] for series in config["series_config"]] == ["secondary", "primary"]
    assert config["style_config"].color_scheme == "aim00"
    assert config["layout_config"].title == "组合2024年以来收益率走势图"
    assert config["layout_config"].secondary_value_axis_config.number_format == "#,##0"
    assert config["layout_config"].date_axis_config.number_format == "yyyy/mm"


def test_financial_chart4_preset_uses_dynamic_secondary_axis_range():
    df = _sample_finance_df()
    config = get_chart4_config(df)
    axis = config["layout_config"].secondary_value_axis_config
    bond_series = df["中债新综合总财富指数(收盘价)"]
    bond_range = bond_series.max() - bond_series.min()
    bond_unit = int(bond_range / 6)

    assert axis.min_value == int(bond_series.min() - (bond_range / 6))
    assert axis.max_value == int(bond_series.max() + (bond_range / 6))
    assert axis.major_unit == bond_unit


def test_finance_registry_aliases_match():
    assert FINANCE_PRESET_FUNCTIONS is CHART_PRESET_FUNCTIONS
    assert set(CHART_PRESET_FUNCTIONS) == {
        "图表1 - 当年收益率走势",
        "图表2 - 成立以来收益率",
        "图表3 - 权益仓位走势",
        "图表4 - 久期走势",
    }


def test_range_snapshot_preset_registry_and_defaults():
    global_df = pd.DataFrame(
        {
            "market": ["China", "Japan"],
            "range_min": [7.8, 10.9],
            "range_max": [24.5, 41.2],
            "average": [11.2, 15.8],
            "current": [13.8, 15.8],
        }
    )
    sector_df = pd.DataFrame(
        {
            "sector": ["Energy", "Info Tech"],
            "range_min": [1.5, 12.8],
            "range_max": [50.0, 46.2],
            "average": [1.7, 22.0],
            "current": [15.2, 29.8],
        }
    )

    global_cfg = get_vertical_global_valuation_snapshot_preset(global_df)
    sector_cfg = get_vertical_sector_valuation_snapshot_preset(
        sector_df,
        title="MSCI EMU valuation",
    )

    assert set(RANGE_SNAPSHOT_PRESET_FUNCTIONS) == {
        "估值快照 - 全市场竖版",
        "估值快照 - 行业竖版带轴断裂",
    }
    assert set(VERTICAL_VALUATION_PRESET_FUNCTIONS) == {
        "ASX 200 valuation",
        "S&P 500 valuation",
        "MSCI EMU valuation",
        "MSCI Japan valuation",
    }
    assert global_cfg["orientation"] == "vertical"
    assert global_cfg["axis_break"] is None
    assert sector_cfg["orientation"] == "vertical"
    assert sector_cfg["axis_break"]["value"] == 30.0
    assert sector_cfg["axis_break"]["tick_values"][-1] >= 45.0


def test_page_specific_vertical_valuation_presets():
    df = pd.DataFrame(
        {
            "sector": ["Energy", "Info Tech", "Materials"],
            "range_min": [1.5, 12.8, 8.6],
            "range_max": [50.0, 46.2, 17.2],
            "average": [1.7, 22.0, 15.6],
            "current": [15.2, 29.8, 19.8],
        }
    )

    asx = get_asx200_sector_valuation_snapshot_preset(df)
    spx = get_sp500_sector_valuation_snapshot_preset(df)
    emu = get_msci_emu_sector_valuation_snapshot_preset(df)
    japan = get_msci_japan_sector_valuation_snapshot_preset(df)

    assert asx["title"] == "ASX 200 valuation"
    assert asx["axis_break"]["value"] == 50.0
    assert asx["axis_break"]["tick_values"] == [0.0, 10.0, 20.0, 30.0, 40.0, 50.0]

    assert spx["title"] == "S&P 500 valuation"
    assert spx["axis_break"]["value"] == 30.0
    assert spx["axis_break"]["tick_values"][-1] == 50.0

    assert emu["title"] == "MSCI EMU valuation"
    assert emu["axis_break"] is None

    assert japan["title"] == "MSCI Japan valuation"
    assert japan["subtitle"] == "Trailing price-to-book ratio"


def test_prepare_waterfall_dataframe():
    df = pd.DataFrame(
        {
            "项目": ["期初", "权益贡献", "债券贡献", "汇率拖累", "期末"],
            "值": [100, 20, 10, -15, 115],
            "measure": ["total", "relative", "relative", "relative", "total"],
        }
    )

    waterfall = prepare_waterfall_dataframe(df, "项目", "值", measure_col="measure")
    assert list(waterfall.columns) == ["项目", "__base__", "__increase__", "__decrease__", "__total__"]
    assert waterfall.iloc[0]["__total__"] == 100
    assert waterfall.iloc[1]["__base__"] == 100
    assert waterfall.iloc[1]["__increase__"] == 20
    assert waterfall.iloc[3]["__base__"] == 115
    assert waterfall.iloc[3]["__decrease__"] == 15
    assert waterfall.iloc[4]["__total__"] == 115


def test_prepare_waterfall_dataframe_supports_total_categories_without_measure_col():
    df = pd.DataFrame(
        {
            "项目": ["起点", "经营改善", "融资拖累", "终点"],
            "值": [80, 25, -10, 95],
        }
    )

    waterfall = prepare_waterfall_dataframe(
        df,
        "项目",
        "值",
        total_categories=["起点", "终点"],
    )

    assert waterfall.iloc[0]["__total__"] == 80
    assert waterfall.iloc[1]["__base__"] == 80
    assert waterfall.iloc[2]["__base__"] == 95
    assert waterfall.iloc[2]["__decrease__"] == 10
    assert waterfall.iloc[3]["__total__"] == 95


def test_prepare_range_snapshot_dataframe():
    df = pd.DataFrame(
        {
            "sector": ["China", "EM", "U.S."],
            "range_min": [7.8, 7.5, 10.2],
            "range_max": [24.5, 18.1, 22.7],
        }
    )

    snapshot = prepare_range_snapshot_dataframe(df, "sector", "range_min", "range_max")
    assert list(snapshot.columns) == ["sector", "__base__", "__range__"]
    assert snapshot.iloc[0]["__base__"] == 7.8
    assert snapshot.iloc[0]["__range__"] == pytest.approx(16.7)
    assert snapshot.iloc[2]["__range__"] == pytest.approx(12.5)


def test_round_trip_recovers_categories_col_and_series_keys(tmp_path):
    df = pd.DataFrame(
        {
            "年份": [2021, 2022, 2023, 2024],
            "revenue": [100.0, 110.0, 120.0, 140.0],
            "profit": [10.0, 12.5, 15.0, 18.0],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide,
        df=df,
        categories_col="年份",
        series_config=[
            {"key": "revenue", "name": "营收(亿元)", "type": "bar", "axis": "primary"},
            {"key": "profit", "name": "净利润(亿元)", "type": "line", "axis": "secondary"},
        ],
    )

    output = tmp_path / "roundtrip.pptx"
    prs.save(output)

    series_config, parsed_df, categories_col, layout_info = parse_chart_from_pptx(str(output))

    assert categories_col == "年份"
    assert [series["key"] for series in series_config] == ["revenue", "profit"]
    assert [series["name"] for series in series_config] == ["营收(亿元)", "净利润(亿元)"]
    assert [series["axis"] for series in series_config] == ["primary", "secondary"]
    assert list(parsed_df.columns) == ["年份", "营收(亿元)", "净利润(亿元)"]
    assert isinstance(layout_info, dict)


def test_round_trip_recovers_categories_col_with_ascii_header(tmp_path):
    # Pins the metadata fallback for the case where the embedded workbook
    # leaves the first header cell blank — the failure mode previously
    # surfaced only because every other combo bar/line round-trip test
    # happened to use a non-ASCII categories_col.
    df = pd.DataFrame(
        {
            "year": [2021, 2022, 2023, 2024],
            "revenue": [100.0, 110.0, 120.0, 140.0],
            "profit": [10.0, 12.5, 15.0, 18.0],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide,
        df=df,
        categories_col="year",
        series_config=[
            {"key": "revenue", "name": "Revenue", "type": "bar", "axis": "primary"},
            {"key": "profit", "name": "Profit", "type": "line", "axis": "secondary"},
        ],
    )

    output = tmp_path / "roundtrip-ascii.pptx"
    prs.save(output)

    series_config, parsed_df, categories_col, _ = parse_chart_from_pptx(str(output))

    assert categories_col == "year"
    assert [series["key"] for series in series_config] == ["revenue", "profit"]
    assert list(parsed_df.columns) == ["year", "Revenue", "Profit"]


def test_round_trip_recovers_stacked_grouping(tmp_path):
    df = pd.DataFrame(
        {
            "年份": [2021, 2022, 2023],
            "domestic": [40.0, 42.0, 44.0],
            "overseas": [60.0, 58.0, 56.0],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide,
        df=df,
        categories_col="年份",
        series_config=[
            {"key": "domestic", "name": "国内收入", "type": "bar", "axis": "primary", "grouping": "stacked"},
            {"key": "overseas", "name": "海外收入", "type": "bar", "axis": "primary", "grouping": "stacked"},
        ],
    )

    output = tmp_path / "stacked-roundtrip.pptx"
    prs.save(output)

    series_config, parsed_df, categories_col, layout_info = parse_chart_from_pptx(str(output))

    assert categories_col == "年份"
    assert all(series["grouping"] == "stacked" for series in series_config)
    assert [series["key"] for series in series_config] == ["domestic", "overseas"]


def test_round_trip_recovers_percent_stacked_grouping(tmp_path):
    df = pd.DataFrame(
        {
            "年份": [2021, 2022, 2023],
            "equity": [55.0, 52.0, 50.0],
            "bond": [45.0, 48.0, 50.0],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide,
        df=df,
        categories_col="年份",
        series_config=[
            {"key": "equity", "name": "权益占比", "type": "bar", "axis": "primary", "grouping": "percent_stacked"},
            {"key": "bond", "name": "固收占比", "type": "bar", "axis": "primary", "grouping": "percent_stacked"},
        ],
    )

    output = tmp_path / "percent-stacked-roundtrip.pptx"
    prs.save(output)

    series_config, parsed_df, categories_col, layout_info = parse_chart_from_pptx(str(output))

    assert categories_col == "年份"
    assert all(series["grouping"] == "percent_stacked" for series in series_config)
    assert [series["key"] for series in series_config] == ["equity", "bond"]


def test_create_waterfall_chart_round_trip(tmp_path):
    df = pd.DataFrame(
        {
            "项目": ["期初", "权益贡献", "债券贡献", "汇率拖累", "期末"],
            "值": [100, 20, 10, -15, 115],
            "measure": ["total", "relative", "relative", "relative", "total"],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_waterfall_chart(
        slide=slide,
        df=df,
        categories_col="项目",
        value_col="值",
        measure_col="measure",
    )

    output = tmp_path / "waterfall.pptx"
    prs.save(output)
    series_config, parsed_df, categories_col, layout_info = parse_chart_from_pptx(str(output))

    assert categories_col == "项目"
    assert all(series["grouping"] == "stacked" for series in series_config)
    assert [series["key"] for series in series_config] == [
        "__base__",
        "__increase__",
        "__decrease__",
        "__total__",
    ]
    assert parsed_df.shape == (5, 5)
    assert layout_info["chart_family"] == "waterfall"
    assert get_waterfall_spec(layout_info)["value_col"] == "值"

    restored = restore_waterfall_dataframe(parsed_df, layout_info)
    assert restored.to_dict(orient="list") == df.to_dict(orient="list")


def test_parse_waterfall_from_pptx_returns_semantic_dataframe(tmp_path):
    df = pd.DataFrame(
        {
            "项目": ["期初", "股票贡献", "债券贡献", "期末"],
            "值": [100, 15, -5, 110],
            "measure": ["total", "relative", "relative", "absolute"],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_waterfall_chart(
        slide=slide,
        df=df,
        categories_col="项目",
        value_col="值",
        measure_col="measure",
    )

    output = tmp_path / "waterfall-semantic-roundtrip.pptx"
    prs.save(output)

    result = parse_waterfall_from_pptx(str(output))

    assert result.categories_col == "项目"
    assert result.value_col == "值"
    assert result.measure_col == "measure"
    assert result.df.to_dict(orient="list") == df.to_dict(orient="list")
    assert result.raw_chart_df.columns.tolist() == [
        "项目",
        "__base__",
        "__increase__",
        "__decrease__",
        "__total__",
    ]

    assert result.series_config[0]["key"] == "__base__"
    assert result.layout_info["chart_family"] == "waterfall"


def test_create_range_snapshot_chart_round_trip(tmp_path):
    df = pd.DataFrame(
        {
            "market": ["China", "EM", "Europe", "U.S."],
            "range_min": [7.8, 7.5, 8.2, 10.2],
            "range_max": [24.5, 18.1, 19.7, 22.7],
            "average": [11.2, 11.7, 13.1, 16.0],
            "current": [13.8, 14.0, 14.6, 22.7],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_range_snapshot_chart(
        slide=slide,
        df=df,
        categories_col="market",
        min_col="range_min",
        max_col="range_max",
        average_col="average",
        current_col="current",
    )

    output = tmp_path / "range-snapshot.pptx"
    prs.save(output)

    result = parse_range_snapshot_from_pptx(str(output))
    assert result.categories_col == "market"
    assert result.min_col == "range_min"
    assert result.max_col == "range_max"
    assert result.average_col == "average"
    assert result.current_col == "current"
    assert result.orientation == "vertical"
    assert result.df.to_dict(orient="list") == df.to_dict(orient="list")
    assert result.raw_chart_df.columns.tolist() == ["market", "__base__", "__range__"]
    assert result.series_config[0]["key"] == "__base__"
    assert result.layout_info["chart_family"] == "range_snapshot"
    assert get_range_snapshot_spec(result.layout_info)["current_col"] == "current"


def test_create_horizontal_range_snapshot_chart_round_trip(tmp_path):
    df = pd.DataFrame(
        {
            "market": ["China", "EM", "Europe", "U.S."],
            "range_min": [7.8, 7.5, 8.2, 10.2],
            "range_max": [24.5, 18.1, 19.7, 22.7],
            "average": [11.2, 11.7, 13.1, 16.0],
            "current": [13.8, 14.0, 14.6, 22.7],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_range_snapshot_chart(
        slide=slide,
        df=df,
        categories_col="market",
        min_col="range_min",
        max_col="range_max",
        average_col="average",
        current_col="current",
        orientation="horizontal",
    )

    output = tmp_path / "range-snapshot-horizontal.pptx"
    prs.save(output)

    result = parse_range_snapshot_from_pptx(str(output))
    assert result.orientation == "horizontal"
    assert result.df.to_dict(orient="list") == df.to_dict(orient="list")
    assert get_range_snapshot_spec(result.layout_info)["orientation"] == "horizontal"


def test_range_snapshot_axis_break_metadata_and_overlays():
    df = pd.DataFrame(
        {
            "market": ["China", "Japan"],
            "range_min": [7.8, 10.9],
            "range_max": [24.5, 41.2],
            "average": [11.2, 15.8],
            "current": [13.8, 15.8],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_range_snapshot_chart(
        slide=slide,
        df=df,
        categories_col="market",
        min_col="range_min",
        max_col="range_max",
        average_col="average",
        current_col="current",
        orientation="horizontal",
        show_average_ticks=False,
        show_current_markers=False,
        show_current_labels=False,
        axis_break={"value": 30.0, "categories": ["Japan"]},
    )

    assert len(slide.shapes) > 5
    chart = next(shape.chart for shape in slide.shapes if getattr(shape, "has_chart", False))
    result = parse_range_snapshot_chart(chart)
    spec = get_range_snapshot_spec(result.layout_info)
    assert spec["axis_break"]["value"] == 30.0
    assert spec["axis_break"]["categories"] == ["Japan"]


def test_restore_range_snapshot_dataframe_falls_back_to_base_range_without_metadata():
    raw_df = pd.DataFrame(
        {
            "market": ["China", "EM"],
            "__base__": [7.8, 7.5],
            "__range__": [16.7, 10.6],
        }
    )
    restored = restore_range_snapshot_dataframe(raw_df, layout_info=None)
    assert restored["market"].tolist() == ["China", "EM"]
    assert restored["min"].tolist() == [7.8, 7.5]
    assert restored["max"].tolist() == [24.5, 18.1]
    assert restored["average"].isna().all()
    assert restored["current"].isna().all()


def test_create_scatter_chart_round_trip(tmp_path):
    df = pd.DataFrame(
        {
            "volatility": [8.1, 9.2, 7.4],
            "return": [10.5, 12.0, 8.8],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_scatter_chart(
        slide=slide,
        df=df,
        x_col="volatility",
        y_col="return",
        series_name="Risk/Return",
    )
    output = tmp_path / "scatter.pptx"
    prs.save(output)

    result = parse_scatter_from_pptx(str(output))
    assert result.chart_family == "scatter"
    assert result.x_col == "volatility"
    assert result.y_col == "return"
    assert result.size_col is None
    assert result.df.to_dict(orient="list") == df.to_dict(orient="list")


def test_create_bubble_chart_round_trip(tmp_path):
    df = pd.DataFrame(
        {
            "volatility": [8.1, 9.2, 7.4],
            "return": [10.5, 12.0, 8.8],
            "aum": [120.0, 210.0, 95.0],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_bubble_chart(
        slide=slide,
        df=df,
        x_col="volatility",
        y_col="return",
        size_col="aum",
        series_name="Risk/Return/AUM",
    )
    output = tmp_path / "bubble.pptx"
    prs.save(output)

    result = parse_bubble_from_pptx(str(output))
    assert result.chart_family == "bubble"
    assert result.x_col == "volatility"
    assert result.y_col == "return"
    assert result.size_col == "aum"
    assert result.df.to_dict(orient="list") == df.to_dict(orient="list")


def test_create_combo_chart_supports_scatter_round_trip(tmp_path):
    df = pd.DataFrame(
        {
            "volatility": [8.1, 9.2, 7.4],
            "return_a": [10.5, 12.0, 8.8],
            "return_b": [9.8, 11.4, 8.1],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide,
        df=df,
        categories_col="volatility",
        series_config=[
            {"key": "return_a", "name": "Risk/Return A", "type": "scatter", "axis": "primary", "x_key": "volatility"},
            {"key": "return_b", "name": "Risk/Return B", "type": "scatter", "axis": "primary", "x_key": "volatility"},
        ],
    )
    output = tmp_path / "combo-scatter.pptx"
    prs.save(output)

    series_config, parsed_df, categories_col, _layout_info = parse_chart_from_pptx(str(output))
    assert categories_col == "volatility"
    assert [series["type"] for series in series_config] == ["scatter", "scatter"]
    assert [series["x_key"] for series in series_config] == ["volatility", "volatility"]
    assert parsed_df.to_dict(orient="list") == df.to_dict(orient="list")


def test_create_combo_chart_supports_bubble_round_trip(tmp_path):
    df = pd.DataFrame(
        {
            "volatility": [8.1, 9.2, 7.4],
            "return": [10.5, 12.0, 8.8],
            "aum": [120.0, 210.0, 95.0],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_combo_chart(
        slide=slide,
        df=df,
        categories_col="volatility",
        series_config=[
            {
                "key": "return",
                "name": "Risk/Return/AUM",
                "type": "bubble",
                "axis": "primary",
                "x_key": "volatility",
                "size_key": "aum",
            }
        ],
    )
    output = tmp_path / "combo-bubble.pptx"
    prs.save(output)

    series_config, parsed_df, categories_col, _layout_info = parse_chart_from_pptx(str(output))
    assert categories_col == "volatility"
    assert series_config[0]["type"] == "bubble"
    assert series_config[0]["x_key"] == "volatility"
    assert series_config[0]["size_key"] == "aum"
    assert parsed_df.to_dict(orient="list") == df.to_dict(orient="list")


def test_create_combo_chart_rejects_xy_and_category_mix():
    df = pd.DataFrame(
        {
            "volatility": [8.1, 9.2, 7.4],
            "return": [10.5, 12.0, 8.8],
            "aum": [120.0, 210.0, 95.0],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    with pytest.raises(ValueError, match="混搭"):
        create_combo_chart(
            slide=slide,
            df=df,
            categories_col="volatility",
            series_config=[
                {"key": "return", "name": "Scatter", "type": "scatter", "axis": "primary", "x_key": "volatility"},
                {"key": "aum", "name": "Column", "type": "bar", "axis": "primary"},
            ],
        )

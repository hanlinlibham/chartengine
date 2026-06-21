"""Contract tests for the semantic-family dispatcher and registry.

These close two pre-0.1.0 coverage gaps:

1. ``create_semantic_chart`` (the string-dispatched entry point) was only ever
   asserted ``callable`` — never actually executed. Here we route a real family
   through it end-to-end and round-trip the result.
2. The registry / FAMILY constants / ``list_semantic_families`` views could
   drift out of sync silently. Here we pin them to each other.
"""
import pandas as pd
from pptx import Presentation

import pptchartengine as pce
from pptchartengine import (
    create_semantic_chart,
    list_semantic_families,
    get_semantic_chart_spec,
    parse_semantic_chart_from_layout_info,
    parse_chart_from_pptx,
    SEMANTIC_FAMILY_REGISTRY,
    PERFORMANCE_COMPARE_FAMILY,
)


def test_create_semantic_chart_dispatcher_routes_and_round_trips(tmp_path):
    """The string dispatcher renders a real chart and survives a round-trip."""
    df = pd.DataFrame(
        {
            "日期": pd.date_range("2025-01-01", periods=5, freq="ME"),
            "基金": [0.01, 0.03, 0.02, 0.04, 0.06],
            "沪深300": [0.00, 0.01, 0.015, 0.02, 0.03],
        }
    )
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Route through the dispatcher by family string (not the direct create_* fn).
    create_semantic_chart(
        slide,
        PERFORMANCE_COMPARE_FAMILY,
        df=df,
        categories_col="日期",
        series_entries=[
            {"key": "基金", "name": "基金", "role": "fund", "type": "line"},
            {"key": "沪深300", "name": "沪深300", "role": "benchmark", "type": "line"},
        ],
        title="绩效对比",
    )

    output = tmp_path / "dispatched-performance.pptx"
    prs.save(output)

    _, _, _, layout_info = parse_chart_from_pptx(output)
    spec = get_semantic_chart_spec(layout_info)
    parsed = parse_semantic_chart_from_layout_info(layout_info)

    assert spec["chart_family"] == PERFORMANCE_COMPARE_FAMILY
    assert parsed.family == PERFORMANCE_COMPARE_FAMILY


def test_create_semantic_chart_rejects_unknown_family():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    try:
        create_semantic_chart(slide, "not_a_real_family", df=pd.DataFrame())
    except (KeyError, ValueError):
        return
    raise AssertionError("dispatcher should reject an unknown family")


def test_semantic_registry_views_stay_in_sync():
    """Registry, list_semantic_families(), and *_FAMILY constants must agree."""
    registry_keys = set(SEMANTIC_FAMILY_REGISTRY)
    listed = list_semantic_families()

    # list_semantic_families() exposes exactly the registry families.
    assert set(listed) == registry_keys

    # Every *_FAMILY string constant points at a real registry family.
    family_constants = {
        getattr(pce, name)
        for name in dir(pce)
        if name.endswith("_FAMILY") and isinstance(getattr(pce, name), str)
    }
    assert family_constants == registry_keys

    # Each entry carries the documented metadata shape.
    for family, meta in SEMANTIC_FAMILY_REGISTRY.items():
        assert {"renderable", "base_geometry", "description"} <= set(meta), family
        assert isinstance(meta["description"], str) and meta["description"]

    # renderable_only is a subset filter, never inventing families.
    renderable = set(list_semantic_families(renderable_only=True))
    assert renderable <= registry_keys

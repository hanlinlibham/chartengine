"""Value-axis resolution by axId, not axPos (issue #9, Gap B).

The primary/secondary value axes must be identified by the axId the series'
plot group references (document order), so number_format / scale land on the
right axis even when axPos is non-standard (e.g. horizontal bar combos put the
value axes at 'b'/'t', not 'l'/'r').
"""
import pandas as pd
from pptx import Presentation

from ablechart import (
    create_combo_chart,
    ChartLayoutConfig,
    ValueAxisConfig,
)
from ablechart.oxml.axes import resolve_value_axes

_C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"


def _combo_with_secondary(slide):
    df = pd.DataFrame({"FY": ["2022", "2023", "2024"],
                       "Revenue": [328.6, 400.9, 362.0],
                       "Margin": [9.4, 11.0, 14.0]})
    return create_combo_chart(
        slide=slide, df=df, categories_col="FY",
        series_config=[
            {"key": "Revenue", "name": "Revenue", "type": "column", "axis": "primary"},
            {"key": "Margin", "name": "Margin", "type": "line", "axis": "secondary"},
        ],
        layout_config=ChartLayoutConfig(
            value_axis_config=ValueAxisConfig(number_format="0.0"),
            secondary_value_axis_config=ValueAxisConfig(number_format='0"%"'),
        ),
    )


def _axid(ax):
    return ax.find(f"{_C}axId").get("val")


def _fmt(ax):
    nf = ax.find(f"{_C}numFmt")
    return nf.get("formatCode") if nf is not None else None


def test_resolve_returns_distinct_primary_and_secondary():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart = _combo_with_secondary(slide)

    primary, secondary = resolve_value_axes(chart._element)
    assert primary is not None and secondary is not None
    assert _axid(primary) != _axid(secondary)
    # primary is the first valAx in document order (builder convention)
    all_val = chart._element.xpath(".//c:valAx")
    assert primary is all_val[0]


def test_number_format_lands_on_correct_axis():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart = _combo_with_secondary(slide)

    primary, secondary = resolve_value_axes(chart._element)
    assert _fmt(primary) == "0.0"
    assert _fmt(secondary) == '0"%"'


def test_resolution_is_independent_of_axpos():
    """Simulate a horizontal layout: flip axPos to 'b'/'t'. axId resolution must
    still find both axes, whereas the old axPos=='r' lookup would find none."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart = _combo_with_secondary(slide)

    primary, secondary = resolve_value_axes(chart._element)
    pri_id, sec_id = _axid(primary), _axid(secondary)

    # rewrite axPos to a horizontal-style layout (no 'r' anywhere)
    primary.find(f"{_C}axPos").set("val", "b")
    secondary.find(f"{_C}axPos").set("val", "t")
    assert not chart._element.xpath(".//c:valAx/c:axPos[@val='r']")  # old heuristic dead

    primary2, secondary2 = resolve_value_axes(chart._element)
    assert _axid(primary2) == pri_id
    assert _axid(secondary2) == sec_id

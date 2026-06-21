"""Render a 4-family chart gallery (all real CATL data, distinct palettes).

Each panel is a real editable PowerPoint chart produced by pptchartengine, then
rendered to PNG. Captions state capabilities as plain facts (no comparisons).

Families shown:
  1. Combo dual-axis  — stacked columns + 2 lines on a second % axis
  2. Waterfall        — FY2024 P&L bridge (total/relative measures, connectors)
  3. Bubble           — margin trajectory, bubble size = revenue
  4. Range snapshot   — margin ranges with min/max/average/current markers

Run:  python scripts/make_gallery.py
Output: docs/assets/gallery.png
"""
from pathlib import Path
import subprocess
import sys

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "src"))

import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont

from pptchartengine import (
    create_combo_chart, create_waterfall_chart, create_bubble_chart,
    create_range_snapshot_chart,
    StyleConfig, ChartLayoutConfig, LegendConfig, CategoryAxisConfig, ValueAxisConfig,
)

SOFFICE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
ASSETS = ROOT / "docs" / "assets"
ASSETS.mkdir(parents=True, exist_ok=True)

# ---- real CATL (300750.SZ) annual data, RMB bn (source: Tushare) -------------
YEARS = ["2019", "2020", "2021", "2022", "2023", "2024"]
REV  = [45.79, 50.32, 130.36, 328.59, 400.92, 362.01]
COGS = [32.48, 36.35, 96.09, 262.05, 309.07, 273.52]
OPP  = [5.76, 6.96, 19.82, 36.82, 53.72, 64.05]   # operating profit
NP   = [4.56, 5.58, 15.93, 30.73, 44.12, 50.74]   # net profit attr. parent
GP   = [round(r - c, 2) for r, c in zip(REV, COGS)]
GM   = [round(g / r * 100, 1) for g, r in zip(GP, REV)]
OM   = [round(o / r * 100, 1) for o, r in zip(OPP, REV)]
NM   = [round(n / r * 100, 1) for n, r in zip(NP, REV)]


def _new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_pptx(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    POS, SIZE = (Inches(0.55), Inches(0.55)), (Inches(12.2), Inches(6.2))

    # 1) Combo dual-axis stacked
    df1 = pd.DataFrame({"FY": YEARS, "Cost of goods sold": COGS, "Gross profit": GP,
                        "Gross margin": GM, "Net margin": NM})
    create_combo_chart(
        slide=_new_slide(prs), df=df1, categories_col="FY",
        series_config=[
            {"key": "Cost of goods sold", "name": "Cost of goods sold", "type": "column", "axis": "primary", "grouping": "stacked"},
            {"key": "Gross profit", "name": "Gross profit", "type": "column", "axis": "primary", "grouping": "stacked"},
            {"key": "Gross margin", "name": "Gross margin", "type": "line", "axis": "secondary"},
            {"key": "Net margin", "name": "Net margin", "type": "line", "axis": "secondary"},
        ],
        position=POS, size=SIZE,
        style_config=StyleConfig(color_scheme="categorical", line_width_pt=2.25, marker_style="circle", marker_size=6),
        layout_config=ChartLayoutConfig(
            legend_config=LegendConfig(font_size_pt=11, font_name="Arial"),
            category_axis_config=CategoryAxisConfig(font_size_pt=11, font_name="Arial"),
            value_axis_config=ValueAxisConfig(number_format='"¥"0" bn"', font_name="Arial", min_value=0, max_value=450),
            secondary_value_axis_config=ValueAxisConfig(number_format='0"%"', font_name="Arial", min_value=0, max_value=35),
        ),
    )

    # 2) Waterfall — FY2024 P&L bridge
    bridge = pd.DataFrame({
        "Item": ["Revenue", "COGS", "Gross profit", "Opex", "Operating profit", "Tax & non-op", "Net profit"],
        "Value": [362.0, -273.5, 88.5, -24.4, 64.1, -10.1, 54.0],
        "measure": ["total", "relative", "total", "relative", "total", "relative", "total"],
    })
    create_waterfall_chart(
        slide=_new_slide(prs), df=bridge, categories_col="Item", value_col="Value",
        measure_col="measure", position=POS, size=SIZE,
        show_value_labels=True, show_connectors=True, label_font_name="Arial",
    )

    # 3) Bubble — margin trajectory (x=gross margin, y=net margin, size=revenue)
    bub = pd.DataFrame({"Gross margin %": GM, "Net margin %": NM, "Revenue": REV})
    create_bubble_chart(
        slide=_new_slide(prs), df=bub, x_col="Gross margin %", y_col="Net margin %",
        size_col="Revenue", series_name="FY2019–2024", position=POS, size=SIZE,
        color="7E57C2",
    )

    # 4) Range snapshot — margin ranges (min/max/avg/current=FY2024)
    def rng(vals):
        return min(vals), max(vals), round(sum(vals) / len(vals), 1), vals[-1]
    rows = []
    for label, vals in [("Gross margin", GM), ("Operating margin", OM), ("Net margin", NM)]:
        lo, hi, avg, cur = rng(vals)
        rows.append({"Metric": label, "min": lo, "max": hi, "avg": avg, "current": cur})
    snap = pd.DataFrame(rows)
    create_range_snapshot_chart(
        slide=_new_slide(prs), df=snap, categories_col="Metric",
        min_col="min", max_col="max", average_col="avg", current_col="current",
        position=POS, size=SIZE, number_format='0.0"%"', label_font_name="Arial",
    )

    prs.save(str(path))


def render_pages(pptx):
    subprocess.run([SOFFICE, "--headless", "--convert-to", "pdf", "--outdir", "/tmp/gallery", str(pptx)],
                   check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pdf = Path("/tmp/gallery") / (pptx.stem + ".pdf")
    subprocess.run(["pdftoppm", "-png", "-r", "150", str(pdf), "/tmp/gallery/page"], check=True)
    return sorted(Path("/tmp/gallery").glob("page-*.png"))


def _font(size, bold=False):
    for p in (["/System/Library/Fonts/Supplemental/Arial Bold.ttf"] if bold else ["/System/Library/Fonts/Supplemental/Arial.ttf"]) + \
             ["/System/Library/Fonts/Helvetica.ttc"]:
        try:
            return ImageFont.truetype(p, size)
        except Exception:
            continue
    return ImageFont.load_default()


CAPTIONS = [
    ("Dual-axis combo", "Stacked columns + two lines on a second % axis — one native chart, per-axis number formats"),
    ("Waterfall bridge", "FY2024 P&L bridge with total / relative measures, connectors and value labels"),
    ("Bubble (XY family)", "Margin trajectory; bubble area = revenue — native scatter / bubble with parse helpers"),
    ("Valuation-range snapshot", "Margin ranges with min / max / average and a current-value marker"),
]


def montage(pages, out):
    panels = []
    cap_h, pad, head_h = 96, 28, 130
    tf, ff = _font(40, bold=True), _font(30)
    for img_path, (title, fact) in zip(pages, CAPTIONS):
        im = Image.open(img_path).convert("RGB")
        # trim to a consistent 16:9-ish crop and add caption band on top
        w, h = im.size
        band = Image.new("RGB", (w, h + cap_h), "white")
        d = ImageDraw.Draw(band)
        d.rectangle([0, 0, w, cap_h], fill="#0F2A43")
        d.text((24, 16), title, font=tf, fill="white")
        d.text((24, 60), fact, font=_font(24), fill="#AFC3D6")
        band.paste(im, (0, cap_h))
        panels.append(band)

    pw, ph = panels[0].size
    cols, rows = 2, 2
    W = cols * pw + (cols + 1) * pad
    H = head_h + rows * ph + (rows + 1) * pad
    canvas = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(canvas)
    d.text((pad, 30), "pptchartengine — chart families & palettes",
           font=_font(52, bold=True), fill="#0F2A43")
    d.text((pad, 88),
           "Every panel is a real, editable PowerPoint chart — and can be parsed back into a DataFrame and updated in place.",
           font=ff, fill="#52606D")
    for i, panel in enumerate(panels):
        r, c = divmod(i, cols)
        x = pad + c * (pw + pad)
        y = head_h + pad + r * (ph + pad)
        canvas.paste(panel, (x, y))

    # downscale to a sensible README width
    target_w = 2400
    canvas = canvas.resize((target_w, int(H * target_w / W)), Image.LANCZOS)
    canvas.save(out)
    print("saved", out, canvas.size)


if __name__ == "__main__":
    pptx = ASSETS / "gallery.pptx"
    build_pptx(pptx)
    pages = render_pages(pptx)
    assert len(pages) == 4, f"expected 4 pages, got {len(pages)}"
    montage(pages, ASSETS / "gallery.png")

"""Generate the README hero chart — a complex dual-axis combo on REAL data.

Source: CATL (Contemporary Amperex Technology, 300750.SZ) consolidated annual
income statements, FY2019-2024, as reported (via Tushare). Figures in RMB bn.

Structure (an honest, information-dense dual-axis combo):
  - primary axis  : stacked columns = Cost of goods sold + Gross profit
                    (the two add up to total revenue)
  - secondary axis: 2 lines sharing one % scale = Gross margin, Net margin
  - 6 fiscal years, bottom legend, per-axis number formats

Run:  python scripts/make_readme_hero.py
Output: docs/assets/hero-combo.pptx
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from pptchartengine import (
    create_combo_chart,
    StyleConfig,
    ChartLayoutConfig,
    LegendConfig,
    CategoryAxisConfig,
    ValueAxisConfig,
)

# CATL (300750.SZ) consolidated annual income statement, RMB bn (亿元 / 10).
# total_revenue, oper_cost, net profit attributable to parent — as reported.
ROWS = [
    # year,  revenue, cogs,    np_attr_parent
    ("2019",  45.79,  32.48,   4.56),
    ("2020",  50.32,  36.35,   5.58),
    ("2021", 130.36,  96.09,  15.93),
    ("2022", 328.59, 262.05,  30.73),
    ("2023", 400.92, 309.07,  44.12),
    ("2024", 362.01, 273.52,  50.74),
]

years = [r[0] for r in ROWS]
revenue = [r[1] for r in ROWS]
cogs = [r[2] for r in ROWS]
gross_profit = [round(r[1] - r[2], 2) for r in ROWS]          # = revenue - cogs
gross_margin = [round((r[1] - r[2]) / r[1] * 100, 1) for r in ROWS]
net_margin = [round(r[3] / r[1] * 100, 1) for r in ROWS]

df = pd.DataFrame(
    {
        "FY": years,
        "Cost of goods sold": cogs,
        "Gross profit": gross_profit,
        "Gross margin": gross_margin,
        "Net margin": net_margin,
    }
)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

TITLE = "CATL (300750.SZ) — Revenue split by cost vs. gross profit, with margins  ·  FY2019–2024, RMB bn"

chart = create_combo_chart(
    slide=slide,
    df=df,
    categories_col="FY",
    series_config=[
        {"key": "Cost of goods sold", "name": "Cost of goods sold", "type": "column", "axis": "primary", "grouping": "stacked"},
        {"key": "Gross profit", "name": "Gross profit", "type": "column", "axis": "primary", "grouping": "stacked"},
        {"key": "Gross margin", "name": "Gross margin", "type": "line", "axis": "secondary"},
        {"key": "Net margin", "name": "Net margin", "type": "line", "axis": "secondary"},
    ],
    position=(Inches(0.7), Inches(0.95)),
    size=(Inches(12.0), Inches(5.85)),
    style_config=StyleConfig(color_scheme="categorical", line_width_pt=2.25, marker_style="circle", marker_size=6),
    layout_config=ChartLayoutConfig(
        title=TITLE,
        legend_config=LegendConfig(font_size_pt=11, font_name="Arial"),
        category_axis_config=CategoryAxisConfig(number_format=None, font_size_pt=11, font_name="Arial"),
        value_axis_config=ValueAxisConfig(number_format='"¥"0" bn"', font_name="Arial", min_value=0, max_value=450),
        secondary_value_axis_config=ValueAxisConfig(number_format='0"%"', font_name="Arial", min_value=0, max_value=35),
    ),
)

# Title font: the engine default is a CJK face (微软雅黑) which renderers without
# that font substitute with a serif. Force a portable system sans for the title.
if chart.has_title:
    for para in chart.chart_title.text_frame.paragraphs:
        for run in (para.runs or [para.add_run()]):
            run.font.name = "Arial"

out = Path(__file__).resolve().parent.parent / "docs" / "assets" / "hero-combo.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print("saved", out)

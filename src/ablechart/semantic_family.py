"""Semantic chart families derived from the demo01 analysis surface.

This layer sits above the low-level editable chart primitives and provides
business-shaped families that map onto the existing combo / scatter /
waterfall capabilities.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Mapping, Sequence

import pandas as pd
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from .api import _write_embedded_metadata, create_combo_chart
from .date_axis import DateAxisConfig
from .layout import ChartLayoutConfig, LegendConfig, ValueAxisConfig
from .plot_area import estimate_chart_plot_area
from .scatter import create_scatter_chart
from .semantic_anchor import create_semantic_anchor, update_chart_semantic_metadata
from .styles import COLOR_SCHEMES, StyleConfig
from .waterfall import create_waterfall_chart


PERFORMANCE_COMPARE_FAMILY = "performance_compare"
DISTRIBUTION_PLUS_HISTORY_FAMILY = "distribution_plus_history"
STYLE_BOX_FAMILY = "style_box"
STYLE_ALLOCATION_FAMILY = "style_allocation"
FACTOR_EXPOSURE_FAMILY = "factor_exposure"
SCORE_OVERLAY_FAMILY = "score_overlay"
CONCENTRATION_FAMILY = "concentration"
EVENT_TIMELINE_FAMILY = "event_timeline"
ATTRIBUTION_DECOMPOSITION_FAMILY = "attribution_decomposition"
HOLDING_DETAIL_FAMILY = "holding_detail"
RANKED_TILE_MATRIX_FAMILY = "ranked_tile_matrix"
HEATMAP_MATRIX_FAMILY = "heatmap_matrix"
TABLE_PLUS_CHART_COMPOSITE_FAMILY = "table_plus_chart_composite"
FACTOR_ATTRIBUTION_PANEL_FAMILY = "factor_attribution_panel"
REGIME_TABLE_PANEL_FAMILY = "regime_table_panel"
MANAGER_TIMELINE_PROFILE_FAMILY = "manager_timeline_profile"
AWARD_TIMELINE_PANEL_FAMILY = "award_timeline_panel"
SELECTION_TIMING_GRID_FAMILY = "selection_timing_grid"
DUAL_CHART_PANEL_FAMILY = "dual_chart_panel"


ROLE_COLORS = {
    "fund": "5679CC",
    "portfolio": "5679CC",
    "benchmark": "FFB57D",
    "peer": "E56060",
    "synthetic": "77C3C7",
    "median": "C292F4",
    "raw": "5679CC",
    "score": "86B6F5",
    "excess": "6E9FE6",
}

SEMANTIC_FAMILY_REGISTRY: dict[str, dict[str, Any]] = {
    PERFORMANCE_COMPARE_FAMILY: {
        "renderable": True,
        "base_geometry": "combo",
        "description": "基金与基准/同类/中位数的时序比较。",
    },
    DISTRIBUTION_PLUS_HISTORY_FAMILY: {
        "renderable": True,
        "base_geometry": "combo",
        "modes": ("snapshot", "history"),
        "description": "某时点构成分布 + 历史配置变化。",
    },
    STYLE_BOX_FAMILY: {
        "renderable": True,
        "base_geometry": "scatter",
        "description": "二维风格定位箱体。",
    },
    STYLE_ALLOCATION_FAMILY: {
        "renderable": True,
        "base_geometry": "combo",
        "modes": ("snapshot", "history", "relative"),
        "description": "风格桶配置比例 / 超配比例。",
    },
    FACTOR_EXPOSURE_FAMILY: {
        "renderable": True,
        "base_geometry": "combo",
        "modes": ("compare", "history"),
        "description": "因子暴露对比与时序演化。",
    },
    SCORE_OVERLAY_FAMILY: {
        "renderable": True,
        "base_geometry": "combo",
        "description": "原值叠加相对得分/分位得分。",
    },
    CONCENTRATION_FAMILY: {
        "renderable": True,
        "base_geometry": "combo",
        "description": "集中度与相对得分。",
    },
    EVENT_TIMELINE_FAMILY: {
        "renderable": True,
        "base_geometry": "combo+overlay",
        "description": "主时序叠加市场阶段/调仓事件带。",
    },
    ATTRIBUTION_DECOMPOSITION_FAMILY: {
        "renderable": True,
        "base_geometry": "waterfall",
        "description": "收益来源 / 贡献拆解。",
    },
    RANKED_TILE_MATRIX_FAMILY: {
        "renderable": True,
        "base_geometry": "composer_shapes",
        "description": "带行列标签的排名/配比 tile 矩阵。",
    },
    HEATMAP_MATRIX_FAMILY: {
        "renderable": True,
        "base_geometry": "composer_shapes",
        "description": "带色阶映射的 heatmap 矩阵。",
    },
    TABLE_PLUS_CHART_COMPOSITE_FAMILY: {
        "renderable": True,
        "base_geometry": "chart+table",
        "description": "左图右表的复合分析卡片。",
    },
    FACTOR_ATTRIBUTION_PANEL_FAMILY: {
        "renderable": True,
        "base_geometry": "chart+insight_panel",
        "description": "左侧归因主图 + 右侧贡献摘要/解释面板。",
    },
    REGIME_TABLE_PANEL_FAMILY: {
        "renderable": True,
        "base_geometry": "event_chart+table",
        "description": "上方市场阶段/事件图，下方区间收益表。",
    },
    MANAGER_TIMELINE_PROFILE_FAMILY: {
        "renderable": True,
        "base_geometry": "profile+compare_chart",
        "description": "基金经理任职概况 + 经理指数对比。",
    },
    AWARD_TIMELINE_PANEL_FAMILY: {
        "renderable": True,
        "base_geometry": "table/empty_state",
        "description": "获奖记录面板，支持空状态。",
    },
    SELECTION_TIMING_GRID_FAMILY: {
        "renderable": True,
        "base_geometry": "judgment_grid",
        "description": "选股能力/择时能力的多周期等级判断栅格。",
    },
    DUAL_CHART_PANEL_FAMILY: {
        "renderable": True,
        "base_geometry": "chart+chart",
        "description": "左右双图并排的分析面板。",
    },
    HOLDING_DETAIL_FAMILY: {
        "renderable": True,
        "base_geometry": "table/detail",
        "description": "持仓明细、前五大重仓债券、持仓基金等表格型信息面板。",
    },
}


COLOR_SCHEMES.setdefault(
    "demo01_compare",
    [
        ROLE_COLORS["fund"],
        ROLE_COLORS["peer"],
        ROLE_COLORS["benchmark"],
        ROLE_COLORS["synthetic"],
        ROLE_COLORS["median"],
        ROLE_COLORS["excess"],
        "5FC7B5",
        "87E491",
    ],
)
COLOR_SCHEMES.setdefault(
    "demo01_distribution",
    [
        "6E9FE6",
        "FFB3BD",
        "FFB57D",
        "A8AAED",
        "56B6F5",
        "5FC7B5",
        "87E491",
        "FFE677",
        "90D4FF",
        "74D6B5",
        "F09E60",
        "D3DB72",
    ],
)
COLOR_SCHEMES.setdefault(
    "demo01_score",
    [
        ROLE_COLORS["raw"],
        ROLE_COLORS["score"],
        ROLE_COLORS["excess"],
        ROLE_COLORS["median"],
        ROLE_COLORS["benchmark"],
        "5FC7B5",
        "87E491",
        "B0C9B0",
    ],
)


@dataclass
class SemanticChartParseResult:
    family: str
    base_chart_family: str | None
    metadata: dict[str, Any]


def list_semantic_families(*, renderable_only: bool = False) -> dict[str, dict[str, Any]]:
    if not renderable_only:
        return {key: dict(value) for key, value in SEMANTIC_FAMILY_REGISTRY.items()}
    return {
        key: dict(value)
        for key, value in SEMANTIC_FAMILY_REGISTRY.items()
        if value.get("renderable", False)
    }


def build_performance_compare_preset(
    df: pd.DataFrame,
    categories_col: str,
    series_entries: Sequence[Mapping[str, Any]],
    *,
    title: str | None = None,
    number_format: str = "0.0%",
    secondary_number_format: str | None = None,
    date_number_format: str = "yyyy-mm-dd",
    color_scheme: str = "demo01_compare",
    primary_axis_font_size_pt: float = 9,
    secondary_axis_font_size_pt: float = 9,
    line_width_pt: float = 1.75,
    metadata_overrides: Mapping[str, Any] | None = None,
) -> dict[str, Any]:
    series_config = [_normalize_series_entry(entry) for entry in series_entries]
    date_axis_config = _auto_date_axis_config(df, categories_col, fallback_number_format=date_number_format)
    has_secondary_axis = any(entry.get("axis", "primary") == "secondary" for entry in series_entries)
    layout_config = ChartLayoutConfig(
        title=title,
        legend_config=LegendConfig(
            position=XL_LEGEND_POSITION.TOP,
            font_size_pt=9,
            font_name="黑体",
        ),
        value_axis_config=ValueAxisConfig(
            number_format=number_format,
            font_size_pt=primary_axis_font_size_pt,
            font_name="黑体",
            has_major_gridlines=True,
        ),
        secondary_value_axis_config=(
            ValueAxisConfig(
                number_format=secondary_number_format or number_format,
                font_size_pt=secondary_axis_font_size_pt,
                font_name="黑体",
                has_major_gridlines=False,
            )
            if has_secondary_axis
            else None
        ),
        date_axis_config=date_axis_config,
    )
    metadata = {
        "chart_family": PERFORMANCE_COMPARE_FAMILY,
        "categories_col": categories_col,
        "semantic_family": PERFORMANCE_COMPARE_FAMILY,
        "series_roles": [
            {
                "key": entry["key"],
                "name": entry["name"],
                "role": entry.get("role", "fund"),
                "type": entry.get("type", "line"),
                "axis": entry.get("axis", "primary"),
            }
            for entry in series_entries
        ],
    }
    if metadata_overrides:
        metadata.update(dict(metadata_overrides))
    return {
        "df": df,
        "categories_col": categories_col,
        "series_config": series_config,
        "style_config": StyleConfig(
            color_scheme=color_scheme,
            line_width_pt=line_width_pt,
            marker_style="none",
        ),
        "layout_config": layout_config,
        "metadata": metadata,
    }


def create_performance_compare_chart(
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    series_entries: Sequence[Mapping[str, Any]],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    title: str | None = None,
    number_format: str = "0.0%",
    secondary_number_format: str | None = None,
    date_number_format: str = "yyyy-mm-dd",
    color_scheme: str = "demo01_compare",
    primary_axis_font_size_pt: float = 9,
    secondary_axis_font_size_pt: float = 9,
    line_width_pt: float = 1.75,
    metadata_overrides: Mapping[str, Any] | None = None,
):
    preset = build_performance_compare_preset(
        df,
        categories_col,
        series_entries,
        title=title,
        number_format=number_format,
        secondary_number_format=secondary_number_format,
        date_number_format=date_number_format,
        color_scheme=color_scheme,
        primary_axis_font_size_pt=primary_axis_font_size_pt,
        secondary_axis_font_size_pt=secondary_axis_font_size_pt,
        line_width_pt=line_width_pt,
        metadata_overrides=metadata_overrides,
    )
    return create_combo_chart(
        slide=slide,
        df=preset["df"],
        categories_col=preset["categories_col"],
        series_config=preset["series_config"],
        position=position,
        size=size,
        style_config=preset["style_config"],
        layout_config=preset["layout_config"],
        metadata=preset["metadata"],
    )


def build_distribution_snapshot_preset(
    df: pd.DataFrame,
    category_col: str,
    value_col: str,
    *,
    snapshot_label: str = "当前分布",
    title: str | None = None,
    number_format: str = "0%",
    color_scheme: str = "demo01_distribution",
    semantic_family: str = DISTRIBUTION_PLUS_HISTORY_FAMILY,
    metadata_overrides: Mapping[str, Any] | None = None,
) -> dict[str, Any]:
    values = pd.to_numeric(df[value_col], errors="raise")
    total = float(values.sum())
    normalized = values / total if total else values
    prepared_row = {"__snapshot__": snapshot_label}
    series_config = []
    for (_, row), share in zip(df.iterrows(), normalized.tolist()):
        series_name = str(row[category_col])
        prepared_row[series_name] = float(share)
        series_config.append(
            {
                "key": series_name,
                "name": series_name,
                "type": "bar",
                "axis": "primary",
                "grouping": "percent_stacked",
            }
        )

    prepared_df = pd.DataFrame([prepared_row])
    layout_config = ChartLayoutConfig(
        title=title,
        legend_config=LegendConfig(
            position=XL_LEGEND_POSITION.RIGHT,
            font_size_pt=8.5,
            font_name="黑体",
        ),
        value_axis_config=ValueAxisConfig(
            number_format=number_format,
            font_size_pt=9,
            font_name="黑体",
            has_major_gridlines=False,
            min_value=0.0,
            max_value=1.0,
        ),
    )
    metadata = {
        "chart_family": semantic_family,
        "semantic_family": semantic_family,
        "mode": "snapshot",
        "source_category_col": category_col,
        "source_value_col": value_col,
        "snapshot_label": snapshot_label,
    }
    if metadata_overrides:
        metadata.update(dict(metadata_overrides))
    return {
        "df": prepared_df,
        "categories_col": "__snapshot__",
        "series_config": series_config,
        "style_config": StyleConfig(
            color_scheme=color_scheme,
            line_width_pt=1.0,
            marker_style="none",
        ),
        "layout_config": layout_config,
        "metadata": metadata,
    }


def create_distribution_snapshot_chart(
    slide: Slide,
    df: pd.DataFrame,
    category_col: str,
    value_col: str,
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    snapshot_label: str = "当前分布",
    title: str | None = None,
    number_format: str = "0%",
    color_scheme: str = "demo01_distribution",
    semantic_family: str = DISTRIBUTION_PLUS_HISTORY_FAMILY,
):
    preset = build_distribution_snapshot_preset(
        df,
        category_col,
        value_col,
        snapshot_label=snapshot_label,
        title=title,
        number_format=number_format,
        color_scheme=color_scheme,
        semantic_family=semantic_family,
    )
    return create_combo_chart(
        slide=slide,
        df=preset["df"],
        categories_col=preset["categories_col"],
        series_config=preset["series_config"],
        position=position,
        size=size,
        style_config=preset["style_config"],
        layout_config=preset["layout_config"],
        metadata=preset["metadata"],
    )


def build_distribution_history_preset(
    df: pd.DataFrame,
    categories_col: str,
    series_columns: Sequence[str],
    *,
    chart_type: str = "area",
    title: str | None = None,
    number_format: str = "0%",
    date_number_format: str = "yyyy-mm-dd",
    color_scheme: str = "demo01_distribution",
    semantic_family: str = DISTRIBUTION_PLUS_HISTORY_FAMILY,
    metadata_overrides: Mapping[str, Any] | None = None,
) -> dict[str, Any]:
    series_config = [
        {
            "key": column,
            "name": column,
            "type": chart_type,
            "axis": "primary",
            "grouping": "percent_stacked",
        }
        for column in series_columns
    ]
    date_axis_config = _auto_date_axis_config(df, categories_col, fallback_number_format=date_number_format)
    layout_config = ChartLayoutConfig(
        title=title,
        legend_config=LegendConfig(
            position=XL_LEGEND_POSITION.RIGHT,
            font_size_pt=8.5,
            font_name="黑体",
        ),
        value_axis_config=ValueAxisConfig(
            number_format=number_format,
            font_size_pt=9,
            font_name="黑体",
            has_major_gridlines=False,
            min_value=0.0,
            max_value=1.0,
        ),
        date_axis_config=date_axis_config,
    )
    metadata = {
        "chart_family": semantic_family,
        "semantic_family": semantic_family,
        "mode": "history",
        "series_columns": list(series_columns),
    }
    if metadata_overrides:
        metadata.update(dict(metadata_overrides))
    return {
        "df": df,
        "categories_col": categories_col,
        "series_config": series_config,
        "style_config": StyleConfig(
            color_scheme=color_scheme,
            line_width_pt=1.0,
            marker_style="none",
        ),
        "layout_config": layout_config,
        "metadata": metadata,
    }


def create_distribution_history_chart(
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    series_columns: Sequence[str],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    chart_type: str = "area",
    title: str | None = None,
    number_format: str = "0%",
    date_number_format: str = "yyyy-mm-dd",
    color_scheme: str = "demo01_distribution",
    semantic_family: str = DISTRIBUTION_PLUS_HISTORY_FAMILY,
):
    preset = build_distribution_history_preset(
        df,
        categories_col,
        series_columns,
        chart_type=chart_type,
        title=title,
        number_format=number_format,
        date_number_format=date_number_format,
        color_scheme=color_scheme,
        semantic_family=semantic_family,
    )
    return create_combo_chart(
        slide=slide,
        df=preset["df"],
        categories_col=preset["categories_col"],
        series_config=preset["series_config"],
        position=position,
        size=size,
        style_config=preset["style_config"],
        layout_config=preset["layout_config"],
        metadata=preset["metadata"],
    )


def build_score_overlay_preset(
    df: pd.DataFrame,
    categories_col: str,
    raw_series: Sequence[Mapping[str, Any]],
    score_series: Sequence[Mapping[str, Any]],
    *,
    title: str | None = None,
    raw_number_format: str = "0.0%",
    score_number_format: str = "0",
    date_number_format: str = "yyyy-mm-dd",
    color_scheme: str = "demo01_score",
    semantic_family: str = SCORE_OVERLAY_FAMILY,
    metadata_overrides: Mapping[str, Any] | None = None,
) -> dict[str, Any]:
    series_config = []
    for entry in raw_series:
        series_config.append(
            {
                "key": entry["key"],
                "name": entry.get("name", entry["key"]),
                "type": entry.get("type", "line"),
                "axis": "primary",
                "role": entry.get("role", "raw"),
            }
        )
    for entry in score_series:
        series_config.append(
            {
                "key": entry["key"],
                "name": entry.get("name", entry["key"]),
                "type": entry.get("type", "line"),
                "axis": "secondary",
                "role": entry.get("role", "score"),
            }
        )
    date_axis_config = _auto_date_axis_config(df, categories_col, fallback_number_format=date_number_format)
    layout_config = ChartLayoutConfig(
        title=title,
        legend_config=LegendConfig(
            position=XL_LEGEND_POSITION.TOP,
            font_size_pt=9,
            font_name="黑体",
        ),
        value_axis_config=ValueAxisConfig(
            number_format=raw_number_format,
            font_size_pt=9,
            font_name="黑体",
            has_major_gridlines=True,
        ),
        secondary_value_axis_config=ValueAxisConfig(
            number_format=score_number_format,
            font_size_pt=9,
            font_name="黑体",
            has_major_gridlines=False,
        ),
        date_axis_config=date_axis_config,
    )
    metadata = {
        "chart_family": semantic_family,
        "semantic_family": semantic_family,
        "raw_series": [dict(item) for item in raw_series],
        "score_series": [dict(item) for item in score_series],
    }
    if metadata_overrides:
        metadata.update(dict(metadata_overrides))
    return {
        "df": df,
        "categories_col": categories_col,
        "series_config": series_config,
        "style_config": StyleConfig(
            color_scheme=color_scheme,
            line_width_pt=1.5,
            marker_style="none",
        ),
        "layout_config": layout_config,
        "metadata": metadata,
    }


def create_score_overlay_chart(
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    raw_series: Sequence[Mapping[str, Any]],
    score_series: Sequence[Mapping[str, Any]],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    title: str | None = None,
    raw_number_format: str = "0.0%",
    score_number_format: str = "0",
    date_number_format: str = "yyyy-mm-dd",
    color_scheme: str = "demo01_score",
    semantic_family: str = SCORE_OVERLAY_FAMILY,
):
    preset = build_score_overlay_preset(
        df,
        categories_col,
        raw_series,
        score_series,
        title=title,
        raw_number_format=raw_number_format,
        score_number_format=score_number_format,
        date_number_format=date_number_format,
        color_scheme=color_scheme,
        semantic_family=semantic_family,
    )
    return create_combo_chart(
        slide=slide,
        df=preset["df"],
        categories_col=preset["categories_col"],
        series_config=preset["series_config"],
        position=position,
        size=size,
        style_config=preset["style_config"],
        layout_config=preset["layout_config"],
        metadata=preset["metadata"],
    )


def create_concentration_chart(
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    raw_series: Sequence[Mapping[str, Any]],
    score_series: Sequence[Mapping[str, Any]],
    **kwargs,
):
    return create_score_overlay_chart(
        slide=slide,
        df=df,
        categories_col=categories_col,
        raw_series=raw_series,
        score_series=score_series,
        semantic_family=CONCENTRATION_FAMILY,
        **kwargs,
    )


def create_style_allocation_chart(
    slide: Slide,
    df: pd.DataFrame,
    *,
    mode: str,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    title: str | None = None,
    categories_col: str | None = None,
    series_columns: Sequence[str] | None = None,
    category_col: str | None = None,
    value_col: str | None = None,
):
    if mode == "snapshot":
        if category_col is None or value_col is None:
            raise ValueError("style_allocation snapshot 模式需要 category_col 和 value_col")
        return create_distribution_snapshot_chart(
            slide=slide,
            df=df,
            category_col=category_col,
            value_col=value_col,
            position=position,
            size=size,
            title=title,
            semantic_family=STYLE_ALLOCATION_FAMILY,
        )

    if categories_col is None or not series_columns:
        raise ValueError("style_allocation history/relative 模式需要 categories_col 和 series_columns")
    return create_distribution_history_chart(
        slide=slide,
        df=df,
        categories_col=categories_col,
        series_columns=series_columns,
        position=position,
        size=size,
        title=title,
        semantic_family=STYLE_ALLOCATION_FAMILY,
    )


def create_factor_exposure_chart(
    slide: Slide,
    df: pd.DataFrame,
    *,
    mode: str,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    title: str | None = None,
    categories_col: str,
    series_entries: Sequence[Mapping[str, Any]],
):
    if mode == "history":
        return create_performance_compare_chart(
            slide=slide,
            df=df,
            categories_col=categories_col,
            series_entries=series_entries,
            position=position,
            size=size,
            title=title,
            color_scheme="demo01_score",
            metadata_overrides={"semantic_family": FACTOR_EXPOSURE_FAMILY, "mode": "history"},
        )

    # compare mode: usually factors on x-axis and portfolio/benchmark/excess as bars
    return create_performance_compare_chart(
        slide=slide,
        df=df,
        categories_col=categories_col,
        series_entries=series_entries,
        position=position,
        size=size,
        title=title,
        number_format="0.0",
        date_number_format="General",
        color_scheme="demo01_score",
        metadata_overrides={"semantic_family": FACTOR_EXPOSURE_FAMILY, "mode": "compare"},
    )


def build_style_box_spec(
    df: pd.DataFrame,
    x_col: str,
    y_col: str,
    *,
    series_name: str | None = None,
    x_label: str | None = None,
    y_label: str | None = None,
) -> dict[str, Any]:
    return {
        "chart_family": "scatter",
        "base_chart_family": "scatter",
        "semantic_family": STYLE_BOX_FAMILY,
        "x_col": x_col,
        "y_col": y_col,
        "series_name": series_name or y_col,
        "x_label": x_label or x_col,
        "y_label": y_label or y_col,
        "data_points": [
            {"x": float(row[x_col]), "y": float(row[y_col])}
            for _, row in df.iterrows()
        ],
    }


def create_style_box_chart(
    slide: Slide,
    df: pd.DataFrame,
    x_col: str,
    y_col: str,
    *,
    series_name: str | None = None,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    color: str = ROLE_COLORS["fund"],
):
    chart = create_scatter_chart(
        slide=slide,
        df=df,
        x_col=x_col,
        y_col=y_col,
        series_name=series_name,
        position=position,
        size=size,
        color=color,
        marker_size=11,
    )
    _write_embedded_metadata(
        chart,
        x_col,
        [
            {
                "key": y_col,
                "name": series_name or y_col,
                "type": "scatter",
                "axis": "primary",
                "x_key": x_col,
            }
        ],
        metadata=build_style_box_spec(df, x_col, y_col, series_name=series_name),
    )
    _add_style_box_quadrants(slide, chart=chart, position=position, size=size)
    return chart


def build_event_timeline_preset(
    df: pd.DataFrame,
    categories_col: str,
    series_entries: Sequence[Mapping[str, Any]],
    events: Sequence[Mapping[str, Any]],
    *,
    title: str | None = None,
    number_format: str = "0.0%",
    date_number_format: str = "yyyy-mm-dd",
    color_scheme: str = "demo01_compare",
    show_event_labels: bool = False,
) -> dict[str, Any]:
    preset = build_performance_compare_preset(
        df,
        categories_col,
        series_entries,
        title=title,
        number_format=number_format,
        date_number_format=date_number_format,
        color_scheme=color_scheme,
        metadata_overrides={
            "semantic_family": EVENT_TIMELINE_FAMILY,
            "events": [dict(item) for item in events],
            "show_event_labels": show_event_labels,
        },
    )
    preset["metadata"]["chart_family"] = EVENT_TIMELINE_FAMILY
    return preset


def create_event_timeline_chart(
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    series_entries: Sequence[Mapping[str, Any]],
    events: Sequence[Mapping[str, Any]],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    title: str | None = None,
    number_format: str = "0.0%",
    date_number_format: str = "yyyy-mm-dd",
    color_scheme: str = "demo01_compare",
    show_event_labels: bool = False,
):
    preset = build_event_timeline_preset(
        df,
        categories_col,
        series_entries,
        events,
        title=title,
        number_format=number_format,
        date_number_format=date_number_format,
        color_scheme=color_scheme,
        show_event_labels=show_event_labels,
    )
    chart = create_combo_chart(
        slide=slide,
        df=preset["df"],
        categories_col=preset["categories_col"],
        series_config=preset["series_config"],
        position=position,
        size=size,
        style_config=preset["style_config"],
        layout_config=preset["layout_config"],
        metadata=preset["metadata"],
    )
    _add_event_overlays(
        slide,
        chart,
        df,
        categories_col,
        events,
        position=position,
        size=size,
        show_event_labels=show_event_labels,
    )
    return chart


def create_attribution_decomposition_chart(
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    value_col: str,
    *,
    measure_col: str | None = None,
    total_categories: Sequence[str] | None = None,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    positive_color: str = "10B981",
    negative_color: str = "EF4444",
    total_color: str = "1E2761",
):
    chart = create_waterfall_chart(
        slide=slide,
        df=df,
        categories_col=categories_col,
        value_col=value_col,
        measure_col=measure_col,
        total_categories=total_categories,
        position=position,
        size=size,
        positive_color=positive_color,
        negative_color=negative_color,
        total_color=total_color,
        show_legend=False,
    )
    _write_embedded_metadata(
        chart,
        categories_col,
        [
            {"key": "__base__", "name": "__base__", "type": "bar", "axis": "primary", "grouping": "stacked"},
            {"key": "__increase__", "name": "__increase__", "type": "bar", "axis": "primary", "grouping": "stacked"},
            {"key": "__decrease__", "name": "__decrease__", "type": "bar", "axis": "primary", "grouping": "stacked"},
            {"key": "__total__", "name": "__total__", "type": "bar", "axis": "primary", "grouping": "stacked"},
        ],
        metadata={
            "chart_family": "waterfall",
            "base_chart_family": "waterfall",
            "semantic_family": ATTRIBUTION_DECOMPOSITION_FAMILY,
            "categories_col": categories_col,
            "value_col": value_col,
            "measure_col": measure_col,
            "total_categories": list(total_categories or []),
        },
    )
    return chart


def create_ranked_tile_matrix_chart(
    slide: Slide,
    df: pd.DataFrame,
    row_col: str,
    column_col: str,
    value_col: str,
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    family: str = RANKED_TILE_MATRIX_FAMILY,
    value_format: str = ".0f",
    low_color: str = "F4F6FB",
    high_color: str = "4F66A8",
    text_dark: str = "2F3542",
    text_light: str = "FFFFFF",
    row_header_fill: str = "FFFFFF",
    col_header_fill: str = "FFFFFF",
    grid_line_color: str = "DCE3EF",
):
    matrix = _pivot_matrix(df, row_col, column_col, value_col)
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)

    row_labels = list(matrix.index)
    col_labels = list(matrix.columns)
    if not row_labels or not col_labels:
        raise ValueError("matrix family 需要至少 1 行 1 列")

    values = matrix.to_numpy(dtype=float)
    min_val = float(values.min())
    max_val = float(values.max())

    row_header_w = min(max(width * 0.12, 0.55), 1.2)
    col_header_h = min(max(height * 0.12, 0.32), 0.55)
    footnote_h = 0.18 if family == RANKED_TILE_MATRIX_FAMILY else 0.0
    grid_left = left + row_header_w
    grid_top = top + col_header_h
    grid_width = width - row_header_w
    grid_height = height - col_header_h - footnote_h
    cell_w = grid_width / len(col_labels)
    cell_h = grid_height / len(row_labels)

    # Column headers
    for idx, label in enumerate(col_labels):
        x = grid_left + idx * cell_w
        _add_rect(slide, x, top, cell_w, col_header_h, fill=_rgb(col_header_fill))
        _add_textbox(
            slide,
            str(label),
            x,
            top + col_header_h * 0.18,
            cell_w,
            col_header_h * 0.64,
            font_name="微软雅黑",
            font_size=10,
            color=text_dark,
            bold=True,
        )

    # Row headers
    for idx, label in enumerate(row_labels):
        y = grid_top + idx * cell_h
        _add_rect(slide, left, y, row_header_w, cell_h, fill=_rgb(row_header_fill))
        _add_textbox(
            slide,
            str(label),
            left + 0.04,
            y + cell_h * 0.28,
            row_header_w - 0.08,
            cell_h * 0.44,
            font_name="微软雅黑",
            font_size=10,
            color=text_dark,
            bold=False,
            align="left",
        )

    for row_idx, row_label in enumerate(row_labels):
        for col_idx, col_label in enumerate(col_labels):
            value = float(matrix.loc[row_label, col_label])
            x = grid_left + col_idx * cell_w
            y = grid_top + row_idx * cell_h
            fill_hex = _interpolate_hex(low_color, high_color, value, min_val, max_val)
            _add_rect(slide, x, y, cell_w, cell_h, fill=_rgb(fill_hex))
            text_color = text_light if _hex_luminance(fill_hex) < 0.53 else text_dark
            _add_textbox(
                slide,
                _format_matrix_value(value, value_format),
                x,
                y + cell_h * 0.22,
                cell_w,
                cell_h * 0.52,
                font_name="Aptos",
                font_size=max(12, min(20, int(cell_h * 10))),
                color=text_color,
                bold=False,
            )

    # Grid lines
    for idx in range(len(col_labels) + 1):
        x = grid_left + idx * cell_w
        _add_rect(slide, x - 0.003, grid_top, 0.006, grid_height, fill=_rgb(grid_line_color))
    for idx in range(len(row_labels) + 1):
        y = grid_top + idx * cell_h
        _add_rect(slide, grid_left, y - 0.003, grid_width, 0.006, fill=_rgb(grid_line_color))

    if family == RANKED_TILE_MATRIX_FAMILY:
        _add_textbox(
            slide,
            "注：矩阵值可表示占比、分数或排名强度。",
            left,
            top + height - 0.15,
            width,
            0.14,
            font_name="微软雅黑",
            font_size=8,
            color="8C95A3",
            bold=False,
            align="left",
        )

    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "row_col": row_col,
        "column_col": column_col,
        "value_col": value_col,
        "records": df[[row_col, column_col, value_col]].to_dict(orient="records"),
        "rows": [str(item) for item in row_labels],
        "columns": [str(item) for item in col_labels],
        "value_format": value_format,
    }
    _persist_semantic_component_anchor(slide, metadata)
    return metadata


def create_heatmap_matrix_chart(
    slide: Slide,
    df: pd.DataFrame,
    row_col: str,
    column_col: str,
    value_col: str,
    **kwargs,
):
    return create_ranked_tile_matrix_chart(
        slide=slide,
        df=df,
        row_col=row_col,
        column_col=column_col,
        value_col=value_col,
        family=HEATMAP_MATRIX_FAMILY,
        **kwargs,
    )


def create_table_plus_chart_composite(
    slide: Slide,
    chart_family: str,
    chart_kwargs: Mapping[str, Any],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(10), Inches(4.8)),
    headers: Sequence[str] | None = None,
    rows: Sequence[Sequence[Any]] | None = None,
    table_df: pd.DataFrame | None = None,
    chart_ratio: float = 0.62,
    gap_inches: float = 0.30,
    left_title: str | None = None,
    right_title: str | None = None,
    family: str = TABLE_PLUS_CHART_COMPOSITE_FAMILY,
) -> dict[str, Any]:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)
    chart_width = width * chart_ratio
    table_width = width - chart_width - gap_inches
    chart_position = (Inches(left), Inches(top))
    chart_size = (Inches(chart_width), Inches(height))

    if left_title:
        _add_textbox(
            slide,
            left_title,
            left,
            max(0.15, top - 0.28),
            chart_width,
            0.22,
            font_name="微软雅黑",
            font_size=10,
            color="5F6772",
            bold=False,
            align="left",
        )

    nested_kwargs = dict(chart_kwargs)
    if isinstance(nested_kwargs.get("df"), list):
        nested_kwargs["df"] = pd.DataFrame(nested_kwargs["df"])
    nested_kwargs["position"] = chart_position
    nested_kwargs["size"] = chart_size
    nested_result = create_semantic_chart(slide=slide, family=chart_family, **nested_kwargs)

    table_x = left + chart_width + gap_inches
    if right_title:
        _add_textbox(
            slide,
            right_title,
            table_x,
            max(0.15, top - 0.28),
            table_width,
            0.22,
            font_name="微软雅黑",
            font_size=10,
            color="5F6772",
            bold=False,
            align="left",
        )

    if table_df is not None:
        headers = list(table_df.columns)
        rows = table_df.values.tolist()
    if headers and rows:
        _add_native_table(
            slide,
            headers=headers,
            rows=rows,
            x=table_x,
            y=top,
            w=table_width,
            h=min(height, 0.42 * len(rows) + 0.42),
        )

    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "nested_chart_family": chart_family,
        "chart_kwargs": chart_kwargs,
        "position": [left, top],
        "size": [width, height],
        "chart_ratio": chart_ratio,
        "gap_inches": gap_inches,
        "left_title": left_title,
        "right_title": right_title,
        "headers": list(headers or []),
        "rows": list(rows or []),
        "row_count": len(rows or []),
    }
    anchor_payload = _persist_semantic_component_anchor(slide, metadata)
    _attach_semantic_parent_to_chart(nested_result, family, anchor_payload["semantic_anchor_id"])
    return metadata


def create_factor_attribution_panel(
    slide: Slide,
    chart_family: str,
    chart_kwargs: Mapping[str, Any],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(10), Inches(4.8)),
    sidebar_title: str | None = None,
    summary_items: Sequence[Mapping[str, Any]] | None = None,
    bullets: Sequence[str] | None = None,
    chart_ratio: float = 0.66,
    gap_inches: float = 0.32,
    family: str = FACTOR_ATTRIBUTION_PANEL_FAMILY,
) -> dict[str, Any]:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)
    chart_width = width * chart_ratio
    panel_width = width - chart_width - gap_inches
    chart_position = (Inches(left), Inches(top))
    chart_size = (Inches(chart_width), Inches(height))

    nested_kwargs = dict(chart_kwargs)
    if isinstance(nested_kwargs.get("df"), list):
        nested_kwargs["df"] = pd.DataFrame(nested_kwargs["df"])
    nested_kwargs["position"] = chart_position
    nested_kwargs["size"] = chart_size
    nested_result = create_semantic_chart(slide=slide, family=chart_family, **nested_kwargs)

    panel_x = left + chart_width + gap_inches
    panel_y = top
    current_y = panel_y

    if sidebar_title:
        _add_textbox(
            slide,
            sidebar_title,
            panel_x,
            current_y,
            panel_width,
            0.28,
            font_name="微软雅黑",
            font_size=11,
            color="2F3542",
            bold=True,
            align="left",
        )
        current_y += 0.34

    for item in summary_items or []:
        card_h = float(item.get("height", 0.66))
        fill = _rgb(item.get("fill", "F6F8FC"))
        accent = item.get("accent", "5679CC")
        _add_rect(slide, panel_x, current_y, panel_width, card_h, fill=fill)
        _add_rect(slide, panel_x, current_y, 0.05, card_h, fill=_rgb(accent))
        _add_textbox(
            slide,
            str(item.get("label", "")),
            panel_x + 0.12,
            current_y + 0.08,
            panel_width - 0.18,
            0.18,
            font_name="微软雅黑",
            font_size=9,
            color=item.get("label_color", "6B7280"),
            bold=False,
            align="left",
        )
        _add_textbox(
            slide,
            str(item.get("value", "")),
            panel_x + 0.12,
            current_y + 0.26,
            panel_width - 0.18,
            0.22,
            font_name="Aptos",
            font_size=float(item.get("value_size", 16)),
            color=item.get("value_color", "2F3542"),
            bold=bool(item.get("bold", True)),
            align="left",
        )
        if item.get("note"):
            _add_textbox(
                slide,
                str(item["note"]),
                panel_x + 0.12,
                current_y + card_h - 0.20,
                panel_width - 0.18,
                0.16,
                font_name="微软雅黑",
                font_size=8,
                color=item.get("note_color", "8C95A3"),
                bold=False,
                align="left",
            )
        current_y += card_h + 0.10

    if bullets:
        if current_y < panel_y + height - 0.22:
            _add_textbox(
                slide,
                "结论",
                panel_x,
                current_y,
                panel_width,
                0.22,
                font_name="微软雅黑",
                font_size=10,
                color="5F6772",
                bold=True,
                align="left",
            )
            current_y += 0.26

        for bullet in bullets:
            if current_y > panel_y + height - 0.22:
                break
            _add_textbox(
                slide,
                f"• {bullet}",
                panel_x,
                current_y,
                panel_width,
                0.30,
                font_name="微软雅黑",
                font_size=9,
                color="4B5563",
                bold=False,
                align="left",
            )
            current_y += 0.28

    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "nested_chart_family": chart_family,
        "chart_kwargs": chart_kwargs,
        "position": [left, top],
        "size": [width, height],
        "chart_ratio": chart_ratio,
        "gap_inches": gap_inches,
        "sidebar_title": sidebar_title,
        "summary_items": list(summary_items or []),
        "bullets": list(bullets or []),
        "summary_count": len(summary_items or []),
        "bullet_count": len(bullets or []),
    }
    anchor_payload = _persist_semantic_component_anchor(slide, metadata)
    _attach_semantic_parent_to_chart(nested_result, family, anchor_payload["semantic_anchor_id"])
    return metadata


def create_regime_table_panel(
    slide: Slide,
    chart_kwargs: Mapping[str, Any],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(11), Inches(5.2)),
    table_headers: Sequence[str] | None = None,
    table_rows: Sequence[Sequence[Any]] | None = None,
    table_df: pd.DataFrame | None = None,
    top_title: str | None = None,
    bottom_title: str | None = None,
    chart_height_ratio: float = 0.42,
    gap_inches: float = 0.24,
    family: str = REGIME_TABLE_PANEL_FAMILY,
) -> dict[str, Any]:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)
    chart_height = height * chart_height_ratio
    table_y = top + chart_height + gap_inches
    table_height = height - chart_height - gap_inches

    if top_title:
        _add_textbox(
            slide,
            top_title,
            left,
            max(0.15, top - 0.26),
            width,
            0.22,
            font_name="微软雅黑",
            font_size=10,
            color="5F6772",
            bold=False,
            align="left",
        )

    nested_kwargs = dict(chart_kwargs)
    if isinstance(nested_kwargs.get("df"), list):
        nested_kwargs["df"] = pd.DataFrame(nested_kwargs["df"])
    nested_kwargs["position"] = (Inches(left), Inches(top))
    nested_kwargs["size"] = (Inches(width), Inches(chart_height))
    nested_result = create_semantic_chart(slide=slide, family=EVENT_TIMELINE_FAMILY, **nested_kwargs)

    if table_df is not None:
        table_headers = list(table_df.columns)
        table_rows = table_df.values.tolist()

    if bottom_title:
        _add_textbox(
            slide,
            bottom_title,
            left,
            table_y - 0.02,
            width,
            0.20,
            font_name="微软雅黑",
            font_size=10,
            color="5F6772",
            bold=False,
            align="left",
        )
        table_y += 0.18
        table_height -= 0.18

    if table_headers and table_rows:
        _add_native_table(
            slide,
            headers=table_headers,
            rows=table_rows,
            x=left,
            y=table_y,
            w=width,
            h=min(table_height, 0.42 * len(table_rows) + 0.42),
        )

    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "nested_chart_family": EVENT_TIMELINE_FAMILY,
        "chart_kwargs": chart_kwargs,
        "position": [left, top],
        "size": [width, height],
        "top_title": top_title,
        "bottom_title": bottom_title,
        "table_headers": list(table_headers or []),
        "table_rows": list(table_rows or []),
        "row_count": len(table_rows or []),
    }
    anchor_payload = _persist_semantic_component_anchor(slide, metadata)
    _attach_semantic_parent_to_chart(nested_result, family, anchor_payload["semantic_anchor_id"])
    return metadata


def create_manager_timeline_profile(
    slide: Slide,
    chart_kwargs: Mapping[str, Any],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(11), Inches(5.2)),
    manager_name: str,
    tenure_start: str,
    tenure_days: str | int,
    tenure_return: str,
    summary_items: Sequence[Mapping[str, Any]] | None = None,
    family: str = MANAGER_TIMELINE_PROFILE_FAMILY,
) -> dict[str, Any]:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)

    profile_h = min(1.10, height * 0.22)
    gap = 0.18
    chart_top = top + profile_h + gap
    chart_h = height - profile_h - gap

    _add_rect(slide, left, top, width, profile_h, fill=_rgb("F6F8FC"))
    _add_rect(slide, left, top, 0.06, profile_h, fill=_rgb("5679CC"))
    _add_textbox(slide, manager_name, left + 0.16, top + 0.10, 2.2, 0.26, font_name="微软雅黑", font_size=16, color="1A2744", bold=True, align="left")
    _add_textbox(slide, f"任职开始时间：{tenure_start}", left + 0.16, top + 0.42, 2.4, 0.18, font_name="微软雅黑", font_size=9, color="6B7C93", bold=False, align="left")
    _add_textbox(slide, f"任职时长：{tenure_days}天", left + 2.75, top + 0.42, 1.8, 0.18, font_name="微软雅黑", font_size=9, color="6B7C93", bold=False, align="left")
    _add_textbox(slide, f"任期回报：{tenure_return}", left + 4.65, top + 0.42, 1.8, 0.18, font_name="微软雅黑", font_size=9, color="6B7C93", bold=False, align="left")

    current_x = left + 6.1
    for item in summary_items or []:
        card_w = float(item.get("width", 1.45))
        _add_rect(slide, current_x, top + 0.12, card_w, profile_h - 0.24, fill=_rgb(item.get("fill", "FFFFFF")))
        _add_textbox(slide, str(item.get("label", "")), current_x + 0.08, top + 0.20, card_w - 0.16, 0.16, font_name="微软雅黑", font_size=8, color=item.get("label_color", "8C95A3"), bold=False, align="left")
        _add_textbox(slide, str(item.get("value", "")), current_x + 0.08, top + 0.44, card_w - 0.16, 0.22, font_name="Aptos", font_size=float(item.get("value_size", 13)), color=item.get("value_color", "2F3542"), bold=True, align="left")
        current_x += card_w + 0.10

    nested_kwargs = dict(chart_kwargs)
    if isinstance(nested_kwargs.get("df"), list):
        nested_kwargs["df"] = pd.DataFrame(nested_kwargs["df"])
    nested_kwargs["position"] = (Inches(left), Inches(chart_top))
    nested_kwargs["size"] = (Inches(width), Inches(chart_h))
    nested_result = create_semantic_chart(slide=slide, family=PERFORMANCE_COMPARE_FAMILY, **nested_kwargs)

    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "nested_chart_family": PERFORMANCE_COMPARE_FAMILY,
        "chart_kwargs": chart_kwargs,
        "position": [left, top],
        "size": [width, height],
        "manager_name": manager_name,
        "tenure_start": tenure_start,
        "tenure_days": tenure_days,
        "tenure_return": tenure_return,
        "summary_items": list(summary_items or []),
    }
    anchor_payload = _persist_semantic_component_anchor(slide, metadata)
    _attach_semantic_parent_to_chart(nested_result, family, anchor_payload["semantic_anchor_id"])
    return metadata


def create_award_timeline_panel(
    slide: Slide,
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(11), Inches(4.2)),
    headers: Sequence[str],
    rows: Sequence[Sequence[Any]] | None = None,
    empty_title: str = "暂无数据",
    empty_subtitle: str = "暂无数据/结果",
    family: str = AWARD_TIMELINE_PANEL_FAMILY,
) -> dict[str, Any]:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)
    rows = list(rows or [])

    if rows:
        _add_native_table(slide, headers=headers, rows=rows, x=left, y=top, w=width, h=min(height, 0.42 * len(rows) + 0.42))
    else:
        _add_rect(slide, left, top, width, height, fill=_rgb("F8FAFD"))
        _add_rect(slide, left, top, width, 0.28, fill=_rgb("F1F5FB"))
        header_w = width / max(len(headers), 1)
        for idx, header in enumerate(headers):
            _add_textbox(slide, str(header), left + idx * header_w, top + 0.04, header_w, 0.16, font_name="微软雅黑", font_size=9, color="5F6772", bold=True, align="center")
        _add_textbox(slide, empty_title, left, top + height * 0.40, width, 0.28, font_name="微软雅黑", font_size=15, color="6B7C93", bold=True, align="center")
        _add_textbox(slide, empty_subtitle, left, top + height * 0.54, width, 0.18, font_name="微软雅黑", font_size=9, color="9AA4B2", bold=False, align="center")

    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "position": [left, top],
        "size": [width, height],
        "headers": list(headers),
        "rows": rows,
        "empty_title": empty_title,
        "empty_subtitle": empty_subtitle,
        "row_count": len(rows),
    }
    _persist_semantic_component_anchor(slide, metadata)
    return metadata


def create_selection_timing_grid(
    slide: Slide,
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(11), Inches(4.2)),
    section_titles: Sequence[str],
    level_labels: Sequence[str],
    row_labels: Sequence[str],
    sections: Sequence[Sequence[str]],
    family: str = SELECTION_TIMING_GRID_FAMILY,
) -> dict[str, Any]:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)

    if len(section_titles) != len(sections):
        raise ValueError("section_titles 与 sections 数量必须一致")
    if not row_labels or not level_labels:
        raise ValueError("selection_timing_grid 需要 row_labels 和 level_labels")

    section_count = len(sections)
    gutter = 0.28
    section_w = (width - gutter * (section_count - 1)) / section_count
    row_header_w = min(max(section_w * 0.30, 0.80), 1.50)
    level_row_h = 0.34
    title_h = 0.24
    body_h = height - title_h - level_row_h - 0.08
    row_h = body_h / len(row_labels)

    level_colors = {
        "不显著": ("F5F7FB", "8C95A3"),
        "一般": ("E6EEF9", "5679CC"),
        "强": ("5679CC", "FFFFFF"),
    }

    for sec_idx, section_title in enumerate(section_titles):
        sec_x = left + sec_idx * (section_w + gutter)
        _add_textbox(
            slide,
            section_title,
            sec_x,
            top,
            section_w,
            title_h,
            font_name="微软雅黑",
            font_size=10,
            color="2F3542",
            bold=True,
            align="left",
        )

        value_grid_x = sec_x + row_header_w
        level_w = (section_w - row_header_w) / len(level_labels)

        for level_idx, level in enumerate(level_labels):
            fill, fg = level_colors.get(level, ("EEF3FB", "4B5563"))
            x = value_grid_x + level_idx * level_w
            _add_rect(slide, x, top + title_h + 0.02, level_w, level_row_h, fill=_rgb(fill))
            _add_textbox(
                slide,
                level,
                x,
                top + title_h + 0.10,
                level_w,
                0.14,
                font_name="微软雅黑",
                font_size=9,
                color=fg,
                bold=True,
                align="center",
            )

        section_values = list(sections[sec_idx])
        if len(section_values) != len(row_labels):
            raise ValueError("每个 section 的 values 数量必须与 row_labels 一致")

        for row_idx, row_label in enumerate(row_labels):
            y = top + title_h + level_row_h + 0.06 + row_idx * row_h
            _add_textbox(
                slide,
                row_label,
                sec_x,
                y + row_h * 0.20,
                row_header_w - 0.06,
                row_h * 0.40,
                font_name="微软雅黑",
                font_size=9,
                color="4B5563",
                bold=False,
                align="left",
            )

            for level_idx, level in enumerate(level_labels):
                x = value_grid_x + level_idx * level_w
                active = section_values[row_idx] == level
                fill, fg = level_colors.get(level, ("EEF3FB", "4B5563"))
                if active:
                    draw_fill = fill
                    draw_fg = fg
                else:
                    draw_fill = "FFFFFF"
                    draw_fg = "D1D7E2"
                _add_rect(slide, x, y, level_w, row_h - 0.02, fill=_rgb(draw_fill))
                _add_textbox(
                    slide,
                    level if active else "",
                    x,
                    y + row_h * 0.18,
                    level_w,
                    row_h * 0.42,
                    font_name="微软雅黑",
                    font_size=9,
                    color=draw_fg,
                    bold=active,
                    align="center",
                )

    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "position": [left, top],
        "size": [width, height],
        "section_titles": list(section_titles),
        "level_labels": list(level_labels),
        "row_labels": list(row_labels),
        "sections": [list(section) for section in sections],
        "section_count": len(section_titles),
        "row_count": len(row_labels),
    }
    _persist_semantic_component_anchor(slide, metadata)
    return metadata


def create_holding_detail_panel(
    slide: Slide,
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(11), Inches(4.8)),
    headers: Sequence[str],
    rows: Sequence[Sequence[Any]] | None = None,
    df: pd.DataFrame | None = None,
    subtitle: str | None = None,
    summary_text: str | None = None,
    family: str = HOLDING_DETAIL_FAMILY,
) -> dict[str, Any]:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)
    current_y = top

    if df is not None and rows is None:
        table_df = df.copy()
        if headers:
            table_df = table_df[headers]
        else:
            headers = list(table_df.columns)
        rows = table_df.values.tolist()
    rows = list(rows or [])
    headers = list(headers)

    if subtitle:
        _add_textbox(
            slide,
            subtitle,
            left,
            current_y,
            width,
            0.24,
            font_name="微软雅黑",
            font_size=10,
            color="5F6772",
            bold=True,
            align="left",
        )
        current_y += 0.28

    if summary_text:
        _add_textbox(
            slide,
            summary_text,
            left,
            current_y,
            width,
            0.28,
            font_name="微软雅黑",
            font_size=9,
            color="6B7C93",
            bold=False,
            align="left",
        )
        current_y += 0.34

    remaining_h = height - (current_y - top)
    _add_native_table(
        slide,
        headers=headers,
        rows=rows,
        x=left,
        y=current_y,
        w=width,
        h=min(remaining_h, 0.40 * len(rows) + 0.42),
    )
    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "position": [left, top],
        "size": [width, height],
        "headers": headers,
        "rows": rows,
        "subtitle": subtitle,
        "summary_text": summary_text,
        "row_count": len(rows),
    }
    _persist_semantic_component_anchor(slide, metadata)
    return metadata


def create_dual_chart_panel(
    slide: Slide,
    left_chart_family: str,
    left_chart_kwargs: Mapping[str, Any],
    right_chart_family: str,
    right_chart_kwargs: Mapping[str, Any],
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(11), Inches(4.8)),
    left_title: str | None = None,
    right_title: str | None = None,
    gap_inches: float = 0.28,
    family: str = DUAL_CHART_PANEL_FAMILY,
) -> dict[str, Any]:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)
    half_w = (width - gap_inches) / 2
    title_h = 0.20 if (left_title or right_title) else 0.0
    chart_top = top + title_h
    chart_h = height - title_h

    if left_title:
        _add_textbox(
            slide, left_title, left, top, half_w, 0.18,
            font_name="微软雅黑", font_size=10, color="5F6772", bold=False, align="left"
        )
    if right_title:
        _add_textbox(
            slide, right_title, left + half_w + gap_inches, top, half_w, 0.18,
            font_name="微软雅黑", font_size=10, color="5F6772", bold=False, align="left"
        )

    lk = dict(left_chart_kwargs)
    if isinstance(lk.get("df"), list):
        lk["df"] = pd.DataFrame(lk["df"])
    lk["position"] = (Inches(left), Inches(chart_top))
    lk["size"] = (Inches(half_w), Inches(chart_h))
    left_result = create_semantic_chart(slide=slide, family=left_chart_family, **lk)

    rk = dict(right_chart_kwargs)
    if isinstance(rk.get("df"), list):
        rk["df"] = pd.DataFrame(rk["df"])
    rk["position"] = (Inches(left + half_w + gap_inches), Inches(chart_top))
    rk["size"] = (Inches(half_w), Inches(chart_h))
    right_result = create_semantic_chart(slide=slide, family=right_chart_family, **rk)

    metadata = {
        "chart_family": family,
        "semantic_family": family,
        "position": [left, top],
        "size": [width, height],
        "left_chart_family": left_chart_family,
        "left_chart_kwargs": left_chart_kwargs,
        "left_title": left_title,
        "right_chart_family": right_chart_family,
        "right_chart_kwargs": right_chart_kwargs,
        "right_title": right_title,
        "gap_inches": gap_inches,
    }
    anchor_payload = _persist_semantic_component_anchor(slide, metadata)
    _attach_semantic_parent_to_chart(left_result, family, anchor_payload["semantic_anchor_id"])
    _attach_semantic_parent_to_chart(right_result, family, anchor_payload["semantic_anchor_id"])
    return metadata


def create_semantic_chart(slide: Slide, family: str, **kwargs):
    if family == PERFORMANCE_COMPARE_FAMILY:
        return create_performance_compare_chart(slide=slide, **kwargs)
    if family == DISTRIBUTION_PLUS_HISTORY_FAMILY:
        mode = kwargs.pop("mode", "history")
        if mode == "snapshot":
            return create_distribution_snapshot_chart(slide=slide, **kwargs)
        return create_distribution_history_chart(slide=slide, **kwargs)
    if family == STYLE_BOX_FAMILY:
        return create_style_box_chart(slide=slide, **kwargs)
    if family == STYLE_ALLOCATION_FAMILY:
        return create_style_allocation_chart(slide=slide, **kwargs)
    if family == FACTOR_EXPOSURE_FAMILY:
        return create_factor_exposure_chart(slide=slide, **kwargs)
    if family == SCORE_OVERLAY_FAMILY:
        return create_score_overlay_chart(slide=slide, **kwargs)
    if family == CONCENTRATION_FAMILY:
        return create_concentration_chart(slide=slide, **kwargs)
    if family == EVENT_TIMELINE_FAMILY:
        return create_event_timeline_chart(slide=slide, **kwargs)
    if family == ATTRIBUTION_DECOMPOSITION_FAMILY:
        return create_attribution_decomposition_chart(slide=slide, **kwargs)
    if family == RANKED_TILE_MATRIX_FAMILY:
        return create_ranked_tile_matrix_chart(slide=slide, **kwargs)
    if family == HEATMAP_MATRIX_FAMILY:
        return create_heatmap_matrix_chart(slide=slide, **kwargs)
    if family == TABLE_PLUS_CHART_COMPOSITE_FAMILY:
        return create_table_plus_chart_composite(slide=slide, **kwargs)
    if family == FACTOR_ATTRIBUTION_PANEL_FAMILY:
        return create_factor_attribution_panel(slide=slide, **kwargs)
    if family == REGIME_TABLE_PANEL_FAMILY:
        return create_regime_table_panel(slide=slide, **kwargs)
    if family == MANAGER_TIMELINE_PROFILE_FAMILY:
        return create_manager_timeline_profile(slide=slide, **kwargs)
    if family == AWARD_TIMELINE_PANEL_FAMILY:
        return create_award_timeline_panel(slide=slide, **kwargs)
    if family == SELECTION_TIMING_GRID_FAMILY:
        return create_selection_timing_grid(slide=slide, **kwargs)
    if family == DUAL_CHART_PANEL_FAMILY:
        return create_dual_chart_panel(slide=slide, **kwargs)
    if family == HOLDING_DETAIL_FAMILY:
        return create_holding_detail_panel(slide=slide, **kwargs)
    raise ValueError(f"未知 semantic family: {family}")


def get_semantic_chart_spec(layout_info: dict[str, Any] | None) -> dict[str, Any] | None:
    if not layout_info:
        return None
    metadata = layout_info.get("chart_metadata")
    if not isinstance(metadata, dict):
        return None
    family = metadata.get("semantic_family") or metadata.get("chart_family")
    if family not in SEMANTIC_FAMILY_REGISTRY:
        return None
    return metadata


def parse_semantic_chart_from_layout_info(layout_info: dict[str, Any] | None) -> SemanticChartParseResult | None:
    metadata = get_semantic_chart_spec(layout_info)
    if metadata is None:
        return None
    family = metadata.get("semantic_family") or metadata.get("chart_family")
    return SemanticChartParseResult(
        family=family,
        base_chart_family=metadata.get("base_chart_family"),
        metadata=metadata,
    )


def _normalize_series_entry(entry: Mapping[str, Any]) -> dict[str, Any]:
    return {
        "key": entry["key"],
        "name": entry.get("name", entry["key"]),
        "type": entry.get("type", "line"),
        "axis": entry.get("axis", "primary"),
        "grouping": entry.get("grouping"),
        "role": entry.get("role", "fund"),
    }


def _add_style_box_quadrants(slide: Slide, *, chart, position: tuple, size: tuple) -> None:
    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)
    plot = estimate_chart_plot_area(
        chart,
        left=left,
        top=top,
        width=width,
        height=height,
        family_hint="style_box",
    )
    mid_x = plot["left"] + plot["width"] / 2
    mid_y = plot["top"] + plot["height"] / 2

    vline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(mid_x),
        Inches(plot["top"]),
        Pt(1),
        Inches(plot["height"]),
    )
    vline.fill.solid()
    vline.fill.fore_color.rgb = _rgb("D9DEE7")
    vline.line.fill.background()

    hline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(plot["left"]),
        Inches(mid_y),
        Inches(plot["width"]),
        Pt(1),
    )
    hline.fill.solid()
    hline.fill.fore_color.rgb = _rgb("D9DEE7")
    hline.line.fill.background()


def _add_event_overlays(
    slide: Slide,
    chart,
    df: pd.DataFrame,
    categories_col: str,
    events: Sequence[Mapping[str, Any]],
    *,
    position: tuple,
    size: tuple,
    show_event_labels: bool,
) -> None:
    categories = [str(value) for value in df[categories_col].tolist()]
    if not categories:
        return

    left, top = (_unit_to_inches(value) for value in position)
    width, height = (_unit_to_inches(value) for value in size)
    plot = estimate_chart_plot_area(
        chart,
        left=left,
        top=top,
        width=width,
        height=height,
        family_hint="event_timeline",
    )
    count = max(len(categories), 1)

    for event in events:
        start_idx = _resolve_event_index(event, categories, "start", "start_index")
        end_idx = _resolve_event_index(event, categories, "end", "end_index")
        start_idx = max(0, min(count - 1, start_idx))
        end_idx = max(start_idx, min(count - 1, end_idx))
        band_left = plot["left"] + plot["width"] * (start_idx / count)
        band_width = plot["width"] * ((end_idx - start_idx + 1) / count)

        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(band_left),
            Inches(plot["top"]),
            Inches(band_width),
            Inches(plot["height"]),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = _rgb(event.get("color", "CCD7EB"))
        rect.fill.transparency = 0.72
        rect.line.fill.background()
        _move_shape_behind_chart(slide, rect, chart)

        if show_event_labels and event.get("label"):
            label_box = slide.shapes.add_textbox(
                Inches(band_left),
                Inches(max(top + 0.02, plot["top"] - 0.20)),
                Inches(max(band_width, 0.8)),
                Pt(18),
            )
            text_frame = label_box.text_frame
            text_frame.text = str(event["label"])
            text_frame.paragraphs[0].runs[0].font.size = Pt(8)
            text_frame.paragraphs[0].runs[0].font.name = "黑体"
            text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(event.get("text_color", "4A576A"))


def _resolve_event_index(
    event: Mapping[str, Any],
    categories: Sequence[str],
    value_key: str,
    index_key: str,
) -> int:
    if index_key in event:
        return int(event[index_key])
    if value_key not in event:
        raise ValueError(f"event 缺少 {value_key} 或 {index_key}")
    target = str(event[value_key])
    try:
        return categories.index(target)
    except ValueError as exc:
        raise ValueError(f"event {value_key}={target!r} 不在 categories 中") from exc


def _rgb(value: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(value)


def _unit_to_inches(value: Any) -> float:
    try:
        numeric = float(value)
    except Exception:
        return value
    return numeric / 914400 if numeric > 1000 else numeric


def _move_shape_behind_chart(slide: Slide, shape, chart) -> None:
    chart_frame = None
    for candidate in slide.shapes:
        if getattr(candidate, "has_chart", False) and candidate.chart.part == chart.part:
            chart_frame = candidate
            break

    if chart_frame is None:
        return

    sp_tree = chart_frame._element.getparent()
    shape_el = shape._element
    chart_el = chart_frame._element
    if sp_tree is None or shape_el is None or chart_el is None:
        return

    try:
        sp_tree.remove(shape_el)
        chart_index = list(sp_tree).index(chart_el)
        sp_tree.insert(chart_index, shape_el)
    except Exception:
        return


def _pivot_matrix(df: pd.DataFrame, row_col: str, column_col: str, value_col: str) -> pd.DataFrame:
    for column in (row_col, column_col, value_col):
        if column not in df.columns:
            raise ValueError(f"matrix family 缺少列: {column}")
    pivot = (
        df[[row_col, column_col, value_col]]
        .pivot(index=row_col, columns=column_col, values=value_col)
        .fillna(0)
    )
    return pivot


def _auto_date_axis_config(df: pd.DataFrame, categories_col: str, *, fallback_number_format: str) -> DateAxisConfig:
    row_count = len(df)
    label_interval = max(1, row_count // 10)

    # Prefer coarser month-style labels for long daily histories.
    if row_count >= 500:
        return DateAxisConfig(major_unit=float(max(1, row_count // 12)), number_format="yyyy-mm")
    if row_count >= 180:
        return DateAxisConfig(major_unit=float(max(1, row_count // 10)), number_format="yyyy-mm")
    if row_count >= 60:
        return DateAxisConfig(major_unit=float(label_interval), number_format="yyyy-mm-dd")
    return DateAxisConfig(major_unit=float(max(1, row_count // 8)), number_format=fallback_number_format)


def _interpolate_hex(low_hex: str, high_hex: str, value: float, min_val: float, max_val: float) -> str:
    if max_val <= min_val:
        return high_hex
    ratio = max(0.0, min(1.0, (value - min_val) / (max_val - min_val)))
    lo = tuple(int(low_hex[i : i + 2], 16) for i in (0, 2, 4))
    hi = tuple(int(high_hex[i : i + 2], 16) for i in (0, 2, 4))
    rgb = tuple(round(lo[i] + (hi[i] - lo[i]) * ratio) for i in range(3))
    return "".join(f"{channel:02X}" for channel in rgb)


def _hex_luminance(hex_color: str) -> float:
    r, g, b = (int(hex_color[i : i + 2], 16) / 255 for i in (0, 2, 4))
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _format_matrix_value(value: float, value_format: str) -> str:
    if value_format.endswith("%"):
        return format(value, value_format[:-1]) + "%"
    return format(value, value_format)


def _add_rect(slide: Slide, x: float, y: float, w: float, h: float, *, fill) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()


def _add_textbox(
    slide: Slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_name: str,
    font_size: float,
    color: str,
    bold: bool,
    align: str = "center",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    if align == "left":
        paragraph.alignment = PP_ALIGN.LEFT
    elif align == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    else:
        paragraph.alignment = PP_ALIGN.CENTER
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _add_native_table(
    slide: Slide,
    *,
    headers: Sequence[Any],
    rows: Sequence[Sequence[Any]],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    table = shape.table
    col_width = Inches(w / max(len(headers), 1))
    for idx in range(len(headers)):
        table.columns[idx].width = col_width

    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb("F5F7FB")
        _style_table_cell(cell, bold=True, align="center")

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = "" if value is None else str(value)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("FFFFFF")
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("FAFBFD")
            align = "right" if isinstance(value, (int, float)) else ("left" if col_idx == 0 else "center")
            _style_table_cell(cell, bold=False, align=align)


def _style_table_cell(cell, *, bold: bool, align: str) -> None:
    paragraph = cell.text_frame.paragraphs[0]
    if align == "left":
        paragraph.alignment = PP_ALIGN.LEFT
    elif align == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    else:
        paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.name = "微软雅黑"
    run.font.size = Pt(10)
    run.font.bold = bold
    run.font.color.rgb = _rgb("2F3542")


def _persist_semantic_component_anchor(slide: Slide, metadata: dict[str, Any]) -> dict[str, Any]:
    _, payload = create_semantic_anchor(slide, metadata)
    return payload


def _attach_semantic_parent_to_chart(result: Any, parent_family: str, anchor_id: str) -> None:
    if not hasattr(result, "part"):
        return
    update_chart_semantic_metadata(
        result,
        semantic_parent_family=parent_family,
        semantic_anchor_id=anchor_id,
    )

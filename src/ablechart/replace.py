"""Chart data replacement.

Belongs to: **replace** lifecycle per ADR-0007 §1.
Realises: ADR-0006 §2 (chart data replacement).

Given an existing chart in a .pptx (located via a :class:`ChartSelector` from
:mod:`ablechart.inspect`), replace its categories / series / embedded
workbook in-place while preserving shape identity, position, size, theme,
style and editability.

ADR constraints:

- Replace ≠ rebuild (ADR-0007 §4). This module **must not** delete the
  existing chart shape and reinsert one. Only data is mutated via
  ``chart.replace_data``.
- Technical-layer only (no business slot, no ``user_id``) — ADR-0007 §2.
- Output is written to ``output_pptx``; ``input_pptx`` is read-only.
- Unsupported / mismatched inputs → structured ``ReplaceResult(status="failed",
  error_code=...)`` with one of the codes from ADR-0006 §3. **Never**
  produce a half-written output file (ADR-0006 §5.8).

Supported chart types (ADR-0006 §3 first batch):

- Category-based: ``bar`` / ``line`` / ``combo`` / ``area`` / ``pie``
  → uses :class:`pptx.chart.data.CategoryChartData`
- XY-based: ``scatter`` → uses :class:`pptx.chart.data.XyChartData`
- Bubble: ``bubble`` → uses :class:`pptx.chart.data.BubbleChartData`
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Any, List, Optional

from pptx import Presentation
from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData

from .inspect import ChartSelector, inspect_pptx_charts


# ---------------------------------------------------------------------------
# Public data model
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class SeriesData:
    """One series for chart replacement.

    Field requirements by chart type:

    - **Category** (bar/line/combo/area/pie): only ``name`` + ``values`` (y).
      ``len(values)`` must equal ``len(categories)``.
    - **Scatter**: ``name`` + ``values`` (y) + ``x_values``. Lengths must match.
    - **Bubble**: ``name`` + ``values`` (y) + ``x_values`` + ``size_values``.
      All three lengths must match.
    """

    name: str
    values: List[float]
    x_values: Optional[List[float]] = None      # required for scatter / bubble
    size_values: Optional[List[float]] = None   # required for bubble


@dataclass(frozen=True)
class ReplaceResult:
    """Outcome of :func:`replace_pptx_chart_data`. Always returned;
    failures reported via ``status="failed"`` + ``error_code``, never raised
    (except missing/corrupt input file). See ADR-0006 §5.8."""

    status: str  # "ok" / "failed"
    selector_resolved: Optional[ChartSelector]
    chart_part: Optional[str]
    series_replaced: int = 0
    categories_replaced: int = 0       # 0 for scatter / bubble (no categories)
    data_points_replaced: int = 0      # always populated: points per series × series count
    warnings: List[str] = field(default_factory=list)
    error_code: Optional[str] = None   # one of ADR-0006 §3 codes
    error_detail: Optional[str] = None


# ---------------------------------------------------------------------------
# Internal — chart_type → ChartData class dispatch
# ---------------------------------------------------------------------------

_CATEGORY_REPLACE_TYPES = frozenset({"bar", "line", "combo", "area", "pie"})
_XY_REPLACE_TYPES = frozenset({"scatter"})
_BUBBLE_REPLACE_TYPES = frozenset({"bubble"})
_ALL_REPLACE_TYPES = _CATEGORY_REPLACE_TYPES | _XY_REPLACE_TYPES | _BUBBLE_REPLACE_TYPES


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _find_chart_shape(prs, selector: ChartSelector):
    """Locate the chart shape matching selector. Returns shape or ``None``."""
    if selector.slide_index >= len(prs.slides):
        return None
    slide = prs.slides[selector.slide_index]
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        if shape.shape_id != selector.shape_id:
            continue
        actual_chart_part = str(shape.chart.part.partname).lstrip("/")
        if actual_chart_part == selector.chart_part:
            return shape
    return None


def _failure(
    selector: Optional[ChartSelector],
    error_code: str,
    error_detail: str,
) -> ReplaceResult:
    return ReplaceResult(
        status="failed",
        selector_resolved=selector,
        chart_part=selector.chart_part if selector else None,
        error_code=error_code,
        error_detail=error_detail,
    )


def _validate_category_series(
    categories: Optional[List[Any]],
    series: List[SeriesData],
) -> Optional[tuple[str, str]]:
    """Returns (error_code, error_detail) if invalid, else None."""
    if categories is None:
        return ("categories_required_for_category_chart",
                "category-based chart (bar/line/combo/area/pie) requires categories")
    for s in series:
        if len(s.values) != len(categories):
            return ("category_length_mismatch",
                    f"Series '{s.name}' has {len(s.values)} values but "
                    f"{len(categories)} categories")
    return None


def _validate_xy_series(series: List[SeriesData]) -> Optional[tuple[str, str]]:
    for s in series:
        if s.x_values is None:
            return ("x_values_required_for_scatter",
                    f"Series '{s.name}' missing x_values for scatter chart")
        if len(s.x_values) != len(s.values):
            return ("xy_length_mismatch",
                    f"Series '{s.name}' x_values len {len(s.x_values)} != "
                    f"y values len {len(s.values)}")
    return None


def _validate_bubble_series(series: List[SeriesData]) -> Optional[tuple[str, str]]:
    for s in series:
        if s.x_values is None or s.size_values is None:
            return ("size_values_required_for_bubble",
                    f"Series '{s.name}' missing x_values or size_values for bubble chart")
        n = len(s.values)
        if len(s.x_values) != n or len(s.size_values) != n:
            return ("bubble_length_mismatch",
                    f"Series '{s.name}' x/y/size lengths must match "
                    f"(got x={len(s.x_values)}, y={n}, size={len(s.size_values)})")
    return None


def _build_category_data(categories: List[Any], series: List[SeriesData]):
    data = CategoryChartData()
    data.categories = categories
    for s in series:
        data.add_series(s.name, s.values)
    return data


def _build_xy_data(series: List[SeriesData]):
    data = XyChartData()
    for s in series:
        series_obj = data.add_series(s.name)
        for x, y in zip(s.x_values, s.values):
            series_obj.add_data_point(x, y)
    return data


def _build_bubble_data(series: List[SeriesData]):
    data = BubbleChartData()
    for s in series:
        series_obj = data.add_series(s.name)
        for x, y, size in zip(s.x_values, s.values, s.size_values):
            series_obj.add_data_point(x, y, size)
    return data


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def replace_pptx_chart_data(
    input_pptx: str,
    output_pptx: str,
    selector: ChartSelector,
    categories: Optional[List[Any]] = None,
    series: Optional[List[SeriesData]] = None,
) -> ReplaceResult:
    """Replace categories/series in an existing chart, preserving shape identity.

    Per ADR-0006 §5 the invariants this function must satisfy:

    1. selector resolves to a unique chart (else ``ambiguous_selector``)
    2. shape_id / position / size unchanged
    3. chart_type unchanged
    4. embedded workbook data matches input
    5. chart XML cache matches workbook
    6. parse-after-replace recovers new categories / series
    7. unsupported chart → fail-loud, no partial output file
    8. ADR-0007 §4: must not delete + reinsert chart

    Error codes (ADR-0006 §3 + this implementation):

    - ``ambiguous_selector`` — selector did not resolve
    - ``unsupported_chart_type`` — chart_type not in first-batch supported set
    - ``missing_embedded_workbook`` — chart references external link only
    - ``categories_required_for_category_chart`` — category chart called without categories
    - ``category_length_mismatch`` — series.values length != len(categories)
    - ``x_values_required_for_scatter`` — scatter chart series missing x_values
    - ``xy_length_mismatch`` — scatter chart series x/y length mismatch
    - ``size_values_required_for_bubble`` — bubble chart series missing x_values or size_values
    - ``bubble_length_mismatch`` — bubble chart series x/y/size length mismatch
    """
    if series is None or not series:
        return _failure(selector, "series_required", "series argument must be non-empty")

    # -----------------------------------------------------------------
    # Pre-validate BEFORE touching the file (ADR-0006 §5.8: no half-write)
    # -----------------------------------------------------------------

    # Inspect to get chart_type + replaceability without mutation.
    inventory = inspect_pptx_charts(input_pptx)
    target_item = next(
        (item for item in inventory
         if item.selector.shape_id == selector.shape_id
         and item.selector.chart_part == selector.chart_part
         and item.selector.slide_index == selector.slide_index),
        None,
    )
    if target_item is None:
        return _failure(
            selector,
            "ambiguous_selector",
            f"Selector did not match any chart in {input_pptx}",
        )

    chart_type = target_item.chart_type
    if chart_type not in _ALL_REPLACE_TYPES:
        return _failure(
            selector,
            "unsupported_chart_type",
            f"chart_type='{chart_type}' not in first-batch supported set "
            f"(ADR-0006 §3: bar/line/combo/area/pie/scatter/bubble)",
        )

    if not target_item.has_embedded_workbook:
        return _failure(
            selector,
            "missing_embedded_workbook",
            "chart has no embedded workbook; cannot perform template-safe replace",
        )

    # Type-specific series shape validation
    err: Optional[tuple[str, str]] = None
    if chart_type in _CATEGORY_REPLACE_TYPES:
        err = _validate_category_series(categories, series)
    elif chart_type in _XY_REPLACE_TYPES:
        err = _validate_xy_series(series)
    elif chart_type in _BUBBLE_REPLACE_TYPES:
        err = _validate_bubble_series(series)
    if err is not None:
        return _failure(selector, err[0], err[1])

    # -----------------------------------------------------------------
    # Perform replace
    # -----------------------------------------------------------------

    prs = Presentation(input_pptx)
    shape = _find_chart_shape(prs, selector)
    if shape is None:
        return _failure(
            selector,
            "ambiguous_selector",
            "Selector resolved in inspect but not after Presentation re-load",
        )

    if chart_type in _CATEGORY_REPLACE_TYPES:
        chart_data = _build_category_data(categories, series)
        categories_count = len(categories)
        data_points = categories_count * len(series)
    elif chart_type in _XY_REPLACE_TYPES:
        chart_data = _build_xy_data(series)
        categories_count = 0
        data_points = sum(len(s.values) for s in series)
    else:  # bubble
        chart_data = _build_bubble_data(series)
        categories_count = 0
        data_points = sum(len(s.values) for s in series)

    # python-pptx replace_data: mutates chart XML + embedded workbook in place,
    # does NOT delete the shape (ADR-0007 §4 satisfied by construction).
    shape.chart.replace_data(chart_data)

    prs.save(output_pptx)

    return ReplaceResult(
        status="ok",
        selector_resolved=selector,
        chart_part=selector.chart_part,
        series_replaced=len(series),
        categories_replaced=categories_count,
        data_points_replaced=data_points,
        warnings=[],
    )

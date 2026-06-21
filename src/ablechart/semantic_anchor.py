"""Persistence helpers for semantic-family slide components.

Chart families can persist metadata inside the embedded workbook. Shape/table
families need a slide-level anchor so they can be rediscovered after save/reopen.
"""

from __future__ import annotations

import json
from datetime import date, datetime
from typing import Any, Iterable
from uuid import uuid4

import pandas as pd
from pptx.util import Pt

SEMANTIC_ANCHOR_NAME_PREFIX = "__pptchartengine_semantic__"
SEMANTIC_ANCHOR_VERSION = 1


def create_semantic_anchor(
    slide,
    metadata: dict[str, Any],
    *,
    anchor_id: str | None = None,
):
    """Persist semantic metadata in an off-canvas invisible textbox."""

    anchor_id = anchor_id or uuid4().hex
    payload = dict(metadata)
    payload.setdefault("semantic_anchor_version", SEMANTIC_ANCHOR_VERSION)
    payload.setdefault("semantic_anchor_id", anchor_id)

    shape = slide.shapes.add_textbox(-1000, -1000, Pt(1), Pt(1))
    shape.name = f"{SEMANTIC_ANCHOR_NAME_PREFIX}:{anchor_id}"
    shape.text_frame.text = json.dumps(_json_safe(payload), ensure_ascii=False, separators=(",", ":"))
    shape.fill.background()
    shape.line.fill.background()
    return shape, payload


def read_semantic_anchor(shape) -> dict[str, Any] | None:
    name = getattr(shape, "name", "")
    if not isinstance(name, str) or not name.startswith(f"{SEMANTIC_ANCHOR_NAME_PREFIX}:"):
        return None

    if not hasattr(shape, "text_frame"):
        return None

    raw = shape.text_frame.text or ""
    if not raw.strip():
        return None

    try:
        metadata = json.loads(raw)
    except json.JSONDecodeError:
        return None

    if not isinstance(metadata, dict):
        return None

    return metadata


def iter_semantic_anchors(slide) -> Iterable[tuple[int, Any, dict[str, Any]]]:
    for shape_idx, shape in enumerate(slide.shapes):
        metadata = read_semantic_anchor(shape)
        if metadata is not None:
            yield shape_idx, shape, metadata


def update_chart_semantic_metadata(chart, **patch: Any) -> dict[str, Any] | None:
    """Merge additional semantic metadata into a chart's hidden workbook sheet."""

    from openpyxl import load_workbook
    import io

    chart_part = chart.part
    xlsx_part = chart_part.chart_workbook.xlsx_part
    xlsx_stream = io.BytesIO(xlsx_part.blob)
    wb = load_workbook(xlsx_stream)

    if "_pptchartengine_meta" not in wb.sheetnames:
        return None

    ws = wb["_pptchartengine_meta"]
    metadata_cell = ws["B5"]
    try:
        metadata = json.loads(metadata_cell.value) if metadata_cell.value else {}
    except json.JSONDecodeError:
        metadata = {}

    if not isinstance(metadata, dict):
        metadata = {}

    metadata.update(_json_safe(patch))
    if not metadata.get("chart_family") and ws["B4"].value:
        metadata["chart_family"] = ws["B4"].value

    metadata_cell.value = json.dumps(metadata, ensure_ascii=False, separators=(",", ":"))
    if metadata.get("chart_family"):
        ws["B4"] = metadata["chart_family"]

    output_stream = io.BytesIO()
    wb.save(output_stream)
    xlsx_part._blob = output_stream.getvalue()
    return metadata


def _json_safe(value: Any) -> Any:
    if isinstance(value, pd.DataFrame):
        return [_json_safe(record) for record in value.to_dict(orient="records")]
    if isinstance(value, pd.Series):
        return [_json_safe(item) for item in value.tolist()]
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    if hasattr(value, "to_pydatetime"):
        return value.to_pydatetime().isoformat()
    if isinstance(value, dict):
        return {str(key): _json_safe(item) for key, item in value.items()}
    if isinstance(value, (list, tuple, set)):
        return [_json_safe(item) for item in value]
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value
    return str(value)

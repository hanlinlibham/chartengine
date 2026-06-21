"""Range snapshot chart family with semantic round-trip support."""

from __future__ import annotations

from dataclasses import dataclass
import math
from typing import Any

import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from .api import create_combo_chart
from .layout import ValueAxisConfig
from .oxml_ns import NAMESPACES
from .parser import ChartParser
from .plot_area import estimate_chart_plot_area

RANGE_SNAPSHOT_CHART_FAMILY = "range_snapshot"
RANGE_SNAPSHOT_SPEC_VERSION = 1

BASE_SERIES_KEY = "__base__"
RANGE_SERIES_KEY = "__range__"

DEFAULT_RANGE_COLOR = "5F6772"
DEFAULT_AVERAGE_COLOR = "87A330"
DEFAULT_CURRENT_COLOR = "1E88E5"


@dataclass
class RangeSnapshotParseResult:
    categories_col: str
    min_col: str
    max_col: str
    average_col: str
    current_col: str
    df: pd.DataFrame
    raw_chart_df: pd.DataFrame
    series_config: list[dict]
    layout_info: dict | None
    orientation: str


def prepare_range_snapshot_dataframe(
    df: pd.DataFrame,
    categories_col: str,
    min_col: str,
    max_col: str,
    *,
    axis_break: dict[str, Any] | None = None,
) -> pd.DataFrame:
    """Transform semantic min/max inputs into stacked base/range series."""

    _validate_range_snapshot_input(df, categories_col, min_col, max_col)
    base_values = pd.to_numeric(df[min_col], errors="raise")
    max_values = pd.to_numeric(df[max_col], errors="raise")

    if (max_values < base_values).any():
        raise ValueError("range snapshot requires max_col >= min_col for every row")

    return pd.DataFrame(
        {
            categories_col: df[categories_col].tolist(),
            BASE_SERIES_KEY: [_transform_axis_value(value, axis_break) for value in base_values.tolist()],
            RANGE_SERIES_KEY: [
                _transform_axis_value(max_val, axis_break) - _transform_axis_value(min_val, axis_break)
                for min_val, max_val in zip(base_values.tolist(), max_values.tolist())
            ],
        }
    )


def build_range_snapshot_spec(
    df: pd.DataFrame,
    categories_col: str,
    min_col: str,
    max_col: str,
    average_col: str,
    current_col: str,
    *,
    orientation: str = "vertical",
    range_color: str = DEFAULT_RANGE_COLOR,
    average_color: str = DEFAULT_AVERAGE_COLOR,
    current_color: str = DEFAULT_CURRENT_COLOR,
    number_format: str = "0.0x",
    show_average_ticks: bool = True,
    show_current_markers: bool = True,
    show_current_labels: bool = True,
    padding_ratio: float = 0.08,
    axis_break: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Build a semantic spec for round-tripping range snapshot charts."""

    _validate_range_snapshot_input(df, categories_col, min_col, max_col, average_col, current_col)
    _validate_orientation(orientation)

    return {
        "chart_family": RANGE_SNAPSHOT_CHART_FAMILY,
        "spec_version": RANGE_SNAPSHOT_SPEC_VERSION,
        "categories_col": categories_col,
        "min_col": min_col,
        "max_col": max_col,
        "average_col": average_col,
        "current_col": current_col,
        "orientation": orientation,
        "style": {
            "range_color": range_color,
            "average_color": average_color,
            "current_color": current_color,
            "number_format": number_format,
            "show_average_ticks": show_average_ticks,
            "show_current_markers": show_current_markers,
            "show_current_labels": show_current_labels,
            "padding_ratio": padding_ratio,
        },
        "axis_break": axis_break,
        "series_config": _range_snapshot_series_config(),
        "data_points": [
            {
                "category": _to_metadata_scalar(row[categories_col]),
                "min": _to_metadata_scalar(row[min_col]),
                "max": _to_metadata_scalar(row[max_col]),
                "average": _to_metadata_scalar(row[average_col]),
                "current": _to_metadata_scalar(row[current_col]),
            }
            for _, row in df.iterrows()
        ],
    }


def create_range_snapshot_chart(
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    min_col: str,
    max_col: str,
    average_col: str,
    current_col: str,
    *,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    layout_config=None,
    orientation: str = "vertical",
    range_color: str = DEFAULT_RANGE_COLOR,
    average_color: str = DEFAULT_AVERAGE_COLOR,
    current_color: str = DEFAULT_CURRENT_COLOR,
    number_format: str = "0.0x",
    show_average_ticks: bool = True,
    show_current_markers: bool = True,
    show_current_labels: bool = True,
    label_font_name: str = "微软雅黑",
    padding_ratio: float = 0.08,
    axis_break: dict[str, Any] | None = None,
):
    """Create an editable range snapshot chart using stacked columns plus overlays."""

    _validate_orientation(orientation)

    spec = build_range_snapshot_spec(
        df,
        categories_col,
        min_col,
        max_col,
        average_col,
        current_col,
        orientation=orientation,
        range_color=range_color,
        average_color=average_color,
        current_color=current_color,
        number_format=number_format,
        show_average_ticks=show_average_ticks,
        show_current_markers=show_current_markers,
        show_current_labels=show_current_labels,
        padding_ratio=padding_ratio,
        axis_break=axis_break,
    )
    axis_state = _build_axis_state(
        df,
        min_col=min_col,
        max_col=max_col,
        average_col=average_col,
        current_col=current_col,
        padding_ratio=padding_ratio,
        axis_break=axis_break,
    )
    prepared_df = prepare_range_snapshot_dataframe(
        df,
        categories_col,
        min_col,
        max_col,
        axis_break=axis_state["axis_break"],
    )

    chart = create_combo_chart(
        slide=slide,
        df=prepared_df,
        categories_col=categories_col,
        series_config=_range_snapshot_series_config(),
        position=position,
        size=size,
        layout_config=layout_config,
        metadata=spec,
    )

    _style_range_snapshot_series(chart, range_color)
    _set_range_snapshot_bar_spacing(chart, orientation)
    chart.has_legend = False

    axis_min = axis_state["display_min"]
    axis_max = axis_state["display_max"]
    ValueAxisConfig(
        number_format=number_format,
        has_major_gridlines=not axis_state["segmented"],
        min_value=axis_min,
        max_value=axis_max,
        font_size_pt=9,
        font_name="黑体",
    ).apply_to_chart(chart)

    _apply_range_snapshot_orientation(chart, orientation)
    if axis_state["segmented"]:
        _configure_segmented_axis(chart)

    if show_average_ticks or show_current_markers or show_current_labels or axis_break:
        _add_range_snapshot_overlays(
            slide=slide,
            chart=chart,
            df=df,
            categories_col=categories_col,
            max_col=max_col,
            average_col=average_col,
            current_col=current_col,
            position=position,
            size=size,
            average_color=average_color,
            current_color=current_color,
            label_font_name=label_font_name,
            show_average_ticks=show_average_ticks,
            show_current_markers=show_current_markers,
            show_current_labels=show_current_labels,
            axis_min=axis_min,
            axis_max=axis_max,
            orientation=orientation,
            axis_break=axis_break,
            axis_state=axis_state,
            number_format=number_format,
        )

    return chart


def get_range_snapshot_spec(layout_info: dict[str, Any] | None) -> dict[str, Any] | None:
    """Return the embedded semantic range snapshot spec when present."""

    if not layout_info:
        return None

    chart_metadata = layout_info.get("chart_metadata")
    if not isinstance(chart_metadata, dict):
        return None

    chart_family = chart_metadata.get("chart_family") or layout_info.get("chart_family")
    if chart_family != RANGE_SNAPSHOT_CHART_FAMILY:
        return None

    return chart_metadata


def parse_range_snapshot_chart(chart) -> RangeSnapshotParseResult:
    series_config, chart_df, categories_col, layout_info = ChartParser(chart).parse()
    raw_chart_df = _rename_snapshot_columns_with_series_keys(chart_df, series_config)
    restored = restore_range_snapshot_dataframe(raw_chart_df, layout_info)
    spec = get_range_snapshot_spec(layout_info) or {}

    return RangeSnapshotParseResult(
        categories_col=categories_col,
        min_col=spec.get("min_col") or "min",
        max_col=spec.get("max_col") or "max",
        average_col=spec.get("average_col") or "average",
        current_col=spec.get("current_col") or "current",
        df=restored,
        raw_chart_df=raw_chart_df,
        series_config=series_config,
        layout_info=layout_info,
        orientation=spec.get("orientation") or "vertical",
    )


def parse_range_snapshot_from_pptx(
    pptx_path: str,
    slide_idx: int = 0,
    shape_idx: int = 0,
) -> RangeSnapshotParseResult:
    from pathlib import Path

    pptx_path = str(Path(pptx_path).expanduser().resolve())
    prs = Presentation(pptx_path)
    if slide_idx >= len(prs.slides):
        raise ValueError(f"幻灯片索引 {slide_idx} 超出范围（共 {len(prs.slides)} 张）")

    slide = prs.slides[slide_idx]
    chart = None
    for shape in slide.shapes:
        if hasattr(shape, "chart"):
            if shape_idx == 0:
                chart = shape.chart
                break
            shape_idx -= 1

    if chart is None:
        raise ValueError(f"在幻灯片 {slide_idx} 上未找到图表")

    return parse_range_snapshot_chart(chart)


def restore_range_snapshot_dataframe(
    raw_chart_df: pd.DataFrame,
    layout_info: dict[str, Any] | None,
) -> pd.DataFrame:
    """Recover semantic range snapshot input from a parsed base/range chart."""

    spec = get_range_snapshot_spec(layout_info)
    if spec:
        rows = spec.get("data_points")
        if isinstance(rows, list):
            categories_col = spec.get("categories_col") or raw_chart_df.columns[0]
            min_col = spec.get("min_col") or "min"
            max_col = spec.get("max_col") or "max"
            average_col = spec.get("average_col") or "average"
            current_col = spec.get("current_col") or "current"
            return pd.DataFrame(
                {
                    categories_col: [row.get("category") for row in rows],
                    min_col: [row.get("min") for row in rows],
                    max_col: [row.get("max") for row in rows],
                    average_col: [row.get("average") for row in rows],
                    current_col: [row.get("current") for row in rows],
                }
            )

    missing = {BASE_SERIES_KEY, RANGE_SERIES_KEY} - set(raw_chart_df.columns)
    if missing:
        raise ValueError(f"range snapshot dataframe is missing required columns: {sorted(missing)}")

    categories_col = raw_chart_df.columns[0]
    min_values = pd.to_numeric(raw_chart_df[BASE_SERIES_KEY], errors="coerce")
    max_values = min_values + pd.to_numeric(raw_chart_df[RANGE_SERIES_KEY], errors="coerce")
    return pd.DataFrame(
        {
            categories_col: raw_chart_df[categories_col],
            "min": min_values,
            "max": max_values,
            "average": pd.NA,
            "current": pd.NA,
        }
    )


def _validate_range_snapshot_input(
    df: pd.DataFrame,
    categories_col: str,
    min_col: str,
    max_col: str,
    average_col: str | None = None,
    current_col: str | None = None,
) -> None:
    required = [categories_col, min_col, max_col]
    if average_col:
        required.append(average_col)
    if current_col:
        required.append(current_col)

    missing = [col for col in required if col not in df.columns]
    if missing:
        raise KeyError(f"range snapshot input is missing required columns: {missing}")

    for col in (min_col, max_col, average_col, current_col):
        if col is None:
            continue
        pd.to_numeric(df[col], errors="raise")


def _validate_orientation(orientation: str) -> None:
    if orientation not in {"vertical", "horizontal"}:
        raise ValueError("range_snapshot orientation must be 'vertical' or 'horizontal'")


def _range_snapshot_series_config() -> list[dict[str, Any]]:
    return [
        {"key": BASE_SERIES_KEY, "name": "Base", "type": "bar", "axis": "primary", "grouping": "stacked"},
        {"key": RANGE_SERIES_KEY, "name": "Range", "type": "bar", "axis": "primary", "grouping": "stacked"},
    ]


def _rename_snapshot_columns_with_series_keys(
    chart_df: pd.DataFrame,
    series_config: list[dict],
) -> pd.DataFrame:
    rename_map = {
        series.get("name"): series.get("key")
        for series in series_config
        if series.get("name") and series.get("key")
    }
    return chart_df.rename(columns=rename_map)


def _resolve_snapshot_axis_bounds(
    df: pd.DataFrame,
    *,
    min_col: str,
    max_col: str,
    average_col: str,
    current_col: str,
    padding_ratio: float,
) -> tuple[float, float]:
    values = []
    for col in (min_col, max_col, average_col, current_col):
        values.extend(pd.to_numeric(df[col], errors="raise").tolist())

    y_min = min(values)
    y_max = max(values)
    if y_min == y_max:
        pad = max(abs(y_min) * padding_ratio, 1.0)
    else:
        pad = (y_max - y_min) * padding_ratio

    y_min -= pad
    y_max += pad
    if y_min > 0:
        y_min = max(0.0, y_min - pad)
    return float(y_min), float(y_max)


def _build_axis_state(
    df: pd.DataFrame,
    *,
    min_col: str,
    max_col: str,
    average_col: str,
    current_col: str,
    padding_ratio: float,
    axis_break: dict[str, Any] | None,
) -> dict[str, Any]:
    semantic_values = []
    for col in (min_col, max_col, average_col, current_col):
        semantic_values.extend(pd.to_numeric(df[col], errors="raise").tolist())

    semantic_min = min(semantic_values)
    semantic_max = max(semantic_values)
    normalized_break = _normalize_axis_break(axis_break, semantic_min=semantic_min, semantic_max=semantic_max)
    transformed_values = [_transform_axis_value(value, normalized_break) for value in semantic_values]

    display_min = min(transformed_values)
    display_max = max(transformed_values)
    if display_min == display_max:
        pad = max(abs(display_min) * padding_ratio, 1.0)
    else:
        pad = (display_max - display_min) * padding_ratio

    display_min -= pad
    display_max += pad
    if display_min > 0:
        display_min = max(0.0, display_min - pad)

    return {
        "segmented": normalized_break is not None,
        "axis_break": normalized_break,
        "semantic_min": float(semantic_min),
        "semantic_max": float(semantic_max),
        "display_min": float(display_min),
        "display_max": float(display_max),
    }


def _normalize_axis_break(
    axis_break: dict[str, Any] | None,
    *,
    semantic_min: float,
    semantic_max: float,
) -> dict[str, Any] | None:
    if not axis_break or axis_break.get("value") is None:
        return None

    break_value = float(axis_break["value"])
    if not (semantic_min < break_value < semantic_max):
        return None

    compress_ratio = float(axis_break.get("compress_ratio", 0.24))
    if not (0 < compress_ratio < 1):
        raise ValueError("axis_break.compress_ratio must be between 0 and 1")

    normalized = dict(axis_break)
    normalized["value"] = break_value
    normalized["compress_ratio"] = compress_ratio
    return normalized


def _transform_axis_value(value: float, axis_break: dict[str, Any] | None) -> float:
    numeric = float(value)
    if not axis_break:
        return numeric

    break_value = axis_break["value"]
    if numeric <= break_value:
        return numeric
    return break_value + (numeric - break_value) * axis_break["compress_ratio"]


def _style_range_snapshot_series(chart, range_color: str) -> None:
    series_elements = chart._element.findall(".//c:ser", namespaces=NAMESPACES)
    if len(series_elements) < 2:
        return

    _hide_bar_series(series_elements[0])
    _set_bar_series_color(series_elements[1], range_color)


def _hide_bar_series(ser_element) -> None:
    spPr = ser_element.find("c:spPr", namespaces=NAMESPACES)
    if spPr is None:
        spPr = etree.SubElement(ser_element, f"{{{NAMESPACES['c']}}}spPr")

    for child in list(spPr):
        spPr.remove(child)

    solid_fill = etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}solidFill")
    srgb = etree.SubElement(solid_fill, f"{{{NAMESPACES['a']}}}srgbClr")
    srgb.set("val", "FFFFFF")
    ln = etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}ln")
    solid_fill_ln = etree.SubElement(ln, f"{{{NAMESPACES['a']}}}solidFill")
    srgb_ln = etree.SubElement(solid_fill_ln, f"{{{NAMESPACES['a']}}}srgbClr")
    srgb_ln.set("val", "FFFFFF")


def _set_bar_series_color(ser_element, color: str) -> None:
    spPr = ser_element.find("c:spPr", namespaces=NAMESPACES)
    if spPr is None:
        spPr = etree.SubElement(ser_element, f"{{{NAMESPACES['c']}}}spPr")

    for child in list(spPr):
        spPr.remove(child)

    solid_fill = etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}solidFill")
    srgb = etree.SubElement(solid_fill, f"{{{NAMESPACES['a']}}}srgbClr")
    srgb.set("val", color)
    ln = etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}ln")
    line_fill = etree.SubElement(ln, f"{{{NAMESPACES['a']}}}solidFill")
    line_srgb = etree.SubElement(line_fill, f"{{{NAMESPACES['a']}}}srgbClr")
    line_srgb.set("val", color)


def _add_range_snapshot_overlays(
    *,
    slide: Slide,
    chart,
    df: pd.DataFrame,
    categories_col: str,
    max_col: str,
    average_col: str,
    current_col: str,
    position,
    size,
    average_color: str,
    current_color: str,
    label_font_name: str,
    show_average_ticks: bool,
    show_current_markers: bool,
    show_current_labels: bool,
    axis_min: float,
    axis_max: float,
    orientation: str,
    axis_break: dict[str, Any] | None,
    axis_state: dict[str, Any],
    number_format: str,
) -> None:
    left = _emu_to_inches(position[0])
    top = _emu_to_inches(position[1])
    width = _emu_to_inches(size[0])
    height = _emu_to_inches(size[1])

    categories = df[categories_col].tolist()
    if not categories:
        return

    plot = estimate_chart_plot_area(
        chart,
        left=left,
        top=top,
        width=width,
        height=height,
        family_hint="range_snapshot_horizontal" if orientation == "horizontal" else "range_snapshot_vertical",
    )

    if orientation == "vertical":
        gap = min(0.22, plot["width"] * 0.02)
        band_width = (plot["width"] - gap * (len(categories) - 1)) / len(categories)

        def value_to_y(value: float) -> float:
            display_value = _transform_axis_value(value, axis_state["axis_break"])
            ratio = (display_value - axis_min) / (axis_max - axis_min)
            return plot["top"] + plot["height"] * (1 - ratio)

        if axis_state["segmented"]:
            _add_segmented_axis_guides_and_labels(
                slide=slide,
                plot=plot,
                orientation=orientation,
                axis_state=axis_state,
                number_format=number_format,
                label_font_name=label_font_name,
            )

        for index, (_, row) in enumerate(df.iterrows()):
            x = plot["left"] + index * (band_width + gap)
            center_x = x + band_width / 2
            average_y = value_to_y(float(row[average_col]))
            current_y = value_to_y(float(row[current_col]))

            if show_average_ticks:
                tick_width = max(band_width * 0.72, 0.18)
                _add_rect(
                    slide,
                    center_x - tick_width / 2,
                    average_y - 0.02,
                    tick_width,
                    0.04,
                    fill=RGBColor.from_string(average_color),
                )

            if show_current_markers:
                marker_size = min(max(band_width * 0.26, 0.12), 0.22)
                _add_diamond_marker(slide, center_x, current_y, marker_size, current_color)

            if show_current_labels:
                label_text = _format_snapshot_value(float(row[current_col]))
                label_y = current_y - 0.28
                _add_textbox(
                    slide,
                    label_text,
                    x,
                    label_y,
                    band_width,
                    0.22,
                    font_name=label_font_name,
                    font_size=10,
                    color=current_color,
                    bold=True,
                )
        _add_axis_break_overlays(
            slide=slide,
            df=df,
            categories_col=categories_col,
            max_col=max_col,
            plot=plot,
            orientation=orientation,
            axis_min=axis_min,
            axis_max=axis_max,
            axis_break=axis_break,
        )
        return

    gap = min(0.14, plot["height"] * 0.018)
    band_height = (plot["height"] - gap * (len(categories) - 1)) / len(categories)

    def value_to_x(value: float) -> float:
        display_value = _transform_axis_value(value, axis_state["axis_break"])
        ratio = (display_value - axis_min) / (axis_max - axis_min)
        return plot["left"] + plot["width"] * ratio

    if axis_state["segmented"]:
        _add_segmented_axis_guides_and_labels(
            slide=slide,
            plot=plot,
            orientation=orientation,
            axis_state=axis_state,
            number_format=number_format,
            label_font_name=label_font_name,
        )

    for index, (_, row) in enumerate(df.iterrows()):
        y = plot["top"] + index * (band_height + gap)
        center_y = y + band_height / 2
        average_x = value_to_x(float(row[average_col]))
        current_x = value_to_x(float(row[current_col]))

        if show_average_ticks:
            tick_height = max(band_height * 0.72, 0.14)
            _add_rect(
                slide,
                average_x - 0.02,
                center_y - tick_height / 2,
                0.04,
                tick_height,
                fill=RGBColor.from_string(average_color),
            )

        if show_current_markers:
            marker_size = min(max(band_height * 0.42, 0.12), 0.22)
            _add_diamond_marker(slide, current_x, center_y, marker_size, current_color)

        if show_current_labels:
            label_text = _format_snapshot_value(float(row[current_col]))
            label_x = min(current_x + 0.10, left + width - 0.72)
            _add_textbox(
                slide,
                label_text,
                label_x,
                center_y - 0.12,
                0.62,
                0.24,
                font_name=label_font_name,
                font_size=10,
                color=current_color,
                bold=True,
                align="left",
            )

    _add_axis_break_overlays(
        slide=slide,
        df=df,
        categories_col=categories_col,
        max_col=max_col,
        plot=plot,
        orientation=orientation,
        axis_min=axis_min,
        axis_max=axis_max,
        axis_break=axis_break,
    )

def _add_diamond_marker(slide: Slide, center_x: float, center_y: float, size: float, color: str) -> None:
    marker = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(center_x - size / 2),
        Inches(center_y - size / 2),
        Inches(size),
        Inches(size),
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor.from_string(color)
    marker.line.color.rgb = RGBColor.from_string("FFFFFF")
    marker.line.width = Pt(1)


def _add_axis_break_overlays(
    *,
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    max_col: str,
    plot: dict[str, float],
    orientation: str,
    axis_min: float,
    axis_max: float,
    axis_break: dict[str, Any] | None,
) -> None:
    if not axis_break:
        return

    break_value = axis_break.get("value")
    if break_value is None:
        return

    break_value = float(break_value)
    if not (axis_min <= break_value <= axis_max):
        return

    axis_color = axis_break.get("axis_color", axis_break.get("line_color", "2F2F2F"))
    bar_color = axis_break.get("bar_color", "FFFFFF")
    line_width_pt = float(axis_break.get("line_width_pt", 2.4))
    slash_size = float(axis_break.get("size_inches", 0.22))
    slash_gap = float(axis_break.get("gap_inches", 0.09))
    categories = axis_break.get("categories")
    if categories is not None:
        categories = {str(item) for item in categories}

    max_values = pd.to_numeric(df[max_col], errors="raise")

    if orientation == "vertical":
        gap = min(0.22, plot["width"] * 0.02)
        band_width = (plot["width"] - gap * (len(df) - 1)) / len(df)

        def value_to_y(value: float) -> float:
            ratio = (value - axis_min) / (axis_max - axis_min)
            return plot["top"] + plot["height"] * (1 - ratio)

        break_y = value_to_y(break_value)
        _draw_axis_break(
            slide,
            orientation="vertical",
            center_x=plot["left"],
            center_y=break_y,
            size=slash_size,
            gap=slash_gap,
            color=axis_color,
            width_pt=line_width_pt,
        )
        for index, (_, row) in enumerate(df.iterrows()):
            category = str(row[categories_col])
            exceeds = float(max_values.iloc[index]) > break_value
            if categories is not None:
                exceeds = category in categories
            if not exceeds:
                continue
            x = plot["left"] + index * (band_width + gap) + band_width / 2
            _draw_axis_break(
                slide,
                orientation="vertical",
                center_x=x,
                center_y=break_y,
                size=slash_size,
                gap=slash_gap,
                color=bar_color,
                width_pt=line_width_pt,
            )
        return

    gap = min(0.14, plot["height"] * 0.018)
    band_height = (plot["height"] - gap * (len(df) - 1)) / len(df)

    def value_to_x(value: float) -> float:
        ratio = (value - axis_min) / (axis_max - axis_min)
        return plot["left"] + plot["width"] * ratio

    break_x = value_to_x(break_value)
    _draw_axis_break(
        slide,
        orientation="horizontal",
        center_x=break_x,
        center_y=plot["top"],
        size=slash_size,
        gap=slash_gap,
        color=axis_color,
        width_pt=line_width_pt,
    )
    for index, (_, row) in enumerate(df.iterrows()):
        category = str(row[categories_col])
        exceeds = float(max_values.iloc[index]) > break_value
        if categories is not None:
            exceeds = category in categories
        if not exceeds:
            continue
        center_y = plot["top"] + index * (band_height + gap) + band_height / 2
        _draw_axis_break(
            slide,
            orientation="horizontal",
            center_x=break_x,
            center_y=center_y,
            size=slash_size,
            gap=slash_gap,
            color=bar_color,
            width_pt=line_width_pt,
        )


def _draw_axis_break(
    slide: Slide,
    *,
    orientation: str,
    center_x: float,
    center_y: float,
    size: float,
    gap: float,
    color: str,
    width_pt: float,
) -> None:
    if orientation == "vertical":
        _add_break_slash(
            slide,
            center_x=center_x - gap / 2,
            center_y=center_y + gap / 2,
            dx=size / 2,
            dy=size / 2,
            color=color,
            width_pt=width_pt,
        )
        _add_break_slash(
            slide,
            center_x=center_x + gap / 2,
            center_y=center_y - gap / 2,
            dx=size / 2,
            dy=size / 2,
            color=color,
            width_pt=width_pt,
        )
        return

    _add_break_slash(
        slide,
        center_x=center_x - gap / 2,
        center_y=center_y + gap / 2,
        dx=size / 2,
        dy=size / 2,
        color=color,
        width_pt=width_pt,
    )
    _add_break_slash(
        slide,
        center_x=center_x + gap / 2,
        center_y=center_y + gap / 2,
        dx=size / 2,
        dy=size / 2,
        color=color,
        width_pt=width_pt,
    )


def _add_break_slash(
    slide: Slide,
    *,
    center_x: float,
    center_y: float,
    dx: float,
    dy: float,
    color: str,
    width_pt: float,
) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(center_x - dx),
        Inches(center_y + dy),
        Inches(center_x + dx),
        Inches(center_y - dy),
    )
    connector.line.color.rgb = RGBColor.from_string(color)
    connector.line.width = Pt(width_pt)


def _apply_range_snapshot_orientation(chart, orientation: str) -> None:
    if orientation == "vertical":
        return

    bar_charts = chart._element.findall(".//c:barChart", namespaces=NAMESPACES)
    for bar_chart in bar_charts:
        bar_dir = bar_chart.find("c:barDir", namespaces=NAMESPACES)
        if bar_dir is None:
            bar_dir = etree.SubElement(bar_chart, f"{{{NAMESPACES['c']}}}barDir")
        bar_dir.set("val", "bar")

    cat_axes = chart._element.findall(".//c:catAx", namespaces=NAMESPACES)
    for cat_ax in cat_axes:
        scaling = cat_ax.find("c:scaling", namespaces=NAMESPACES)
        if scaling is not None:
            orientation = scaling.find("c:orientation", namespaces=NAMESPACES)
            if orientation is None:
                orientation = etree.SubElement(scaling, f"{{{NAMESPACES['c']}}}orientation")
            orientation.set("val", "maxMin")
        ax_pos = cat_ax.find("c:axPos", namespaces=NAMESPACES)
        if ax_pos is not None:
            ax_pos.set("val", "l")

    val_axes = chart._element.findall(".//c:valAx", namespaces=NAMESPACES)
    if val_axes:
        ax_pos = val_axes[0].find("c:axPos", namespaces=NAMESPACES)
        if ax_pos is not None:
            ax_pos.set("val", "b")


def _set_range_snapshot_bar_spacing(chart, orientation: str) -> None:
    bar_charts = chart._element.findall(".//c:barChart", namespaces=NAMESPACES)
    gap_width_value = "220" if orientation == "vertical" else "120"
    for bar_chart in bar_charts:
        gap_width = bar_chart.find("c:gapWidth", namespaces=NAMESPACES)
        if gap_width is None:
            gap_width = etree.SubElement(bar_chart, f"{{{NAMESPACES['c']}}}gapWidth")
        gap_width.set("val", gap_width_value)


def _configure_segmented_axis(chart) -> None:
    val_axes = chart._element.findall(".//c:valAx", namespaces=NAMESPACES)
    for axis in val_axes:
        _set_axis_attr(axis, "c:tickLblPos", "val", "none")
        _set_axis_attr(axis, "c:majorTickMark", "val", "none")
        _set_axis_attr(axis, "c:minorTickMark", "val", "none")
        for gridline in axis.findall("c:majorGridlines", namespaces=NAMESPACES):
            axis.remove(gridline)


def _set_axis_attr(axis_element, selector: str, attr: str, value: str) -> None:
    elem = axis_element.find(selector, namespaces=NAMESPACES)
    if elem is not None:
        elem.set(attr, value)


def _add_segmented_axis_guides_and_labels(
    *,
    slide: Slide,
    plot: dict[str, float],
    orientation: str,
    axis_state: dict[str, Any],
    number_format: str,
    label_font_name: str,
) -> None:
    ticks = _build_segmented_ticks(axis_state)
    gridline_color = axis_state["axis_break"].get("gridline_color", "E1E4EA")
    label_color = axis_state["axis_break"].get("label_color", "2F2F2F")

    for tick in ticks:
        if math.isclose(tick, axis_state["axis_break"]["value"], rel_tol=0, abs_tol=1e-6):
            continue

        display_tick = _transform_axis_value(tick, axis_state["axis_break"])
        if orientation == "vertical":
            y = plot["top"] + plot["height"] * (
                1 - (display_tick - axis_state["display_min"]) / (axis_state["display_max"] - axis_state["display_min"])
            )
            _add_rect(slide, plot["left"], y - 0.004, plot["width"], 0.008, fill=RGBColor.from_string(gridline_color))
            _add_textbox(
                slide,
                _format_axis_tick(tick, number_format),
                plot["left"] - 0.52,
                y - 0.12,
                0.44,
                0.24,
                font_name=label_font_name,
                font_size=9,
                color=label_color,
                bold=False,
                align="right",
            )
            continue

        x = plot["left"] + plot["width"] * (
            (display_tick - axis_state["display_min"]) / (axis_state["display_max"] - axis_state["display_min"])
        )
        _add_rect(slide, x - 0.004, plot["top"], 0.008, plot["height"], fill=RGBColor.from_string(gridline_color))
        _add_textbox(
            slide,
            _format_axis_tick(tick, number_format),
            x - 0.30,
            plot["top"] - 0.22,
            0.60,
            0.18,
            font_name=label_font_name,
            font_size=9,
            color=label_color,
            bold=False,
            align="center",
        )


def _build_segmented_ticks(axis_state: dict[str, Any]) -> list[float]:
    axis_break = axis_state["axis_break"]
    custom_tick_values = axis_break.get("tick_values")
    if custom_tick_values:
        return [float(value) for value in custom_tick_values]

    semantic_min = axis_state["semantic_min"]
    semantic_max = axis_state["semantic_max"]
    tick_step = axis_break.get("tick_step")
    if tick_step is None:
        tick_step = _nice_step((semantic_max - semantic_min) / 8 if semantic_max != semantic_min else 1.0)
    tick_step = float(tick_step)

    tick_start = axis_break.get("tick_start")
    if tick_start is None:
        tick_start = round(semantic_min, 1)
    tick_start = float(tick_start)

    ticks: list[float] = []
    tick = tick_start
    guard = 0
    while tick <= semantic_max + (tick_step * 0.001) and guard < 200:
        if tick >= semantic_min - (tick_step * 0.001):
            ticks.append(round(tick, 3))
        tick += tick_step
        guard += 1

    if not ticks:
        ticks.append(round(semantic_min, 3))
        ticks.append(round(semantic_max, 3))
    return ticks


def _nice_step(value: float) -> float:
    if value <= 0:
        return 1.0
    exponent = math.floor(math.log10(value))
    fraction = value / (10 ** exponent)
    if fraction <= 1:
        nice_fraction = 1
    elif fraction <= 2:
        nice_fraction = 2
    elif fraction <= 5:
        nice_fraction = 5
    else:
        nice_fraction = 10
    return nice_fraction * (10 ** exponent)


def _format_axis_tick(value: float, number_format: str) -> str:
    if "%" in number_format:
        return f"{value:.0%}"
    if "x" in number_format.lower():
        return f"{value:.1f}x"
    if ".0" in number_format or ".00" in number_format:
        decimals = 2 if ".00" in number_format else 1
        return f"{value:.{decimals}f}"
    return f"{value:.0f}"


def _format_snapshot_value(value: float) -> str:
    if abs(value) >= 100:
        return f"{value:.0f}"
    if abs(value) >= 10:
        return f"{value:.1f}"
    return f"{value:.2f}"


def _emu_to_inches(value) -> float:
    if hasattr(value, "inches"):
        return float(value.inches)
    return float(value) / 914400.0


def _add_rect(slide: Slide, x: float, y: float, w: float, h: float, *, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()


def _add_textbox(
    slide: Slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_name: str,
    font_size: float,
    color: str,
    bold: bool,
    align: str = "center",
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = {
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
    }.get(align, PP_ALIGN.CENTER)
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.bold = bold


def _to_metadata_scalar(value: Any) -> Any:
    if hasattr(value, "item"):
        try:
            value = value.item()
        except Exception:
            pass
    if hasattr(value, "isoformat"):
        try:
            return value.isoformat()
        except Exception:
            return str(value)
    return value

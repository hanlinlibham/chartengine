"""Scatter / bubble chart family for editable PowerPoint output."""

from __future__ import annotations

from dataclasses import dataclass
import io
from typing import Any

import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import BubbleChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide
from pptx.util import Inches, Pt

from .api import _normalize_position_tuple, _normalize_size_tuple, _write_embedded_metadata
from .oxml_ns import NAMESPACES
from .polish import polish_xy_chart
from .tokens import get_chart_token  # 颜色真源

SCATTER_CHART_FAMILY = "scatter"
BUBBLE_CHART_FAMILY = "bubble"
SCATTER_SPEC_VERSION = 1


@dataclass
class ScatterParseResult:
    chart_family: str
    x_col: str
    y_col: str
    size_col: str | None
    df: pd.DataFrame
    raw_chart_df: pd.DataFrame
    metadata: dict[str, Any]


def create_scatter_chart(
    slide: Slide,
    df: pd.DataFrame,
    x_col: str,
    y_col: str,
    *,
    series_name: str | None = None,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    color: str = None,
    marker_size: int = 9,
):
    _validate_xy_input(df, x_col, y_col)
    position = _normalize_position_tuple(position)
    size = _normalize_size_tuple(size)
    chart_data = XyChartData()
    series = chart_data.add_series(series_name or y_col)
    for _, row in df.iterrows():
        series.add_data_point(float(row[x_col]), float(row[y_col]))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        position[0],
        position[1],
        size[0],
        size[1],
        chart_data,
    ).chart
    _style_xy_or_bubble_chart(chart, color=color, marker_size=marker_size)
    polish_xy_chart(chart, df[x_col].tolist(), df[y_col].tolist())
    _write_embedded_metadata(
        chart,
        x_col,
        [{
            "key": y_col,
            "name": series_name or y_col,
            "type": "scatter",
            "axis": "primary",
            "x_key": x_col,
        }],
        metadata=_build_xy_metadata(
            chart_family=SCATTER_CHART_FAMILY,
            x_col=x_col,
            y_col=y_col,
            size_col=None,
            series_name=series_name or y_col,
            df=df,
        ),
    )
    return chart


def create_bubble_chart(
    slide: Slide,
    df: pd.DataFrame,
    x_col: str,
    y_col: str,
    size_col: str,
    *,
    series_name: str | None = None,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    color: str = None,
    marker_size: int = 9,
):
    _validate_xy_input(df, x_col, y_col, size_col)
    position = _normalize_position_tuple(position)
    size = _normalize_size_tuple(size)
    chart_data = BubbleChartData()
    series = chart_data.add_series(series_name or y_col)
    for _, row in df.iterrows():
        series.add_data_point(float(row[x_col]), float(row[y_col]), float(row[size_col]))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE,
        position[0],
        position[1],
        size[0],
        size[1],
        chart_data,
    ).chart
    _style_xy_or_bubble_chart(chart, color=color, marker_size=marker_size, is_bubble=True)
    polish_xy_chart(chart, df[x_col].tolist(), df[y_col].tolist())
    _write_embedded_metadata(
        chart,
        x_col,
        [{
            "key": y_col,
            "name": series_name or y_col,
            "type": "bubble",
            "axis": "primary",
            "x_key": x_col,
            "size_key": size_col,
        }],
        metadata=_build_xy_metadata(
            chart_family=BUBBLE_CHART_FAMILY,
            x_col=x_col,
            y_col=y_col,
            size_col=size_col,
            series_name=series_name or y_col,
            df=df,
        ),
    )
    return chart


def parse_scatter_chart(chart) -> ScatterParseResult:
    return _parse_xy_like_chart(chart, expected_family=SCATTER_CHART_FAMILY)


def parse_bubble_chart(chart) -> ScatterParseResult:
    return _parse_xy_like_chart(chart, expected_family=BUBBLE_CHART_FAMILY)


def parse_scatter_from_pptx(pptx_path: str, slide_idx: int = 0, shape_idx: int = 0) -> ScatterParseResult:
    return _parse_xy_like_from_pptx(pptx_path, slide_idx, shape_idx, SCATTER_CHART_FAMILY)


def parse_bubble_from_pptx(pptx_path: str, slide_idx: int = 0, shape_idx: int = 0) -> ScatterParseResult:
    return _parse_xy_like_from_pptx(pptx_path, slide_idx, shape_idx, BUBBLE_CHART_FAMILY)


def _parse_xy_like_from_pptx(pptx_path: str, slide_idx: int, shape_idx: int, family: str) -> ScatterParseResult:
    from pathlib import Path

    pptx_path = str(Path(pptx_path).expanduser().resolve())
    prs = Presentation(pptx_path)
    if slide_idx >= len(prs.slides):
        raise ValueError(f"幻灯片索引 {slide_idx} 超出范围（共 {len(prs.slides)} 张）")
    slide = prs.slides[slide_idx]
    chart = None
    for shape in slide.shapes:
        if hasattr(shape, "chart"):
            if shape_idx == 0:
                chart = shape.chart
                break
            shape_idx -= 1
    if chart is None:
        raise ValueError(f"在幻灯片 {slide_idx} 上未找到图表")
    return _parse_xy_like_chart(chart, expected_family=family)


def _parse_xy_like_chart(chart, *, expected_family: str) -> ScatterParseResult:
    wb = load_workbook(io.BytesIO(chart.part.chart_workbook.xlsx_part.blob))
    metadata = _read_xy_metadata(wb)
    family = metadata.get("chart_family")
    if family != expected_family:
        raise ValueError(f"expected chart_family={expected_family}, got {family!r}")

    df = pd.DataFrame(metadata["data_points"])
    df = df.rename(columns={"x": metadata["x_col"], "y": metadata["y_col"], "size": metadata.get("size_col") or "size"})
    if metadata.get("size_col") is None and "size" in df.columns:
        df = df.drop(columns=["size"])

    raw_df = _read_xy_workbook_dataframe(wb, metadata)
    return ScatterParseResult(
        chart_family=family,
        x_col=metadata["x_col"],
        y_col=metadata["y_col"],
        size_col=metadata.get("size_col"),
        df=df,
        raw_chart_df=raw_df,
        metadata=metadata,
    )


def _build_xy_metadata(
    *,
    chart_family: str,
    x_col: str,
    y_col: str,
    size_col: str | None,
    series_name: str,
    df: pd.DataFrame,
) -> dict[str, Any]:
    data_points = []
    for _, row in df.iterrows():
        item = {"x": float(row[x_col]), "y": float(row[y_col])}
        if size_col:
            item["size"] = float(row[size_col])
        data_points.append(item)

    return {
        "chart_family": chart_family,
        "spec_version": SCATTER_SPEC_VERSION,
        "x_col": x_col,
        "y_col": y_col,
        "size_col": size_col,
        "series_name": series_name,
        "data_points": data_points,
    }


def _read_xy_metadata(workbook) -> dict[str, Any]:
    from .parser import METADATA_SHEET_NAME
    import json

    if METADATA_SHEET_NAME not in workbook.sheetnames:
        raise ValueError("missing embedded chart metadata sheet")
    ws = workbook[METADATA_SHEET_NAME]
    metadata = {}
    for row in ws.iter_rows(values_only=True):
        values = list(row)
        if not any(value is not None for value in values):
            continue
        field_name = values[0]
        field_value = values[1] if len(values) > 1 else None
        if field_name == "chart_metadata_json" and field_value:
            metadata = json.loads(field_value)
            break
    return metadata


def _read_xy_workbook_dataframe(workbook, metadata: dict[str, Any]) -> pd.DataFrame:
    ws = workbook.active
    rows = list(ws.iter_rows(values_only=True))
    if metadata["chart_family"] == BUBBLE_CHART_FAMILY:
        columns = [metadata["x_col"], metadata["y_col"], metadata["size_col"]]
        data = rows[1:]
    else:
        columns = [metadata["x_col"], metadata["y_col"]]
        data = rows[1:]
    return pd.DataFrame(data, columns=columns)


def _style_xy_or_bubble_chart(chart, *, color: str, marker_size: int, is_bubble: bool = False) -> None:
    color = color or get_chart_token("scatter_default")
    chart.has_legend = False
    series = chart.series[0]
    # OOXML schema: CT_BubbleSer does not allow <c:marker>; setting it makes
    # PowerPoint flag the file as needing repair. Bubble size is encoded via
    # <c:bubbleSize> alone, so the visible dot is sized from the data, not
    # from a marker style.
    if not is_bubble:
        series.marker.size = marker_size
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor.from_string(color)
    try:
        series.format.line.color.rgb = RGBColor.from_string(color)
        series.format.line.width = Pt(1)
    except Exception:
        pass
    if is_bubble:
        _set_fill_alpha(series, 72)


def _set_fill_alpha(series, alpha_percent: int) -> None:
    """气泡填充加透明度，重叠气泡可分辨（alpha 元素挂在 srgbClr 下）。"""
    from lxml import etree

    ser_el = series._element
    srgb = ser_el.find("c:spPr/a:solidFill/a:srgbClr", namespaces=NAMESPACES)
    if srgb is None:
        return
    alpha = etree.SubElement(srgb, f"{{{NAMESPACES['a']}}}alpha")
    alpha.set("val", str(alpha_percent * 1000))


def _validate_xy_input(df: pd.DataFrame, x_col: str, y_col: str, size_col: str | None = None) -> None:
    required = [x_col, y_col] + ([size_col] if size_col else [])
    missing = [col for col in required if col not in df.columns]
    if missing:
        raise KeyError(f"xy/bubble input is missing required columns: {missing}")

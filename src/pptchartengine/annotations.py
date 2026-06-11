"""GTM 级图表注释层 — 提炼自国际投行市场指南类报告的签名元素。

覆盖的模式（对应该类报告的惯用手法）：

- **均值/参考虚线 + 行内彩色标签**（"Average"、"Trend growth: 2.9%"）
- **目标区间色带**（RBA Target band 2-3% 灰色横带）
- **末点圆点 + 日期/数值标注**（"Aug '25 4.2%"，文字颜色跟随系列色）
- **预测分隔虚线 + 斜纹预测柱**（2025F/2026F 用 ltUpDiag 斜纹与实际值区分）
- **柱上数值标签**（原生 dLbls，可条内白字或条端系列色）
- **单类目高亮**（一组灰柱里把 "S&P 500" 涂成主题色）

几何策略与 waterfall 一致：manualLayout 钉定绘图区 + 读取显式值轴范围，
slide 级 overlay 与图表元素精确对位。数值标签 / 高亮 / 斜纹用原生 XML
（dLbls / dPt），保持可编辑性。
"""

from __future__ import annotations

import math
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .oxml_ns import NAMESPACES
from .polish import AXIS_FONT, pin_plot_area

C = NAMESPACES["c"]
A = NAMESPACES["a"]

ANNOT_LINE_COLOR = "29ABE2"   # GTM 参考线默认青色
BAND_COLOR = "BFBFBF"
FORECAST_PATTERN = "ltUpDiag"


# ============================================================================
# 几何：钉定绘图区 + 坐标换算
# ============================================================================

class PlotGeometry:
    """钉定后的绘图区几何，负责 数据值 ↔ 幻灯片英寸 的换算。"""

    def __init__(self, position, size, fractions, n_categories: int,
                 y_left: Optional[Tuple[float, float]],
                 y_right: Optional[Tuple[float, float]] = None):
        self.left = _emu_to_inches(position[0])
        self.top = _emu_to_inches(position[1])
        self.width = _emu_to_inches(size[0])
        self.height = _emu_to_inches(size[1])
        fx, fy, fw, fh = fractions
        self.plot_left = self.left + fx * self.width
        self.plot_top = self.top + fy * self.height
        self.plot_width = fw * self.width
        self.plot_height = fh * self.height
        self.n = max(1, n_categories)
        self.y_left = y_left
        self.y_right = y_right

    def y(self, value: float, axis: str = "left") -> float:
        rng = self.y_right if axis == "right" else self.y_left
        if rng is None:
            raise ValueError(f"{axis} 轴没有显式范围，无法定位注释")
        vmin, vmax = rng
        # 裁剪到轴范围：越界注释不允许画出绘图区
        ratio = max(0.0, min(1.0, (value - vmin) / (vmax - vmin)))
        return self.plot_top + self.plot_height * (1 - ratio)

    def in_range(self, value: float, axis: str = "left") -> bool:
        rng = self.y_right if axis == "right" else self.y_left
        if rng is None:
            return False
        vmin, vmax = rng
        return vmin <= value <= vmax

    def x_center(self, index: int) -> float:
        slot = self.plot_width / self.n
        return self.plot_left + (index + 0.5) * slot

    def x_slot_left(self, index: int) -> float:
        return self.plot_left + index * (self.plot_width / self.n)


def read_axis_ranges(chart) -> Dict[str, Optional[Tuple[float, float]]]:
    """从 XML scaling 读左右值轴显式范围（polish 保证默认已设置）。"""
    out: Dict[str, Optional[Tuple[float, float]]] = {"left": None, "right": None}
    plot_area = chart._chartSpace.find(f".//{{{C}}}plotArea")
    for ax in plot_area.findall(f"{{{C}}}valAx"):
        pos_el = ax.find(f"{{{C}}}axPos")
        side = "right" if (pos_el is not None and pos_el.get("val") == "r") else "left"
        scaling = ax.find(f"{{{C}}}scaling")
        if scaling is None:
            continue
        vmin = scaling.find(f"{{{C}}}min")
        vmax = scaling.find(f"{{{C}}}max")
        if vmin is not None and vmax is not None:
            out[side] = (float(vmin.get("val")), float(vmax.get("val")))
    return out


def combo_plot_fractions(chart, *, has_secondary: bool) -> Tuple[float, float, float, float]:
    """combo 图钉定比例：按标题/图例位置留边。"""
    top, bottom = 0.06, 0.14
    if chart.has_title:
        top += 0.11
        if len(chart.chart_title.text_frame.paragraphs) > 1:
            top += 0.05  # 副标题行
    try:
        legend_pos = chart.legend.position if chart.has_legend else None
    except Exception:
        legend_pos = None
    if legend_pos is not None:
        from pptx.enum.chart import XL_LEGEND_POSITION
        if legend_pos == XL_LEGEND_POSITION.TOP:
            top += 0.08
        elif legend_pos == XL_LEGEND_POSITION.BOTTOM:
            bottom += 0.08
    left = 0.085
    right = 0.075 if has_secondary else 0.025
    return (left, top, 1.0 - left - right, 1.0 - top - bottom)


# ============================================================================
# 注释渲染（slide overlay）
# ============================================================================

def apply_annotations(
    chart,
    slide,
    *,
    df: pd.DataFrame,
    categories_col: str,
    series_config: List[Dict],
    annotations: List[Dict],
    position,
    size,
    date_format: str = "%Y/%m",
) -> None:
    """钉定绘图区并渲染注释列表。需在 polish（显式轴范围）之后调用。"""
    if not annotations:
        return

    ranges = read_axis_ranges(chart)
    has_secondary = ranges["right"] is not None
    fractions = combo_plot_fractions(chart, has_secondary=has_secondary)
    pin_plot_area(chart, x=fractions[0], y=fractions[1], w=fractions[2], h=fractions[3])

    geo = PlotGeometry(position, size, fractions, len(df), ranges["left"], ranges["right"])

    for ann in annotations:
        kind = str(ann.get("type", "hline")).lower()
        if kind in ("hline", "average", "reference"):
            _draw_hline(slide, geo, ann)
        elif kind == "band":
            _draw_band(slide, geo, ann)
        elif kind == "vline":
            _draw_vline(slide, geo, ann, df, categories_col)
        elif kind == "vband":
            _draw_vband(slide, geo, ann, df, categories_col)
        elif kind == "last_point":
            _draw_last_point(chart, slide, geo, ann, df, series_config, date_format)


def _draw_hline(slide, geo: PlotGeometry, ann: Dict) -> None:
    axis = "right" if str(ann.get("axis", "left")).lower() in ("right", "secondary", "y2") else "left"
    value = float(ann["value"])
    if not geo.in_range(value, axis):
        return  # 越界参考线直接跳过，避免画在图外
    color = ann.get("color") or ANNOT_LINE_COLOR
    dashed = str(ann.get("style", "dashed")).lower() != "solid"
    y = geo.y(value, axis)

    _add_line(slide, geo.plot_left, y, geo.plot_left + geo.plot_width, y,
              color=color, width_pt=1.25, dashed=dashed)

    label = ann.get("label")
    if label:
        # 标签贴线上方，默认偏右（GTM 惯例），颜色与线一致
        at = str(ann.get("label_at", "right")).lower()
        x = {"left": geo.plot_left + 0.05,
             "center": geo.plot_left + geo.plot_width / 2 - 0.8,
             "right": geo.plot_left + geo.plot_width - 1.85}.get(at, geo.plot_left + geo.plot_width - 1.85)
        _add_label(slide, str(label), x, y - 0.215, 1.8, 0.18,
                   color=color, bold=True, align=PP_ALIGN.RIGHT if at == "right" else PP_ALIGN.LEFT)


def _draw_band(slide, geo: PlotGeometry, ann: Dict) -> None:
    axis = "right" if str(ann.get("axis", "left")).lower() in ("right", "secondary", "y2") else "left"
    lo = float(ann.get("from", ann.get("low")))
    hi = float(ann.get("to", ann.get("high")))
    if not (geo.in_range(lo, axis) or geo.in_range(hi, axis)):
        return  # 整个区间都在轴范围外 → 跳过（y() 同时兜底裁剪部分越界）
    color = ann.get("color") or BAND_COLOR
    alpha = int(ann.get("alpha", 30))
    y_top = geo.y(max(lo, hi), axis)
    y_bot = geo.y(min(lo, hi), axis)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(geo.plot_left), Inches(y_top),
        Inches(geo.plot_width), Inches(max(y_bot - y_top, 0.02)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(_norm_hex(color))
    shape.line.fill.background()
    shape.shadow.inherit = False
    _set_shape_alpha(shape, alpha)

    label = ann.get("label")
    if label:
        _add_label(slide, str(label), geo.plot_left + 0.06, y_top + 0.03, 1.6, 0.18,
                   color="595959", bold=True, align=PP_ALIGN.LEFT)


def _draw_vline(slide, geo: PlotGeometry, ann: Dict, df, categories_col) -> None:
    index = ann.get("index")
    if index is None and ann.get("category") is not None:
        matches = df.index[df[categories_col].astype(str) == str(ann["category"])].tolist()
        index = matches[0] if matches else None
    if index is None:
        return
    x = geo.x_slot_left(int(index))
    color = ann.get("color") or "595959"
    dashed = str(ann.get("style", "dashed")).lower() != "solid"
    _add_line(slide, x, geo.plot_top, x, geo.plot_top + geo.plot_height,
              color=color, width_pt=1.0, dashed=dashed)
    label = ann.get("label")
    if label:
        _add_label(slide, str(label), x + 0.04, geo.plot_top, 1.2, 0.18,
                   color=color, bold=False, align=PP_ALIGN.LEFT)


def _draw_vband(slide, geo: PlotGeometry, ann: Dict, df, categories_col) -> None:
    """竖向阴影带（GTM 的衰退期/事件期灰色竖带）。"""
    def to_index(key_cat, key_idx):
        idx = ann.get(key_idx)
        if idx is None and ann.get(key_cat) is not None:
            matches = df.index[df[categories_col].astype(str) == str(ann[key_cat])].tolist()
            idx = matches[0] if matches else None
        return idx

    i_from = to_index("from_category", "from_index")
    i_to = to_index("to_category", "to_index")
    if i_from is None and i_to is None:
        return
    i_from = int(i_from if i_from is not None else 0)
    i_to = int(i_to if i_to is not None else len(df) - 1)
    if i_from > i_to:
        i_from, i_to = i_to, i_from

    x_left = geo.x_slot_left(i_from)
    x_right = geo.x_slot_left(i_to) + geo.plot_width / geo.n
    color = ann.get("color") or BAND_COLOR
    alpha = int(ann.get("alpha", 25))

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_left), Inches(geo.plot_top),
        Inches(max(x_right - x_left, 0.02)), Inches(geo.plot_height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(_norm_hex(color))
    shape.line.fill.background()
    shape.shadow.inherit = False
    _set_shape_alpha(shape, alpha)

    label = ann.get("label")
    if label:
        _add_label(slide, str(label), x_left + 0.04, geo.plot_top + 0.02, 1.4, 0.18,
                   color="595959", bold=True, align=PP_ALIGN.LEFT)


def _draw_last_point(chart, slide, geo: PlotGeometry, ann: Dict, df, series_config, date_format) -> None:
    """末点圆点（原生 dPt marker）+ 「日期 + 数值」标注。"""
    series_name = ann.get("series")
    cfg = next((s for s in series_config if s["name"] == series_name or s["key"] == series_name), None)
    if cfg is None:
        return
    values = pd.to_numeric(df[cfg["key"]], errors="coerce")
    idx = int(values.last_valid_index()) if values.last_valid_index() is not None else len(df) - 1
    value = float(values.iloc[idx])
    axis = "right" if cfg.get("axis") == "secondary" else "left"

    ser = find_series_element(chart, cfg["name"])
    color = ann.get("color") or (_series_color(ser) if ser is not None else None) or ANNOT_LINE_COLOR
    if ser is not None and cfg.get("type", "bar") == "line":
        _add_last_point_marker(ser, idx, color)

    x = geo.x_center(idx)
    y = geo.y(value, axis)

    cat_value = df.iloc[idx][_categories_col_of(df, series_config)]
    if hasattr(cat_value, "strftime"):
        cat_text = cat_value.strftime(date_format)
    else:
        cat_text = str(cat_value)
    value_text = format_value(value, ann.get("format"))
    text = f"{cat_text}  {value_text}" if ann.get("show_category", True) else value_text

    above = bool(ann.get("above", True))
    label_y = y - 0.30 if above else y + 0.08
    _add_label(slide, text, min(x - 1.0, geo.plot_left + geo.plot_width - 1.6), label_y,
               1.6, 0.18, color=color, bold=True, align=PP_ALIGN.RIGHT)


def _categories_col_of(df, series_config) -> str:
    keys = {s["key"] for s in series_config}
    for col in df.columns:
        if col not in keys:
            return col
    return df.columns[0]


# ============================================================================
# 原生 XML 元素：数值标签 / 单点填充 / 斜纹预测 / 图例项隐藏
# ============================================================================

def find_series_element(chart, name: str):
    for ser in chart._chartSpace.findall(f".//{{{C}}}ser"):
        for path in ("c:tx/c:strRef/c:strCache/c:pt/c:v", "c:tx/c:v"):
            node = ser.find(path, namespaces=NAMESPACES)
            if node is not None and node.text == name:
                return ser
    return None


def _series_index(ser) -> int:
    idx = ser.find(f"{{{C}}}idx")
    return int(idx.get("val")) if idx is not None else 0


def _series_color(ser) -> Optional[str]:
    for path in ("c:spPr/a:solidFill/a:srgbClr", "c:spPr/a:ln/a:solidFill/a:srgbClr"):
        node = ser.find(path, namespaces=NAMESPACES)
        if node is not None:
            return node.get("val")
    return None


def _insert_before_cat(ser, element) -> None:
    """按 CT_*Ser 规范，dPt/dLbls 必须位于 <c:cat>/<c:val> 之前。"""
    for tag in ("cat", "val", "xVal"):
        anchor = ser.find(f"{{{C}}}{tag}")
        if anchor is not None:
            ser.insert(list(ser).index(anchor), element)
            return
    ser.append(element)


def add_value_labels(
    chart,
    series_name: str,
    *,
    number_format: Optional[str] = None,
    position: str = "outside",
    color: Optional[str] = None,
    font_size_pt: float = 9,
    bold: bool = True,
    font: str = AXIS_FONT,
) -> None:
    """为某系列添加原生数值标签（c:dLbls）。

    position: outside（条端外，默认）| inside | center（条内白字用 center + color=white）
    color: None 时跟随系列色。
    """
    ser = find_series_element(chart, series_name)
    if ser is None:
        return
    old = ser.find(f"{{{C}}}dLbls")
    if old is not None:
        ser.remove(old)

    parent_tag = ser.getparent().tag.split("}")[-1]
    grouping_el = ser.getparent().find(f"{{{C}}}grouping")
    stacked = grouping_el is not None and grouping_el.get("val") in ("stacked", "percentStacked")
    pos_map = {"outside": "outEnd", "inside": "inEnd", "center": "ctr"}
    dlbl_pos = pos_map.get(position, "outEnd")
    if parent_tag == "barChart" and stacked and dlbl_pos == "outEnd":
        dlbl_pos = "ctr"  # OOXML: 堆叠柱不允许 outEnd
    if parent_tag == "lineChart" and dlbl_pos in ("outEnd", "inEnd"):
        dlbl_pos = "t"

    label_color = _norm_hex(color) if color else (_series_color(ser) or "404040")

    dLbls = etree.Element(f"{{{C}}}dLbls")
    if number_format:
        num_fmt = etree.SubElement(dLbls, f"{{{C}}}numFmt")
        num_fmt.set("formatCode", number_format)
        num_fmt.set("sourceLinked", "0")
    spPr = etree.SubElement(dLbls, f"{{{C}}}spPr")
    etree.SubElement(spPr, f"{{{A}}}noFill")
    ln = etree.SubElement(spPr, f"{{{A}}}ln")
    etree.SubElement(ln, f"{{{A}}}noFill")
    txPr = etree.SubElement(dLbls, f"{{{C}}}txPr")
    etree.SubElement(txPr, f"{{{A}}}bodyPr")
    etree.SubElement(txPr, f"{{{A}}}lstStyle")
    p = etree.SubElement(txPr, f"{{{A}}}p")
    pPr = etree.SubElement(p, f"{{{A}}}pPr")
    defRPr = etree.SubElement(pPr, f"{{{A}}}defRPr")
    defRPr.set("sz", str(int(font_size_pt * 100)))
    if bold:
        defRPr.set("b", "1")
    fill = etree.SubElement(defRPr, f"{{{A}}}solidFill")
    etree.SubElement(fill, f"{{{A}}}srgbClr").set("val", label_color)
    latin = etree.SubElement(defRPr, f"{{{A}}}latin")
    latin.set("typeface", font)
    ea = etree.SubElement(defRPr, f"{{{A}}}ea")
    ea.set("typeface", font)
    etree.SubElement(p, f"{{{A}}}endParaRPr")
    pos_el = etree.SubElement(dLbls, f"{{{C}}}dLblPos")
    pos_el.set("val", dlbl_pos)
    for tag, val in (("showLegendKey", "0"), ("showVal", "1"), ("showCatName", "0"),
                     ("showSerName", "0"), ("showPercent", "0"), ("showBubbleSize", "0")):
        el = etree.SubElement(dLbls, f"{{{C}}}{tag}")
        el.set("val", val)

    _insert_before_cat(ser, dLbls)


def highlight_category(chart, series_name: str, df, categories_col: str,
                       category, color: str) -> None:
    """单类目高亮：一组灰柱中把指定类目涂成强调色（原生 c:dPt）。"""
    matches = df.index[df[categories_col].astype(str) == str(category)].tolist()
    if not matches:
        return
    set_point_fill(chart, series_name, int(matches[0]), color)


def set_point_fill(chart, series_name: str, point_index: int, color: str) -> None:
    ser = find_series_element(chart, series_name)
    if ser is None:
        return
    dPt = etree.Element(f"{{{C}}}dPt")
    idx = etree.SubElement(dPt, f"{{{C}}}idx")
    idx.set("val", str(point_index))
    inv = etree.SubElement(dPt, f"{{{C}}}invertIfNegative")
    inv.set("val", "0")
    spPr = etree.SubElement(dPt, f"{{{C}}}spPr")
    fill = etree.SubElement(spPr, f"{{{A}}}solidFill")
    etree.SubElement(fill, f"{{{A}}}srgbClr").set("val", _norm_hex(color))
    ln = etree.SubElement(spPr, f"{{{A}}}ln")
    etree.SubElement(ln, f"{{{A}}}noFill")
    _insert_before_cat(ser, dPt)


def apply_forecast_pattern(chart, series_name: str, from_index: int, n_points: int) -> None:
    """从 from_index 起的柱体改为斜纹填充（GTM 预测期惯例）。"""
    ser = find_series_element(chart, series_name)
    if ser is None:
        return
    base_color = _series_color(ser) or "595959"
    for i in range(from_index, n_points):
        dPt = etree.Element(f"{{{C}}}dPt")
        idx = etree.SubElement(dPt, f"{{{C}}}idx")
        idx.set("val", str(i))
        inv = etree.SubElement(dPt, f"{{{C}}}invertIfNegative")
        inv.set("val", "0")
        spPr = etree.SubElement(dPt, f"{{{C}}}spPr")
        patt = etree.SubElement(spPr, f"{{{A}}}pattFill")
        patt.set("prst", FORECAST_PATTERN)
        fg = etree.SubElement(patt, f"{{{A}}}fgClr")
        etree.SubElement(fg, f"{{{A}}}srgbClr").set("val", base_color)
        bg = etree.SubElement(patt, f"{{{A}}}bgClr")
        etree.SubElement(bg, f"{{{A}}}srgbClr").set("val", "FFFFFF")
        ln = etree.SubElement(spPr, f"{{{A}}}ln")
        etree.SubElement(ln, f"{{{A}}}noFill")
        _insert_before_cat(ser, dPt)


def delete_legend_entry(chart, series_index: int) -> None:
    """隐藏某系列的图例项（如 range 图的透明基底系列）。"""
    legend = chart._chartSpace.find(f".//{{{C}}}legend")
    if legend is None:
        return
    entry = etree.Element(f"{{{C}}}legendEntry")
    idx = etree.SubElement(entry, f"{{{C}}}idx")
    idx.set("val", str(series_index))
    delete = etree.SubElement(entry, f"{{{C}}}delete")
    delete.set("val", "1")
    pos = legend.find(f"{{{C}}}legendPos")
    if pos is not None:
        legend.insert(list(legend).index(pos) + 1, entry)
    else:
        legend.insert(0, entry)


def style_marker_only_series(ser, *, symbol: str, size: int, color: str,
                             border_color: str = None, border_width_pt: float = 1.0) -> None:
    """折线系列改为「只显示标记点」：线 noFill + 指定 marker（range 图的当前值/均值刻度）。

    border_color: marker 描边色（如白色描边让菱形从深色区间条上跳出来）。
    """
    spPr = ser.find(f"{{{C}}}spPr")
    if spPr is not None:
        ser.remove(spPr)
    spPr = etree.Element(f"{{{C}}}spPr")
    ln = etree.SubElement(spPr, f"{{{A}}}ln")
    etree.SubElement(ln, f"{{{A}}}noFill")
    tx = ser.find(f"{{{C}}}tx")
    ser.insert(list(ser).index(tx) + 1 if tx is not None else 0, spPr)

    old_marker = ser.find(f"{{{C}}}marker")
    if old_marker is not None:
        ser.remove(old_marker)
    marker = etree.Element(f"{{{C}}}marker")
    sym = etree.SubElement(marker, f"{{{C}}}symbol")
    sym.set("val", symbol)
    size_el = etree.SubElement(marker, f"{{{C}}}size")
    size_el.set("val", str(size))
    mk_spPr = etree.SubElement(marker, f"{{{C}}}spPr")
    fill = etree.SubElement(mk_spPr, f"{{{A}}}solidFill")
    etree.SubElement(fill, f"{{{A}}}srgbClr").set("val", _norm_hex(color))
    mk_ln = etree.SubElement(mk_spPr, f"{{{A}}}ln")
    if border_color:
        mk_ln.set("w", str(int(border_width_pt * 12700)))
        ln_fill = etree.SubElement(mk_ln, f"{{{A}}}solidFill")
        etree.SubElement(ln_fill, f"{{{A}}}srgbClr").set("val", _norm_hex(border_color))
    else:
        etree.SubElement(mk_ln, f"{{{A}}}noFill")
    ser.insert(list(ser).index(spPr) + 1, marker)


def _add_last_point_marker(ser, point_index: int, color: str) -> None:
    """末点 dPt：圆形 marker（线系列默认无 marker，仅末点显示）。"""
    dPt = etree.Element(f"{{{C}}}dPt")
    idx = etree.SubElement(dPt, f"{{{C}}}idx")
    idx.set("val", str(point_index))
    inv = etree.SubElement(dPt, f"{{{C}}}invertIfNegative")
    inv.set("val", "0")
    marker = etree.SubElement(dPt, f"{{{C}}}marker")
    sym = etree.SubElement(marker, f"{{{C}}}symbol")
    sym.set("val", "circle")
    size_el = etree.SubElement(marker, f"{{{C}}}size")
    size_el.set("val", "7")
    spPr = etree.SubElement(marker, f"{{{C}}}spPr")
    fill = etree.SubElement(spPr, f"{{{A}}}solidFill")
    etree.SubElement(fill, f"{{{A}}}srgbClr").set("val", _norm_hex(color))
    ln = etree.SubElement(spPr, f"{{{A}}}ln")
    etree.SubElement(ln, f"{{{A}}}noFill")
    _insert_before_cat(ser, dPt)


# ============================================================================
# 工具
# ============================================================================

def format_value(value: float, fmt: Optional[str]) -> str:
    """极简数值格式化：支持 0% / 0.0% / 0.00 / #,##0 等常用码。"""
    if fmt:
        token = fmt.strip()
        if token.endswith("%"):
            decimals = token.count("0", token.find(".")) if "." in token else 0
            return f"{value * 100:.{decimals}f}%"
        if "#,##" in token:
            decimals = len(token.split(".")[1]) if "." in token else 0
            return f"{value:,.{decimals}f}"
        if token.startswith("0") and "." in token:
            decimals = len(token.split(".")[1])
            return f"{value:.{decimals}f}"
        if token == "0":
            return f"{value:.0f}"
    if abs(value) >= 1000:
        return f"{value:,.0f}"
    return f"{value:.2f}".rstrip("0").rstrip(".")


def _norm_hex(color: str) -> str:
    return str(color).lstrip("#").upper()


def _emu_to_inches(value) -> float:
    if hasattr(value, "inches"):
        return float(value.inches)
    return float(value) / 914400.0


def _add_line(slide, x1, y1, x2, y2, *, color: str, width_pt: float, dashed: bool) -> None:
    from pptx.enum.shapes import MSO_CONNECTOR

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = RGBColor.from_string(_norm_hex(color))
    conn.line.width = Pt(width_pt)
    conn.shadow.inherit = False
    if dashed:
        ln = conn.line._get_or_add_ln()
        dash = etree.SubElement(ln, f"{{{A}}}prstDash")
        dash.set("val", "dash")


def _add_label(slide, text, x, y, w, h, *, color, bold, align, font_size: float = 9.5) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = False
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, margin, Pt(0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = AXIS_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(_norm_hex(color))


def _set_shape_alpha(shape, alpha_percent: int) -> None:
    srgb = shape.fill.fore_color._xFill.find(f".//{{{A}}}srgbClr")
    if srgb is not None:
        alpha = etree.SubElement(srgb, f"{{{A}}}alpha")
        alpha.set("val", str(alpha_percent * 1000))

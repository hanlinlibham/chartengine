"""Chart inventory inspection.

Belongs to: **inspect** lifecycle per ADR-0007 §1.
Realises: ADR-0006 §1 (chart inventory).

Given an existing .pptx (possibly authored outside this engine, possibly
without engine metadata), return a structured technical inventory of every
chart shape — without touching data or layout. The inventory is what upper
layers (``pptfi``, ``ablemind``) use to build template manifests and decide
whether a chart is safely replaceable.

ADR constraints:

- Technical-layer only (no business slot, no ``user_id``, no prompt) — ADR-0007 §2
- Read-only (no XML mutation, no .pptx side effects)
- Unsupported charts get ``replaceable=False`` + warning, not exception
- Fail-loud only on missing/corrupt .pptx
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import List, Optional, Tuple

from pptx import Presentation


# ---------------------------------------------------------------------------
# Public data model
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class ChartSelector:
    """Stable identifier for a chart inside a .pptx.

    Selector priority per ADR-0006 §1:

    1. ``explicit_name`` (shape name / business tag) — preferred when available
    2. ``shape_id`` + ``chart_part`` — stable across re-saves
    3. ``(slide_index, chart_index_on_slide)`` — fallback only

    ``slide_index``, ``shape_id``, ``chart_part`` are always populated.
    ``explicit_name`` is populated only when the shape has a non-default name.
    """

    slide_index: int
    shape_id: int
    chart_part: str  # e.g. "ppt/charts/chart1.xml"
    explicit_name: Optional[str] = None


@dataclass(frozen=True)
class ChartInventoryItem:
    """One chart's technical facts. No business semantics (ADR-0007 §2)."""

    selector: ChartSelector
    shape_name: Optional[str]
    chart_index_on_slide: int
    chart_type: str  # "bar", "line", "combo", "scatter", "bubble", "pie", "area", or "unknown"
    category_count: int
    series_count: int
    series_names: List[str] = field(default_factory=list)
    has_embedded_workbook: bool = False
    replaceable: bool = False
    warnings: List[str] = field(default_factory=list)  # e.g. "external_workbook_link"


# ---------------------------------------------------------------------------
# Internal helpers (technical layer — no business semantics)
# ---------------------------------------------------------------------------

# ADR-0006 §3 first-batch supported chart types.
_REPLACEABLE_CHART_TYPES = frozenset({
    "line", "bar", "combo", "area", "pie", "scatter", "bubble",
})

# python-pptx default shape name prefixes. Anything else is treated as a
# user-provided business tag and surfaced as ``explicit_name``.
_DEFAULT_SHAPE_NAME_PREFIXES: Tuple[str, ...] = (
    "Chart ", "Placeholder ", "Picture ", "Group ", "TextBox ", "Title ",
    "Content Placeholder ", "Rectangle ", "Oval ", "Freeform ",
)


def _classify_chart_type(chart) -> str:
    """Map a python-pptx chart to a kernel chart_type string.

    Single plot → that plot's type. Multiple plots, or plots with mixed types,
    → ``"combo"``. Unknown plot class → ``"unknown"`` for that plot, which can
    still yield ``combo`` if mixed.
    """
    plot_types: List[str] = []
    for plot in chart.plots:
        cls = plot.__class__.__name__.lower()
        if "bar" in cls or "column" in cls:
            plot_types.append("bar")
        elif "line" in cls:
            plot_types.append("line")
        elif "area" in cls:
            plot_types.append("area")
        elif "pie" in cls or "doughnut" in cls:
            plot_types.append("pie")
        elif "scatter" in cls or "xy" in cls:
            plot_types.append("scatter")
        elif "bubble" in cls:
            plot_types.append("bubble")
        else:
            plot_types.append("unknown")

    if not plot_types:
        return "unknown"
    if len(set(plot_types)) > 1:
        return "combo"
    return plot_types[0]


def _extract_chart_part(chart) -> str:
    """Return chart part path like ``'ppt/charts/chart1.xml'`` (strip leading ``/``)."""
    return str(chart.part.partname).lstrip("/")


def _has_embedded_workbook(chart) -> bool:
    """True iff chart references an embedded xlsx package (not an external link)."""
    for rel in chart.part.rels.values():
        # Embedded workbook is exposed as an Office Open XML package relationship.
        if "package" in rel.reltype.lower():
            return True
    return False


def _get_category_count(chart) -> int:
    """Number of categories on the first plot. ``0`` for chart types with no
    categorical axis (e.g. pure scatter/bubble)."""
    if not chart.plots:
        return 0
    try:
        cats = chart.plots[0].categories
        if cats is None:
            return 0
        return len(list(cats))
    except (AttributeError, KeyError, ValueError):
        return 0


def _get_series_info(chart) -> Tuple[int, List[str]]:
    """Return ``(series_count, series_names)``. Empty list if chart has no series."""
    series = list(chart.series)
    return len(series), [s.name for s in series]


def _detect_explicit_name(shape_name: Optional[str]) -> Optional[str]:
    """Return shape_name only if it is NOT a python-pptx default like ``'Chart 1'``."""
    if not shape_name:
        return None
    if any(shape_name.startswith(p) for p in _DEFAULT_SHAPE_NAME_PREFIXES):
        return None
    return shape_name


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def inspect_pptx_charts(pptx_path: str) -> List[ChartInventoryItem]:
    """Scan a .pptx and return a chart inventory.

    Per ADR-0006 §1 the inventory item includes:

    - ``slide_index``
    - ``shape_id`` / ``shape_name``
    - ``chart_index_on_slide``
    - ``chart_part`` (e.g. ``"ppt/charts/chart1.xml"``)
    - ``has_embedded_workbook``
    - ``chart_type``
    - ``category_count``, ``series_count``, ``series_names``
    - ``replaceable`` flag (true for ADR-0006 §3 first-batch + has embedded workbook)
    - ``warnings`` (e.g. ``"unsupported_chart_type"``, ``"missing_embedded_workbook"``)

    The returned list is ordered by ``(slide_index, chart_index_on_slide)``
    — natural traversal order.

    Empty .pptx (no slides or no charts) returns an empty list, **not** an exception.
    """
    prs = Presentation(pptx_path)
    inventory: List[ChartInventoryItem] = []

    for slide_index, slide in enumerate(prs.slides):
        chart_index_on_slide = 0
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            chart = shape.chart

            chart_type = _classify_chart_type(chart)
            cat_count = _get_category_count(chart)
            series_count, series_names = _get_series_info(chart)
            has_wb = _has_embedded_workbook(chart)

            warnings: List[str] = []
            replaceable = chart_type in _REPLACEABLE_CHART_TYPES
            if not replaceable:
                warnings.append("unsupported_chart_type")
            if not has_wb:
                warnings.append("missing_embedded_workbook")
                # No workbook → cannot do template-safe replace per ADR-0006 §3.
                replaceable = False

            selector = ChartSelector(
                slide_index=slide_index,
                shape_id=shape.shape_id,
                chart_part=_extract_chart_part(chart),
                explicit_name=_detect_explicit_name(shape.name),
            )
            inventory.append(ChartInventoryItem(
                selector=selector,
                shape_name=shape.name,
                chart_index_on_slide=chart_index_on_slide,
                chart_type=chart_type,
                category_count=cat_count,
                series_count=series_count,
                series_names=series_names,
                has_embedded_workbook=has_wb,
                replaceable=replaceable,
                warnings=warnings,
            ))
            chart_index_on_slide += 1

    return inventory

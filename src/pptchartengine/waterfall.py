"""Waterfall / bridge chart helpers with semantic round-trip support."""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Sequence

import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from .api import create_combo_chart
from .oxml_ns import NAMESPACES
from .parser import ChartParser

DEFAULT_POSITIVE = "10B981"
DEFAULT_NEGATIVE = "EF4444"
DEFAULT_TOTAL = "1E2761"

WATERFALL_CHART_FAMILY = "waterfall"
WATERFALL_SPEC_VERSION = 1

BASE_SERIES_KEY = "__base__"
INCREASE_SERIES_KEY = "__increase__"
DECREASE_SERIES_KEY = "__decrease__"
TOTAL_SERIES_KEY = "__total__"

DEFAULT_RELATIVE_MEASURE = "relative"
TOTAL_MEASURES = {"total", "subtotal", "absolute"}
RELATIVE_MEASURE_ALIASES = {
    "",
    "change",
    "delta",
    "relative",
    "rel",
}
TOTAL_MEASURE_ALIASES = {
    "abs": "absolute",
    "absolute": "absolute",
    "subtotal": "subtotal",
    "sum": "total",
    "total": "total",
}


@dataclass
class WaterfallParseResult:
    categories_col: str
    value_col: str
    measure_col: str | None
    df: pd.DataFrame
    raw_chart_df: pd.DataFrame
    series_config: list[dict]
    layout_info: dict | None


def prepare_waterfall_dataframe(
    df: pd.DataFrame,
    categories_col: str,
    value_col: str,
    *,
    measure_col: str | None = None,
    total_categories: Sequence[str] | None = None,
) -> pd.DataFrame:
    """Transform semantic waterfall input into stacked bridge series."""

    rows = _collect_waterfall_rows(
        df,
        categories_col,
        value_col,
        measure_col=measure_col,
        total_categories=total_categories,
    )

    bridge_rows = []
    cumulative = 0.0

    for row in rows:
        label = row["category"]
        value = row["value"]

        if row["is_total"]:
            bridge_rows.append(
                {
                    categories_col: label,
                    BASE_SERIES_KEY: 0.0,
                    INCREASE_SERIES_KEY: 0.0,
                    DECREASE_SERIES_KEY: 0.0,
                    TOTAL_SERIES_KEY: value,
                }
            )
            cumulative = value
            continue

        if value >= 0:
            base = cumulative
            increase = value
            decrease = 0.0
        else:
            base = cumulative + value
            increase = 0.0
            decrease = abs(value)

        bridge_rows.append(
            {
                categories_col: label,
                BASE_SERIES_KEY: base,
                INCREASE_SERIES_KEY: increase,
                DECREASE_SERIES_KEY: decrease,
                TOTAL_SERIES_KEY: 0.0,
            }
        )
        cumulative += value

    return pd.DataFrame(bridge_rows)


def build_waterfall_spec(
    df: pd.DataFrame,
    categories_col: str,
    value_col: str,
    *,
    measure_col: str | None = None,
    total_categories: Sequence[str] | None = None,
    positive_color: str = DEFAULT_POSITIVE,
    negative_color: str = DEFAULT_NEGATIVE,
    total_color: str = DEFAULT_TOTAL,
    show_legend: bool = False,
) -> dict[str, Any]:
    """Build a semantic spec for round-tripping waterfall charts."""

    rows = _collect_waterfall_rows(
        df,
        categories_col,
        value_col,
        measure_col=measure_col,
        total_categories=total_categories,
    )

    return {
        "chart_family": WATERFALL_CHART_FAMILY,
        "spec_version": WATERFALL_SPEC_VERSION,
        "categories_col": categories_col,
        "value_col": value_col,
        "measure_col": measure_col,
        "total_categories": [_to_metadata_scalar(item) for item in (total_categories or [])],
        "style": {
            "positive_color": positive_color,
            "negative_color": negative_color,
            "total_color": total_color,
            "show_legend": show_legend,
        },
        "series_config": _waterfall_series_config(),
        "data_points": [
            {
                "category": _to_metadata_scalar(row["category"]),
                "value": row["value"],
                "measure": row["measure"],
                "is_total": row["is_total"],
            }
            for row in rows
        ],
    }


def create_waterfall_chart(
    slide: Slide,
    df: pd.DataFrame,
    categories_col: str,
    value_col: str,
    *,
    measure_col: str | None = None,
    total_categories: Sequence[str] | None = None,
    position: tuple = (Inches(1), Inches(2)),
    size: tuple = (Inches(8), Inches(4.5)),
    layout_config=None,
    positive_color: str = DEFAULT_POSITIVE,
    negative_color: str = DEFAULT_NEGATIVE,
    total_color: str = DEFAULT_TOTAL,
    show_legend: bool = False,
    show_connectors: bool = False,
    show_value_labels: bool = True,
    label_font_name: str = "微软雅黑",
):
    """Create an editable waterfall chart using stacked columns.

    ``show_connectors`` defaults to False because the slide-level connector
    overlay is computed against a hardcoded plot-area inset that does not
    match Windows PowerPoint's actual auto-layout, so the rectangles drift
    above the bar tops by ~0.5 chart units. Aligning them properly would
    require pinning the chart's plot area and value-axis range, which is
    a larger refactor tracked separately. Callers who explicitly opt back
    in via ``show_connectors=True`` will still get the (slightly drifted)
    rectangles — that path is preserved for macOS users where the inset
    happened to line up.
    """

    waterfall_spec = build_waterfall_spec(
        df,
        categories_col,
        value_col,
        measure_col=measure_col,
        total_categories=total_categories,
        positive_color=positive_color,
        negative_color=negative_color,
        total_color=total_color,
        show_legend=show_legend,
    )
    waterfall_df = prepare_waterfall_dataframe(
        df,
        categories_col,
        value_col,
        measure_col=measure_col,
        total_categories=total_categories,
    )

    chart = create_combo_chart(
        slide=slide,
        df=waterfall_df,
        categories_col=categories_col,
        series_config=_waterfall_series_config(),
        position=position,
        size=size,
        layout_config=layout_config,
        metadata=waterfall_spec,
    )

    _style_waterfall_series(chart, positive_color, negative_color, total_color)
    chart.has_legend = show_legend
    if show_connectors or show_value_labels:
        _add_waterfall_overlays(
            slide=slide,
            semantic_rows=_collect_waterfall_rows(
                df,
                categories_col,
                value_col,
                measure_col=measure_col,
                total_categories=total_categories,
            ),
            position=position,
            size=size,
            positive_color=positive_color,
            negative_color=negative_color,
            total_color=total_color,
            show_connectors=show_connectors,
            show_value_labels=show_value_labels,
            label_font_name=label_font_name,
        )
    return chart


def get_waterfall_spec(layout_info: dict[str, Any] | None) -> dict[str, Any] | None:
    """Return the embedded semantic waterfall spec when present."""

    if not layout_info:
        return None

    chart_metadata = layout_info.get("chart_metadata")
    if not isinstance(chart_metadata, dict):
        return None

    chart_family = chart_metadata.get("chart_family") or layout_info.get("chart_family")
    if chart_family != WATERFALL_CHART_FAMILY:
        return None

    return chart_metadata


def parse_waterfall_chart(chart) -> WaterfallParseResult:
    series_config, chart_df, categories_col, layout_info = ChartParser(chart).parse()
    bridge_df = _rename_bridge_columns_with_series_keys(chart_df, series_config)
    restored = restore_waterfall_dataframe(bridge_df, layout_info)
    spec = get_waterfall_spec(layout_info) or {}
    return WaterfallParseResult(
        categories_col=categories_col,
        value_col=spec.get("value_col") or "value",
        measure_col=spec.get("measure_col"),
        df=restored,
        raw_chart_df=bridge_df,
        series_config=series_config,
        layout_info=layout_info,
    )


def parse_waterfall_from_pptx(pptx_path: str, slide_idx: int = 0, shape_idx: int = 0) -> WaterfallParseResult:
    from pathlib import Path

    pptx_path = str(Path(pptx_path).expanduser().resolve())
    prs = Presentation(pptx_path)
    if slide_idx >= len(prs.slides):
        raise ValueError(f"幻灯片索引 {slide_idx} 超出范围（共 {len(prs.slides)} 张）")

    slide = prs.slides[slide_idx]
    chart = None
    for shape in slide.shapes:
        if hasattr(shape, "chart"):
            if shape_idx == 0:
                chart = shape.chart
                break
            shape_idx -= 1

    if chart is None:
        raise ValueError(f"在幻灯片 {slide_idx} 上未找到图表")

    return parse_waterfall_chart(chart)


def restore_waterfall_dataframe(
    bridge_df: pd.DataFrame,
    layout_info: dict[str, Any] | None,
) -> pd.DataFrame:
    """Recover semantic waterfall input from a parsed bridge chart."""

    spec = get_waterfall_spec(layout_info)
    if spec:
        return _restore_waterfall_from_spec(bridge_df, spec)

    return _restore_waterfall_from_bridge(bridge_df)


def _collect_waterfall_rows(
    df: pd.DataFrame,
    categories_col: str,
    value_col: str,
    *,
    measure_col: str | None = None,
    total_categories: Sequence[str] | None = None,
) -> list[dict[str, Any]]:
    _validate_input_columns(df, categories_col, value_col, measure_col)

    total_set = set(total_categories or [])
    rows: list[dict[str, Any]] = []

    for _, row in df.iterrows():
        label = row[categories_col]
        value = _coerce_numeric_value(row[value_col], value_col)
        measure_value = row[measure_col] if measure_col else None
        normalized_measure = _normalize_measure(label, measure_value, total_set)
        rows.append(
            {
                "category": label,
                "value": value,
                "measure": normalized_measure,
                "is_total": normalized_measure in TOTAL_MEASURES,
            }
        )

    return rows


def _waterfall_series_config() -> list[dict[str, str]]:
    return [
        {"key": BASE_SERIES_KEY, "name": "Base", "type": "bar", "axis": "primary", "grouping": "stacked"},
        {"key": INCREASE_SERIES_KEY, "name": "Increase", "type": "bar", "axis": "primary", "grouping": "stacked"},
        {"key": DECREASE_SERIES_KEY, "name": "Decrease", "type": "bar", "axis": "primary", "grouping": "stacked"},
        {"key": TOTAL_SERIES_KEY, "name": "Total", "type": "bar", "axis": "primary", "grouping": "stacked"},
    ]


def _validate_input_columns(
    df: pd.DataFrame,
    categories_col: str,
    value_col: str,
    measure_col: str | None,
) -> None:
    missing = [column for column in [categories_col, value_col, measure_col] if column and column not in df.columns]
    if missing:
        raise KeyError(f"waterfall input is missing required columns: {missing}")


def _coerce_numeric_value(value: Any, value_col: str) -> float:
    numeric_value = pd.to_numeric(pd.Series([value]), errors="coerce").iloc[0]
    if pd.isna(numeric_value):
        raise ValueError(f"waterfall value column '{value_col}' must be numeric, got {value!r}")
    return float(numeric_value)


def _normalize_measure(label: Any, measure: Any, total_categories: set[Any]) -> str:
    if label in total_categories:
        return "total"

    if measure is None or (isinstance(measure, str) and not measure.strip()):
        return DEFAULT_RELATIVE_MEASURE

    normalized = str(measure).strip().lower().replace("-", "_").replace(" ", "_")
    if normalized in RELATIVE_MEASURE_ALIASES:
        return DEFAULT_RELATIVE_MEASURE
    if normalized in TOTAL_MEASURE_ALIASES:
        return TOTAL_MEASURE_ALIASES[normalized]

    return normalized


def _restore_waterfall_from_spec(
    bridge_df: pd.DataFrame,
    spec: dict[str, Any],
) -> pd.DataFrame:
    rows = spec.get("data_points")
    if not isinstance(rows, list) or not rows:
        return _restore_waterfall_from_bridge(bridge_df, spec)

    categories_col = spec.get("categories_col") or bridge_df.columns[0]
    value_col = spec.get("value_col") or "value"
    measure_col = spec.get("measure_col")

    restored = {
        categories_col: [row.get("category") for row in rows],
        value_col: [float(row.get("value", 0.0)) for row in rows],
    }
    if measure_col:
        restored[measure_col] = [row.get("measure", DEFAULT_RELATIVE_MEASURE) for row in rows]

    return pd.DataFrame(restored)


def _restore_waterfall_from_bridge(
    bridge_df: pd.DataFrame,
    spec: dict[str, Any] | None = None,
) -> pd.DataFrame:
    categories_col = spec.get("categories_col") if spec else bridge_df.columns[0]
    value_col = spec.get("value_col") if spec else "value"
    measure_col = spec.get("measure_col") if spec else None

    required_columns = {categories_col, INCREASE_SERIES_KEY, DECREASE_SERIES_KEY, TOTAL_SERIES_KEY}
    missing = required_columns - set(bridge_df.columns)
    if missing:
        raise ValueError(f"bridge dataframe is missing waterfall columns: {sorted(missing)}")

    rows = []
    for _, row in bridge_df.iterrows():
        total_value = _safe_float(row.get(TOTAL_SERIES_KEY))
        increase_value = _safe_float(row.get(INCREASE_SERIES_KEY))
        decrease_value = _safe_float(row.get(DECREASE_SERIES_KEY))

        if total_value:
            value = total_value
            measure = "total"
        elif increase_value:
            value = increase_value
            measure = DEFAULT_RELATIVE_MEASURE
        else:
            value = -decrease_value
            measure = DEFAULT_RELATIVE_MEASURE

        restored_row = {
            categories_col: row[categories_col],
            value_col: value,
        }
        if measure_col:
            restored_row[measure_col] = measure
        rows.append(restored_row)

    return pd.DataFrame(rows)


def _safe_float(value: Any) -> float:
    if value is None or pd.isna(value):
        return 0.0
    return float(value)


def _rename_bridge_columns_with_series_keys(
    bridge_df: pd.DataFrame,
    series_config: list[dict],
) -> pd.DataFrame:
    rename_map = {
        series.get("name"): series.get("key")
        for series in series_config
        if series.get("name") and series.get("key")
    }
    return bridge_df.rename(columns=rename_map)


def _to_metadata_scalar(value: Any) -> Any:
    if hasattr(value, "item"):
        try:
            value = value.item()
        except Exception:
            pass
    if hasattr(value, "isoformat"):
        try:
            return value.isoformat()
        except Exception:
            return str(value)
    return value


def _style_waterfall_series(chart, positive_color: str, negative_color: str, total_color: str) -> None:
    series_elements = chart._element.findall(".//c:ser", namespaces=NAMESPACES)
    if len(series_elements) < 4:
        return

    _hide_bar_series(series_elements[0])
    _set_bar_series_color(series_elements[1], positive_color)
    _set_bar_series_color(series_elements[2], negative_color)
    _set_bar_series_color(series_elements[3], total_color)


def _hide_bar_series(ser_element) -> None:
    spPr = ser_element.find("c:spPr", namespaces=NAMESPACES)
    if spPr is None:
        spPr = etree.SubElement(ser_element, f"{{{NAMESPACES['c']}}}spPr")

    for child in list(spPr):
        spPr.remove(child)

    # OOXML CT_ShapeProperties requires the fill choice (noFill/solidFill/...)
    # to precede <a:ln>. With the order swapped, PowerPoint silently ignores
    # the fill and falls back to the theme color, making the "invisible" Base
    # series visible — which breaks the waterfall stacked-bridge illusion.
    etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}noFill")
    ln = etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}ln")
    etree.SubElement(ln, f"{{{NAMESPACES['a']}}}noFill")


def _set_bar_series_color(ser_element, color: str) -> None:
    spPr = ser_element.find("c:spPr", namespaces=NAMESPACES)
    if spPr is None:
        spPr = etree.SubElement(ser_element, f"{{{NAMESPACES['c']}}}spPr")

    for child in list(spPr):
        spPr.remove(child)

    solid_fill = etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}solidFill")
    srgb = etree.SubElement(solid_fill, f"{{{NAMESPACES['a']}}}srgbClr")
    srgb.set("val", color)
    ln = etree.SubElement(spPr, f"{{{NAMESPACES['a']}}}ln")
    etree.SubElement(ln, f"{{{NAMESPACES['a']}}}noFill")


def _add_waterfall_overlays(
    *,
    slide: Slide,
    semantic_rows: list[dict[str, Any]],
    position,
    size,
    positive_color: str,
    negative_color: str,
    total_color: str,
    show_connectors: bool,
    show_value_labels: bool,
    label_font_name: str,
) -> None:
    left = _emu_to_inches(position[0])
    top = _emu_to_inches(position[1])
    width = _emu_to_inches(size[0])
    height = _emu_to_inches(size[1])

    plot_left = left + 0.55
    plot_top = top + 0.35
    plot_width = max(width - 1.0, 1.0)
    plot_height = max(height - 0.95, 1.0)

    value_bounds = [0.0]
    cumulative = 0.0
    cumulative_before_rows = []
    for row in semantic_rows:
        cumulative_before_rows.append(cumulative)
        if row["is_total"]:
            cumulative = row["value"]
        else:
            cumulative += row["value"]
        value_bounds.extend([cumulative_before_rows[-1], cumulative])

    y_max = max(value_bounds) if value_bounds else 0.0
    y_min = min(value_bounds) if value_bounds else 0.0
    if y_min >= 0:
        y_min = 0.0
    padding = (y_max - y_min) * 0.12 if y_max != y_min else max(abs(y_max) * 0.12, 1.0)
    y_max += padding
    y_min -= padding
    if y_max == y_min:
        y_max += 1.0
        y_min -= 1.0

    def value_to_y(value: float) -> float:
        ratio = (value - y_min) / (y_max - y_min)
        return plot_top + plot_height * (1 - ratio)

    n = len(semantic_rows)
    if n == 0:
        return
    gap = min(0.22, plot_width * 0.02)
    bar_width = (plot_width - gap * (n - 1)) / n

    previous_end_value = None
    previous_right = None
    for index, row in enumerate(semantic_rows):
        x = plot_left + index * (bar_width + gap)
        start_value = cumulative_before_rows[index]
        end_value = row["value"] if row["is_total"] else start_value + row["value"]

        if row["is_total"]:
            top_value = max(end_value, 0.0)
            bottom_value = min(end_value, 0.0)
            color = total_color
        else:
            top_value = max(start_value, end_value)
            bottom_value = min(start_value, end_value)
            color = positive_color if row["value"] >= 0 else negative_color

        top_y = value_to_y(top_value)
        bottom_y = value_to_y(bottom_value)

        if show_connectors and previous_end_value is not None and previous_right is not None:
            connector_y = value_to_y(previous_end_value)
            _add_rect(
                slide,
                previous_right,
                connector_y - 0.006,
                gap,
                0.012,
                fill=RGBColor.from_string("B0B7C3"),
            )

        if show_value_labels:
            label_text = _format_waterfall_value(row["value"], row["is_total"])
            label_y = top_y - 0.28 if (row["is_total"] or row["value"] >= 0) else bottom_y + 0.04
            _add_textbox(
                slide,
                label_text,
                x,
                label_y,
                bar_width,
                0.22,
                font_name=label_font_name,
                font_size=11,
                color=color,
                bold=True,
            )

        previous_end_value = end_value
        previous_right = x + bar_width


def _format_waterfall_value(value: float, is_total: bool) -> str:
    if value > 0 and not is_total:
        return f"+{value:.2f}"
    return f"{value:.2f}"


def _emu_to_inches(value) -> float:
    if hasattr(value, "inches"):
        return float(value.inches)
    return float(value) / 914400.0


def _add_rect(slide: Slide, x: float, y: float, w: float, h: float, *, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()


def _add_textbox(
    slide: Slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_name: str,
    font_size: float,
    color: str,
    bold: bool,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.bold = bold

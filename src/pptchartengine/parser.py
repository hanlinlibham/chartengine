"""
图表解析器 - ChartBuilder 的反向操作

这是 ChartBuilder 的逆向工程：
- ChartBuilder: Config -> XML
- ChartParser: XML -> Config

核心职责：
1. 解析坐标轴映射（主轴/次轴）
2. 解析绘图区中的所有系列
3. 提取嵌入的 Excel 数据
4. 生成可用于 ChartBuilder 的配置
"""

import json
from typing import Dict, List, Tuple, Optional
from lxml import etree
from pptx import Presentation
from pptx.chart.chart import Chart
import pandas as pd
import io
from openpyxl import load_workbook

from ._log import debug_print as print
from .oxml_ns import NAMESPACES
from .metadata import METADATA_SHEET_NAME
from .semantic_anchor import iter_semantic_anchors


class ChartParser:
    """
    图表解析器 - 从现有图表提取配置
    
    使用场景：
    1. 从"黄金标准"PPTX 提取图表配置
    2. 批量复现图表
    3. 数据更新和刷新
    """
    
    def __init__(self, chart: Chart):
        """
        初始化解析器
        
        Args:
            chart: python-pptx Chart 对象
        """
        self.chart = chart
        self.chart_element = chart._element
        self.plotArea = self.chart_element.find('.//c:plotArea', namespaces=NAMESPACES)
        self._embedded_metadata: Dict[str, object] = {}
        
        # 解析坐标轴映射
        self.axis_map = self._parse_axes()
        
        print("\n📖 ChartParser 初始化完成")
        print(f"  - 发现坐标轴: {len(self.axis_map)} 个")
        for ax_id, ax_info in self.axis_map.items():
            print(f"    • {ax_id}: {ax_info['type']} ({ax_info.get('position', 'N/A')})")
    
    def _parse_axes(self) -> Dict[str, Dict]:
        """
        解析所有坐标轴，建立 ID 映射表
        
        Returns:
            {
                "123456": {"type": "category", "position": None},
                "789001": {"type": "primary", "position": "l"},
                "789002": {"type": "secondary", "position": "r"}
            }
        """
        axis_map = {}
        
        cat_axes = self.plotArea.findall('.//c:catAx', namespaces=NAMESPACES)
        has_category_axis = bool(cat_axes)

        # 1. 解析所有值轴 (Y轴 / XY 图表中的 X、Y 轴)
        val_axes = self.plotArea.findall('.//c:valAx', namespaces=NAMESPACES)
        for val_ax in val_axes:
            ax_id_elem = val_ax.find('.//c:axId', namespaces=NAMESPACES)
            if ax_id_elem is None:
                continue
            
            ax_id = ax_id_elem.get('val')
            
            # 查找轴的位置
            ax_pos_elem = val_ax.find('.//c:axPos', namespaces=NAMESPACES)
            position = ax_pos_elem.get('val') if ax_pos_elem is not None else 'l'
            
            if not has_category_axis and position in ('b', 't'):
                axis_map[ax_id] = {"type": "category", "position": position}
            elif position == 'r':
                axis_map[ax_id] = {"type": "secondary", "position": position}
            else:
                axis_map[ax_id] = {"type": "primary", "position": position}
        
        # 2. 解析所有分类轴 (X轴)
        for cat_ax in cat_axes:
            ax_id_elem = cat_ax.find('.//c:axId', namespaces=NAMESPACES)
            if ax_id_elem is None:
                continue
            
            ax_id = ax_id_elem.get('val')
            axis_map[ax_id] = {"type": "category", "position": None}
        
        return axis_map
    
    def parse(self) -> Tuple[List[Dict], pd.DataFrame, str, Optional[Dict]]:
        """
        解析图表配置和数据
        
        Returns:
            (series_config, df, categories_col, layout_info)
            
            series_config: 系列配置列表
            df: 数据 DataFrame
            categories_col: 分类列名
            layout_info: 布局信息（图例、轴配置等）
        """
        print("\n" + "=" * 80)
        print("🔍 开始解析图表")
        print("=" * 80)
        
        # 1. 提取嵌入的 Excel 数据
        df, categories_col = self._extract_embedded_data()
        
        print(f"\n📊 数据提取完成:")
        print(f"  - 行数: {len(df)}")
        print(f"  - 列数: {len(df.columns)}")
        print(f"  - 分类列: {categories_col}")
        print(f"  - 数据列: {list(df.columns[1:])}")
        
        # 2. 解析系列配置
        series_config = self._parse_series_config()
        
        print(f"\n📋 系列配置解析完成:")
        for i, cfg in enumerate(series_config):
            print(f"  {i+1}. {cfg['name']}: type={cfg['type']}, axis={cfg['axis']}, key={cfg['key']}")
        
        # 3. ⭐ 新增：解析布局配置
        layout_info = self._parse_layout_info()
        layout_info = layout_info or {}
        self._attach_embedded_metadata(layout_info)

        if layout_info:
            print(f"\n⚙️ 布局配置:")
            if 'legend' in layout_info:
                leg = layout_info['legend']
                print(f"  - 图例: 位置={leg.get('position')}, 字体={leg.get('font_size_pt')}pt")
            if 'category_axis' in layout_info:
                cat = layout_info['category_axis']
                print(f"  - 横轴: 类型={cat.get('type')}, 格式={cat.get('number_format')}")
        
        print("\n" + "=" * 80)
        print("✅ 图表解析完成！")
        print("=" * 80)
        
        return series_config, df, categories_col, layout_info
    
    def _parse_layout_info(self) -> Dict:
        """
        解析布局信息（图例、轴配置、标题、尺寸等）
        
        Returns:
            布局信息字典
        """
        layout_info = {}
        
        try:
            # 1. 解析图表标题
            try:
                if self.chart.has_title:
                    title_frame = self.chart.chart_title.text_frame
                    if title_frame and title_frame.text:
                        layout_info["title"] = {
                            "text": title_frame.text,
                            "has_title": True,
                        }
            except Exception as e:
                print(f"  ⚠️ 标题解析失败: {e}")
            
            # 2. 解析图例
            if self.chart.has_legend:
                legend = self.chart.legend
                layout_info["legend"] = {
                    "position": legend.position,
                    "font_size_pt": legend.font.size.pt if legend.font.size else None,
                    "include_in_layout": legend.include_in_layout,
                }
            
            # 3. 解析横轴（分类轴）
            try:
                category_axis = self.chart.category_axis
                layout_info["category_axis"] = {
                    "type": category_axis.category_type,
                    "major_unit": category_axis.major_unit if hasattr(category_axis, 'major_unit') else None,
                    "number_format": category_axis.tick_labels.number_format if hasattr(category_axis.tick_labels, 'number_format') else None,
                    "font_size_pt": category_axis.tick_labels.font.size.pt if category_axis.tick_labels.font.size else None,
                }
            except Exception as e:
                print(f"  ⚠️ 横轴解析失败: {e}")
            
            # 4. 解析纵轴（值轴）
            try:
                value_axis = self.chart.value_axis
                layout_info["value_axis"] = {
                    "number_format": value_axis.tick_labels.number_format if hasattr(value_axis.tick_labels, 'number_format') else None,
                    "font_size_pt": value_axis.tick_labels.font.size.pt if value_axis.tick_labels.font.size else None,
                    "has_major_gridlines": value_axis.has_major_gridlines if hasattr(value_axis, 'has_major_gridlines') else None,
                }
            except Exception as e:
                print(f"  ⚠️ 纵轴解析失败: {e}")
                
        except Exception as e:
            print(f"  ⚠️ 布局信息解析失败: {e}")
        
        return layout_info

    def _attach_embedded_metadata(self, layout_info: Dict) -> None:
        if not self._embedded_metadata:
            return

        embedded_metadata = {
            key: value
            for key, value in self._embedded_metadata.items()
            if key not in {"categories_col", "series"} and value not in (None, [], {})
        }

        if embedded_metadata:
            layout_info["embedded_metadata"] = embedded_metadata

        chart_family = self._embedded_metadata.get("chart_family")
        if chart_family:
            layout_info["chart_family"] = chart_family

        chart_metadata = self._embedded_metadata.get("chart_metadata")
        if chart_metadata is not None:
            layout_info["chart_metadata"] = chart_metadata

    def _extract_embedded_data(self) -> Tuple[pd.DataFrame, str]:
        """
        提取嵌入在图表中的 Excel 数据
        
        Returns:
            (df, categories_col)
        """
        try:
            # 方法1: 尝试从 chart 的内部部分获取 Excel 数据
            # python-pptx 的图表对象包含嵌入的 Excel workbook
            chart_part = self.chart.part
            
            # 查找 Excel workbook 关系
            for rel in chart_part.rels.values():
                if 'package' in rel.reltype or 'oleObject' in rel.reltype or 'embeddings' in rel.reltype:
                    try:
                        xlsx_blob = rel.target_part.blob
                        wb = load_workbook(io.BytesIO(xlsx_blob))
                        self._embedded_metadata = self._read_embedded_metadata(wb)

                        if self._is_xy_chart():
                            print(f"  ✅ XY 图表改用 XML + 元数据提取数据")
                            return self._extract_xy_data_from_xml()

                        ws = wb.active
                        
                        # 转换为 DataFrame
                        data = []
                        for row in ws.iter_rows(values_only=True):
                            data.append(row)
                        
                        if data:
                            # 第一行是表头
                            df = pd.DataFrame(data[1:], columns=data[0])
                            df, categories_col = self._normalize_embedded_dataframe(
                                df,
                                self._embedded_metadata.get("categories_col"),
                            )
                            print(f"  ✅ 成功从嵌入的 Excel 提取数据")
                            return df, categories_col
                    except Exception:
                        continue
            
            # 如果上面都失败，尝试从 XML 中提取
            print(f"  ⚠️ 未找到嵌入的 Excel 数据，从 XML 提取")
            return self._extract_data_from_xml()
            
        except Exception as e:
            print(f"  ⚠️ 提取嵌入数据失败: {e}")
            # 如果失败，尝试从 XML 中提取数据点
            return self._extract_data_from_xml()

    def _normalize_embedded_dataframe(
        self,
        df: pd.DataFrame,
        metadata_categories_col: Optional[str] = None,
    ) -> Tuple[pd.DataFrame, str]:
        """Normalize embedded workbook data for reverse parsing.

        PowerPoint embedded chart workbooks often leave the first header cell
        blank, which otherwise makes `categories_col` come back as `None`.
        """

        if df.empty:
            return df, "分类"

        columns = list(df.columns)
        first_column = columns[0] if columns else None

        if self._is_missing_header(first_column):
            inferred = metadata_categories_col or self._infer_categories_col_name(df.iloc[:, 0])
            columns[0] = inferred
            df.columns = columns
            return df, inferred

        return df, str(first_column)

    def _read_embedded_metadata(self, workbook) -> Dict[str, object]:
        if METADATA_SHEET_NAME not in workbook.sheetnames:
            return {}

        ws = workbook[METADATA_SHEET_NAME]
        metadata: Dict[str, object] = {"series": []}
        in_series_section = False

        for row in ws.iter_rows(values_only=True):
            values = list(row)

            if not any(value is not None for value in values):
                continue

            if values[0] == "series_index":
                in_series_section = True
                continue

            if in_series_section:
                index, key, name, chart_type, axis, grouping, x_key, size_key = (values[:8] + [None] * 8)[:8]
                metadata["series"].append(
                    {
                        "index": index,
                        "key": key,
                        "name": name,
                        "type": chart_type,
                        "axis": axis,
                        "grouping": grouping,
                        "x_key": x_key,
                        "size_key": size_key,
                    }
                )
                continue

            field_name = values[0]
            field_value = values[1] if len(values) > 1 else None
            if field_name == "categories_col":
                metadata["categories_col"] = field_value
            elif field_name == "chart_family":
                metadata["chart_family"] = field_value
            elif field_name == "chart_metadata_json" and field_value:
                try:
                    metadata["chart_metadata"] = json.loads(field_value)
                except (TypeError, json.JSONDecodeError):
                    metadata["chart_metadata"] = None

        return metadata

    @staticmethod
    def _is_missing_header(value) -> bool:
        if value is None:
            return True
        # pandas surfaces None column names as float NaN (and pd.NA as <NA>);
        # treat any null-like header as missing so metadata fallback kicks in.
        try:
            if pd.isna(value):
                return True
        except (TypeError, ValueError):
            pass
        text = str(value).strip()
        if text.lower() in ("nan", "none", "<na>"):
            return True
        return text == "" or text.lower().startswith("unnamed:")

    def _infer_categories_col_name(self, series: pd.Series) -> str:
        non_null = series.dropna()
        if non_null.empty:
            return "分类"

        parsed = pd.to_datetime(non_null, errors="coerce")
        if parsed.notna().all():
            # If every parsed timestamp is Jan-01 style and the original values
            # looked year-like, expose a more helpful semantic name.
            if self._looks_like_year_series(non_null):
                return "年份"
            return "日期"

        if self._looks_like_year_series(non_null):
            return "年份"

        return "分类"

    @staticmethod
    def _looks_like_year_series(series: pd.Series) -> bool:
        try:
            numeric = pd.to_numeric(series, errors="coerce")
            if numeric.notna().all():
                years = numeric.astype(int)
                return bool(((years >= 1900) & (years <= 2100)).all())
        except Exception:
            pass

        try:
            text = series.astype(str).str.strip()
            if text.str.fullmatch(r"\d{4}").all():
                years = text.astype(int)
                return bool(((years >= 1900) & (years <= 2100)).all())
        except Exception:
            pass

        return False
    
    def _extract_data_from_xml(self) -> Tuple[pd.DataFrame, str]:
        """
        从 XML 中提取数据点（备用方法）
        
        Returns:
            (df, categories_col)
        """
        if self._is_xy_chart():
            return self._extract_xy_data_from_xml()

        print("  → 尝试从 XML 提取数据...")
        from datetime import datetime, timedelta
        
        # 查找所有绘图元素
        plot_types = ['barChart', 'lineChart', 'areaChart', 'scatterChart', 'bubbleChart']
        
        data_dict = {}
        categories = []
        categories_col = "分类"
        max_data_points = 0
        
        # 第一遍：找到所有系列和最大数据点数量
        for plot_tag in plot_types:
            plot_elements = self.plotArea.findall(f'.//c:{plot_tag}', namespaces=NAMESPACES)
            
            for plot_element in plot_elements:
                # 提取分类数据（只需提取一次）
                if not categories:
                    cat_elem = plot_element.find('.//c:cat', namespaces=NAMESPACES)
                    if cat_elem is not None:
                        # ⭐ 优先尝试 strCache（文本）
                        cat_cache = cat_elem.find('.//c:strCache', namespaces=NAMESPACES)
                        if cat_cache is not None:
                            # ⭐ 修复 A: 强制将字符串转为 datetime
                            print("  → 发现 strCache，尝试转换为 datetime...")
                            for pt in cat_cache.findall('.//c:pt', namespaces=NAMESPACES):
                                v = pt.find('.//c:v', namespaces=NAMESPACES)
                                if v is not None:
                                    try:
                                        # 尝试用 pandas 智能解析多种日期格式
                                        dt = pd.to_datetime(v.text).to_pydatetime()
                                        categories.append(dt)
                                        print(f"    • 转换: {v.text} → {dt}")
                                    except (ValueError, TypeError):
                                        categories.append(v.text)  # 回退到字符串
                            if categories:
                                print(f"  → 成功将 {len(categories)} 个字符串日期转换为 datetime")
                        else:
                            # ⭐ 如果没有 strCache，尝试 numCache（日期轴）
                            num_cache = cat_elem.find('.//c:numCache', namespaces=NAMESPACES)
                            if num_cache is not None:
                                # 提取格式代码（如 "yyyy/mm"）
                                format_code = num_cache.find('.//c:formatCode', namespaces=NAMESPACES)
                                date_format = format_code.text if format_code is not None else None
                                
                                # ⭐ 修复 B: 强制将 Excel 序列号转为 datetime
                                print(f"  → 发现 numCache（格式: {date_format}），转换为 datetime...")
                                for pt in num_cache.findall('.//c:pt', namespaces=NAMESPACES):
                                    v = pt.find('.//c:v', namespaces=NAMESPACES)
                                    if v is not None:
                                        try:
                                            excel_date_num = float(v.text)
                                            # Excel 基准是 1899-12-30
                                            dt = datetime(1899, 12, 30) + timedelta(days=excel_date_num)
                                            categories.append(dt)
                                        except (ValueError, OverflowError):
                                            categories.append(v.text)  # 回退到字符串
                                
                                if categories:
                                    print(f"  → 成功将 {len(categories)} 个 Excel 序列号转换为 datetime")
                
                # 提取每个系列的数据
                series_elements = plot_element.findall('.//c:ser', namespaces=NAMESPACES)
                for ser in series_elements:
                    # 获取系列名称
                    name = self._extract_series_name(ser)
                    
                    # 获取数值
                    values = self._extract_series_values(ser)
                    
                    if name and values:
                        data_dict[name] = values
                        max_data_points = max(max_data_points, len(values))
        
        # 如果没有分类数据，使用最大数据点数量生成索引
        if not categories and max_data_points > 0:
            categories = [f"项目{i+1}" for i in range(max_data_points)]
            print(f"  → 未找到分类数据，生成 {len(categories)} 个默认分类")
        elif not categories and not data_dict:
            # 如果完全没有数据，返回空 DataFrame
            print("  ⚠️ 未找到任何数据，返回空 DataFrame")
            return pd.DataFrame({categories_col: []}), categories_col
        
        # 确保所有系列的长度与分类数量一致
        target_length = len(categories)
        for name, values in data_dict.items():
            if len(values) < target_length:
                # 补齐 None
                data_dict[name] = values + [None] * (target_length - len(values))
            elif len(values) > target_length:
                # 截断
                data_dict[name] = values[:target_length]
                print(f"  → 系列 '{name}' 数据点从 {len(values)} 截断到 {target_length}")
        
        # 构建 DataFrame
        df_data = {categories_col: categories}
        df_data.update(data_dict)
        df = pd.DataFrame(df_data)
        
        print(f"  → 从 XML 提取了 {len(categories)} 行数据，{len(data_dict)} 个系列")
        
        return df, categories_col
    
    def _parse_series_config(self) -> List[Dict]:
        """
        解析系列配置
        
        Returns:
            系列配置列表
        """
        series_config = []
        
        # 图表类型映射
        plot_type_map = {
            'barChart': 'bar',
            'lineChart': 'line',
            'areaChart': 'area',
            'scatterChart': 'scatter',
            'bubbleChart': 'bubble',
        }
        
        # 遍历所有图表类型
        series_index = 0
        for plot_tag, chart_type in plot_type_map.items():
            plot_elements = self.plotArea.findall(f'.//c:{plot_tag}', namespaces=NAMESPACES)
            
            for plot_element in plot_elements:
                # ⭐ 关键修复：在绘图元素级别获取轴引用
                plot_axis_type = self._extract_plot_axis(plot_element)
                plot_grouping = self._extract_plot_grouping(plot_element)
                
                # 遍历该绘图中的所有系列
                series_elements = plot_element.findall('.//c:ser', namespaces=NAMESPACES)
                
                for ser in series_elements:
                    # 解析单个系列
                    config_entry = self._parse_single_series(
                        ser,
                        chart_type,
                        plot_axis_type,
                        series_index,
                        plot_grouping,
                    )
                    if config_entry:
                        series_config.append(config_entry)
                    series_index += 1
        
        return series_config
    
    def _extract_plot_axis(self, plot_element) -> str:
        """
        提取绘图元素使用的坐标轴（主轴/次轴）
        
        Args:
            plot_element: 绘图元素 (<c:barChart>, <c:lineChart> 等)
            
        Returns:
            'primary' 或 'secondary'
        """
        # 查找绘图元素引用的轴 ID
        ax_id_refs = plot_element.findall('.//c:axId', namespaces=NAMESPACES)
        
        for ax_id_ref in ax_id_refs:
            ax_id = ax_id_ref.get('val')
            
            if ax_id in self.axis_map:
                axis_info = self.axis_map[ax_id]
                axis_type = axis_info['type']
                
                # 如果是值轴（主轴或次轴），返回它
                if axis_type in ('primary', 'secondary'):
                    return axis_type
        
        # 默认返回主轴
        return 'primary'
    
    def _extract_plot_grouping(self, plot_element) -> Optional[str]:
        grouping = plot_element.find('.//c:grouping', namespaces=NAMESPACES)
        if grouping is None:
            return None

        value = grouping.get('val')
        if value == 'percentStacked':
            return 'percent_stacked'
        return value

    def _parse_single_series(
        self,
        ser,
        chart_type: str,
        axis_type: str,
        series_index: int,
        grouping: Optional[str],
    ) -> Optional[Dict]:
        """
        解析单个系列
        
        Args:
            ser: 系列元素 (<c:ser>)
            chart_type: 图表类型
            axis_type: 轴类型（从父绘图元素获取）
            
        Returns:
            系列配置字典，或 None
        """
        config_entry = {}
        
        # A. 获取系列名称
        name = self._extract_series_name(ser)
        if not name:
            return None
        
        metadata = self._match_series_metadata(series_index, name)
        config_entry["name"] = metadata.get("name", name) if metadata else name
        config_entry["key"] = metadata.get("key", name) if metadata else name

        # B. 设置图表类型
        config_entry["type"] = metadata.get("type", chart_type) if metadata else chart_type

        # C. 使用传入的轴分配
        config_entry["axis"] = metadata.get("axis", axis_type) if metadata else axis_type
        grouping_value = metadata.get("grouping", grouping) if metadata else grouping
        if grouping_value:
            config_entry["grouping"] = grouping_value
        if metadata and metadata.get("x_key"):
            config_entry["x_key"] = metadata["x_key"]
        if metadata and metadata.get("size_key"):
            config_entry["size_key"] = metadata["size_key"]

        return config_entry

    def _match_series_metadata(self, series_index: int, parsed_name: str) -> Optional[Dict]:
        metadata_series = self._embedded_metadata.get("series", [])
        if not metadata_series:
            return None

        for series in metadata_series:
            if series.get("index") == series_index:
                return series

        for series in metadata_series:
            if series.get("name") == parsed_name:
                return series

        return None
    
    def _extract_series_name(self, ser) -> Optional[str]:
        """
        提取系列名称
        
        Args:
            ser: 系列元素 (<c:ser>)
            
        Returns:
            系列名称，或 None
        """
        # 方法1: 从 <c:tx> 中提取
        tx_elem = ser.find('.//c:tx', namespaces=NAMESPACES)
        if tx_elem is not None:
            # 尝试从 strRef 中提取
            v_elem = tx_elem.find('.//c:v', namespaces=NAMESPACES)
            if v_elem is not None and v_elem.text:
                return v_elem.text
        
        return None
    
    def _extract_series_values(self, ser) -> List:
        """
        提取系列数值
        
        Args:
            ser: 系列元素 (<c:ser>)
            
        Returns:
            数值列表
        """
        values = []
        
        # 尝试从 <c:val> 提取（柱状图、折线图、面积图）
        val_elem = ser.find('.//c:val', namespaces=NAMESPACES)
        if val_elem is not None:
            num_cache = val_elem.find('.//c:numCache', namespaces=NAMESPACES)
            if num_cache is not None:
                for pt in num_cache.findall('.//c:pt', namespaces=NAMESPACES):
                    v = pt.find('.//c:v', namespaces=NAMESPACES)
                    if v is not None:
                        try:
                            values.append(float(v.text))
                        except (ValueError, TypeError):
                            values.append(0)
        
        # 尝试从 <c:yVal> 提取（散点图）
        if not values:
            yVal_elem = ser.find('.//c:yVal', namespaces=NAMESPACES)
            if yVal_elem is not None:
                num_cache = yVal_elem.find('.//c:numCache', namespaces=NAMESPACES)
                if num_cache is not None:
                    for pt in num_cache.findall('.//c:pt', namespaces=NAMESPACES):
                        v = pt.find('.//c:v', namespaces=NAMESPACES)
                        if v is not None:
                            try:
                                values.append(float(v.text))
                            except (ValueError, TypeError):
                                values.append(0)
        
        return values

    def _is_xy_chart(self) -> bool:
        return bool(
            self.plotArea.find('.//c:scatterChart', namespaces=NAMESPACES) is not None
            or self.plotArea.find('.//c:bubbleChart', namespaces=NAMESPACES) is not None
        )

    def _extract_xy_data_from_xml(self) -> Tuple[pd.DataFrame, str]:
        print("  → 尝试从 XY 图表 XML 提取数据...")

        data_dict: Dict[str, List] = {}
        x_values: List[float] = []
        categories_col = str(self._embedded_metadata.get("categories_col") or "x")
        series_index = 0

        for plot_tag in ('scatterChart', 'bubbleChart'):
            plot_elements = self.plotArea.findall(f'.//c:{plot_tag}', namespaces=NAMESPACES)

            for plot_element in plot_elements:
                series_elements = plot_element.findall('.//c:ser', namespaces=NAMESPACES)

                for ser in series_elements:
                    name = self._extract_series_name(ser) or f"series_{series_index}"
                    metadata = self._match_series_metadata(series_index, name) or {}

                    current_x_values = self._extract_numeric_points(ser, 'xVal')
                    y_values = self._extract_numeric_points(ser, 'yVal')
                    if not x_values:
                        x_values = current_x_values
                    elif current_x_values != x_values:
                        raise ValueError("当前解析器仅支持共享同一组 X 数据的 scatter/bubble 图表")

                    y_key = metadata.get("key", name)
                    data_dict[y_key] = y_values

                    bubble_sizes = self._extract_numeric_points(ser, 'bubbleSize')
                    if bubble_sizes:
                        size_key = metadata.get("size_key") or f"{y_key}_size"
                        data_dict[size_key] = bubble_sizes

                    series_index += 1

        if not x_values:
            return pd.DataFrame({categories_col: []}), categories_col

        df_data = {categories_col: x_values}
        df_data.update(data_dict)
        return pd.DataFrame(df_data), categories_col

    def _extract_numeric_points(self, ser, tag_name: str) -> List[float]:
        values: List[float] = []
        element = ser.find(f'.//c:{tag_name}', namespaces=NAMESPACES)
        if element is None:
            return values

        num_cache = element.find('.//c:numCache', namespaces=NAMESPACES)
        if num_cache is None:
            return values

        for pt in num_cache.findall('.//c:pt', namespaces=NAMESPACES):
            v = pt.find('.//c:v', namespaces=NAMESPACES)
            if v is None:
                continue
            try:
                values.append(float(v.text))
            except (TypeError, ValueError):
                values.append(0.0)

        return values


def parse_chart_from_pptx(
    pptx_path: str,
    slide_idx: int = 0,
    shape_idx: int = 0,
) -> Tuple[List[Dict], pd.DataFrame, str, Optional[Dict]]:
    """
    从 PPTX 文件解析图表
    
    Args:
        pptx_path: PPTX 文件路径
        slide_idx: 幻灯片索引（0-based）
        shape_idx: 图表形状索引（0-based）
        
    Returns:
        (series_config, df, categories_col, layout_info)
        
    Example:
        >>> series_config, df, categories_col, layout_info = parse_chart_from_pptx("template.pptx")
        >>> # 现在可以用这些配置来复现图表
        >>> builder.build(series_config)
    """
    # 加载 PPTX
    from pathlib import Path

    pptx_path = str(Path(pptx_path).expanduser().resolve())
    prs = Presentation(pptx_path)
    
    if slide_idx >= len(prs.slides):
        raise ValueError(f"幻灯片索引 {slide_idx} 超出范围（共 {len(prs.slides)} 张）")
    
    slide = prs.slides[slide_idx]
    
    # 查找图表
    chart = None
    chart_shape = None
    
    for shape in slide.shapes:
        if hasattr(shape, 'chart'):
            if shape_idx == 0:
                chart = shape.chart
                break
            shape_idx -= 1
    
    if chart is None:
        raise ValueError(f"在幻灯片 {slide_idx} 上未找到图表")
    
    # 解析图表
    parser = ChartParser(chart)
    return parser.parse()


def parse_all_charts_from_pptx(pptx_path: str) -> List[Dict]:
    """
    解析 PPTX 文件中的所有图表
    
    Args:
        pptx_path: PPTX 文件路径
        
    Returns:
        图表信息列表
        [
            {
                "slide_idx": 0,
                "shape_idx": 0,
                "shape_name": "图表 1",
                "series_config": [...],
                "df": DataFrame,
                "categories_col": "日期"
            },
            ...
        ]
    """
    prs = Presentation(pptx_path)
    all_charts = []
    
    print(f"\n🔍 扫描 PPTX 文件: {pptx_path}")
    print(f"  - 共 {len(prs.slides)} 张幻灯片\n")
    
    for slide_idx, slide in enumerate(prs.slides):
        chart_count = 0
        
        for shape_idx, shape in enumerate(slide.shapes):
            # 使用 has_chart 属性检查，而不是直接访问 chart
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart_count += 1
                
                try:
                    # 解析图表
                    parser = ChartParser(shape.chart)
                    series_config, df, categories_col, layout_info = parser.parse()
                    
                    chart_info = {
                        "slide_index": slide_idx,
                        "chart_index": chart_count - 1,  # 0-based
                        "shape_idx": shape_idx,
                        "shape_name": shape.name if hasattr(shape, 'name') else f"图表 {chart_count}",
                        "series_config": series_config,
                        "df": df,
                        "categories_col": categories_col,
                        "layout_info": layout_info  # 新增：布局信息
                    }
                    
                    all_charts.append(chart_info)
                    
                    print(f"✅ 幻灯片 {slide_idx + 1}, 图表 {chart_count}: {chart_info['shape_name']}")
                    print(f"   - {len(series_config)} 个系列")
                    
                except Exception as e:
                    print(f"⚠️ 幻灯片 {slide_idx + 1}, 图表 {chart_count} 解析失败: {e}")
        
        if chart_count > 0:
            print(f"  → 幻灯片 {slide_idx + 1}: 发现 {chart_count} 个图表\n")
    
    print(f"\n📊 总计解析 {len(all_charts)} 个图表")
    
    return all_charts


def parse_semantic_component_from_pptx(
    pptx_path: str,
    slide_idx: int = 0,
    component_idx: int = 0,
) -> Dict:
    """Parse one semantic component (chart-backed or anchor-backed) from a PPTX."""

    components = parse_all_semantic_components_from_pptx(pptx_path, slide_idx=slide_idx)
    if component_idx >= len(components):
        raise ValueError(f"语义组件索引 {component_idx} 超出范围（共 {len(components)} 个）")
    return components[component_idx]


def parse_all_semantic_components_from_pptx(
    pptx_path: str,
    slide_idx: int | None = None,
) -> List[Dict]:
    """Discover semantic-family components, including non-chart anchor-backed ones."""

    prs = Presentation(pptx_path)
    components: list[dict] = []

    slide_iter = enumerate(prs.slides)
    for current_slide_idx, slide in slide_iter:
        if slide_idx is not None and current_slide_idx != slide_idx:
            continue

        for shape_idx, shape, metadata in iter_semantic_anchors(slide):
            family = metadata.get("semantic_family") or metadata.get("chart_family")
            if not family:
                continue
            components.append(
                {
                    "slide_index": current_slide_idx,
                    "component_index": len(components),
                    "shape_idx": shape_idx,
                    "shape_name": getattr(shape, "name", f"semantic-anchor-{shape_idx}"),
                    "family": family,
                    "base_chart_family": metadata.get("base_chart_family"),
                    "metadata": metadata,
                    "source": "anchor",
                }
            )

        chart_count = 0
        for shape_idx, shape in enumerate(slide.shapes):
            if not hasattr(shape, "has_chart") or not shape.has_chart:
                continue

            chart_count += 1
            try:
                parser = ChartParser(shape.chart)
                series_config, df, categories_col, layout_info = parser.parse()
            except Exception as exc:
                print(f"⚠️ 幻灯片 {current_slide_idx + 1}, 图表 {chart_count} 语义解析失败: {exc}")
                continue

            chart_metadata = layout_info.get("chart_metadata") if layout_info else None
            if not isinstance(chart_metadata, dict):
                continue

            family = chart_metadata.get("semantic_family") or chart_metadata.get("chart_family")
            if not family:
                continue
            if chart_metadata.get("semantic_parent_family"):
                continue

            components.append(
                {
                    "slide_index": current_slide_idx,
                    "component_index": len(components),
                    "shape_idx": shape_idx,
                    "shape_name": getattr(shape, "name", f"图表 {chart_count}"),
                    "family": family,
                    "base_chart_family": chart_metadata.get("base_chart_family"),
                    "metadata": chart_metadata,
                    "layout_info": layout_info,
                    "series_config": series_config,
                    "df": df,
                    "categories_col": categories_col,
                    "source": "chart",
                }
            )

    return components
